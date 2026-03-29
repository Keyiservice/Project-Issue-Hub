from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ppt_corporate_assets"
OUT_FILE = ROOT / "非标设备项目现场问题协同系统-企业风格增强版.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
IMG_W = 1600
IMG_H = 900

FONT_REG = Path(r"C:\Windows\Fonts\msyh.ttc")
FONT_BOLD = Path(r"C:\Windows\Fonts\msyhbd.ttc")
FONT_ALT = Path(r"C:\Windows\Fonts\arialbd.ttf")


PALETTE = {
    "navy": (20, 62, 128),
    "navy_dark": (15, 49, 105),
    "cyan": (38, 214, 234),
    "cyan_soft": (204, 247, 250),
    "white": (255, 255, 255),
    "paper": (250, 252, 255),
    "ink": (23, 48, 95),
    "muted": (84, 101, 129),
    "line": (220, 231, 243),
    "pale": (232, 245, 252),
    "panel": (255, 255, 255),
    "panel_alt": (244, 249, 252),
    "orange": (210, 115, 56),
    "teal": (37, 171, 179),
    "green": (87, 162, 84),
    "red": (198, 88, 72),
    "gray_bg": (242, 246, 250),
}


def f(size: int, bold: bool = False):
    font_path = FONT_BOLD if bold and FONT_BOLD.exists() else FONT_REG
    if not font_path.exists():
        font_path = FONT_ALT
    return ImageFont.truetype(str(font_path), size)


def ensure_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def rounded(draw: ImageDraw.ImageDraw, box, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def draw_text(draw, xy, text, size, color, bold=False, anchor="la"):
    draw.text(xy, text, font=f(size, bold), fill=color, anchor=anchor)


def wrap_text(draw, text, font_obj, max_width):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        trial = word if not current else f"{current} {word}"
        if draw.textbbox((0, 0), trial, font=font_obj)[2] <= max_width:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def add_multiline(draw, box, text, size, color, bold=False, line_gap=10):
    font_obj = f(size, bold)
    max_width = int(box[2] - box[0])
    lines = []
    for paragraph in text.split("\n"):
        if not paragraph:
            lines.append("")
            continue
        lines.extend(wrap_text(draw, paragraph, font_obj, max_width))
    x, y = box[0], box[1]
    _, _, _, h = draw.textbbox((0, 0), "Ag", font=font_obj)
    step = h + line_gap
    for line in lines:
        draw.text((x, y), line, font=font_obj, fill=color)
        y += step


def draw_stripes(img: Image.Image, bbox, color, spacing=22, width=3, angle=28):
    overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
    d = ImageDraw.Draw(overlay)
    x0, y0, x1, y1 = bbox
    for x in range(int(x0 - (y1 - y0)), int(x1 + (y1 - y0)), spacing):
        d.line((x, y1, x + (y1 - y0), y0), fill=color, width=width)
    rotated = overlay.rotate(angle, resample=Image.Resampling.BICUBIC, center=((x0 + x1) / 2, (y0 + y1) / 2))
    img.alpha_composite(rotated)


def make_cover_bg(path: Path):
    img = Image.new("RGBA", (IMG_W, IMG_H), PALETTE["navy_dark"] + (255,))
    base = ImageDraw.Draw(img)
    for x in range(IMG_W):
        ratio = x / IMG_W
        r = int(PALETTE["navy_dark"][0] * (1 - ratio) + PALETTE["navy"][0] * ratio)
        g = int(PALETTE["navy_dark"][1] * (1 - ratio) + PALETTE["navy"][1] * ratio)
        b = int(PALETTE["navy_dark"][2] * (1 - ratio) + PALETTE["navy"][2] * ratio)
        base.line((x, 0, x, IMG_H), fill=(r, g, b, 255))
    arc_layer = Image.new("RGBA", (IMG_W, IMG_H), (255, 255, 255, 0))
    arc = ImageDraw.Draw(arc_layer)
    arc.ellipse((720, -330, 1820, 920), fill=(44, 124, 204, 120))
    arc.ellipse((350, 270, 1680, 1350), fill=(18, 78, 165, 120))
    draw_stripes(arc_layer, (1010, -40, 1600, 420), (160, 212, 255, 78), spacing=18, width=2, angle=0)
    img.alpha_composite(arc_layer)
    img = img.filter(ImageFilter.GaussianBlur(0.2))
    img.convert("RGB").save(path)


def make_inner_bg(path: Path):
    img = Image.new("RGBA", (IMG_W, IMG_H), PALETTE["paper"] + (255,))
    layer = Image.new("RGBA", (IMG_W, IMG_H), (255, 255, 255, 0))
    draw = ImageDraw.Draw(layer)
    draw.ellipse((720, -120, 1800, 760), fill=(228, 244, 253, 185))
    draw.ellipse((-120, 450, 1620, 1080), fill=(255, 255, 255, 180))
    draw_stripes(layer, (1180, 70, 1590, 890), (160, 206, 241, 62), spacing=16, width=2, angle=0)
    img.alpha_composite(layer)
    img.convert("RGB").save(path)


def mock_window_base(title: str, size=(1120, 640)):
    img = Image.new("RGB", size, PALETTE["panel"])
    draw = ImageDraw.Draw(img)
    rounded(draw, (0, 0, size[0] - 1, size[1] - 1), 26, PALETTE["panel"], outline=(214, 226, 236), width=2)
    rounded(draw, (24, 20, size[0] - 24, 78), 18, (242, 247, 252), outline=(226, 235, 243), width=1)
    draw_text(draw, (52, 48), title, 28, PALETTE["ink"], True, "lm")
    for i, col in enumerate([(239, 108, 95), (248, 200, 71), (101, 204, 116)]):
        draw.ellipse((size[0] - 132 + i * 28, 36, size[0] - 112 + i * 28, 56), fill=col)
    return img, draw


def panel(draw, box, title="", subtitle="", tone="navy"):
    tone_map = {
        "navy": PALETTE["navy"],
        "cyan": PALETTE["teal"],
        "teal": PALETTE["teal"],
        "orange": PALETTE["orange"],
        "red": PALETTE["red"],
        "green": PALETTE["green"],
    }
    rounded(draw, box, 20, PALETTE["panel"], outline=(223, 232, 239), width=2)
    x0, y0, x1, y1 = box
    draw.rounded_rectangle((x0 + 12, y0 + 12, x0 + 18, y1 - 12), radius=4, fill=tone_map[tone])
    if title:
        draw_text(draw, (x0 + 36, y0 + 28), title, 22, PALETTE["ink"], True)
    if subtitle:
        add_multiline(draw, (x0 + 36, y0 + 58, x1 - 22, y1 - 22), subtitle, 15, PALETTE["muted"])


def draw_kpi(draw, box, title, value, note, tone="navy"):
    panel(draw, box, "", "", tone)
    x0, y0, x1, y1 = box
    draw_text(draw, (x0 + 28, y0 + 26), title, 16, PALETTE["muted"])
    draw_text(draw, (x0 + 28, y0 + 64), value, 34, PALETTE["ink"], True)
    draw_text(draw, (x0 + 28, y1 - 26), note, 14, PALETTE["muted"], False, "lb")


def bullet_lines(draw, start_xy, items, color=None, size=16, gap=12):
    x, y = start_xy
    color = color or PALETTE["ink"]
    font_obj = f(size)
    _, _, _, h = draw.textbbox((0, 0), "Ag", font=font_obj)
    for item in items:
        draw.ellipse((x, y + 8, x + 8, y + 16), fill=PALETTE["teal"])
        draw.text((x + 18, y), item, font=font_obj, fill=color)
        y += h + gap


def render_dashboard(path: Path):
    img, draw = mock_window_base("驾驶舱 / Control Tower")
    rounded(draw, (24, 98, 220, 620), 22, PALETTE["navy_dark"])
    draw_text(draw, (58, 142), "Field Issue Ops", 18, (170, 224, 248), False)
    for i, label in enumerate(["驾驶舱", "项目问题库", "项目协同", "统计分析", "用户管理"]):
        y = 190 + i * 82
        fill = PALETTE["orange"] if i == 1 else (33, 49, 72)
        rounded(draw, (44, y, 200, y + 54), 18, fill)
        draw_text(draw, (66, y + 27), label, 18, PALETTE["white"], True, "lm")
    draw_text(draw, (260, 132), "FS21 项目驾驶舱", 30, PALETTE["ink"], True)
    draw_text(draw, (260, 168), "支持项目选择器；未选择时可自动轮巡。", 18, PALETTE["muted"])
    for idx, (title, value) in enumerate([("问题总数", "67"), ("未关闭", "60"), ("高优先级", "54"), ("已超期", "53")]):
        x = 260 + idx * 198
        draw_kpi(draw, (x, 206, x + 182, 318), title, value, "整项目口径")
    panel(draw, (260, 338, 698, 604), "新增 / 关闭趋势", "支持按项目看趋势变化、关注超期与风险。", "cyan")
    chart = [(300, 540), (360, 490), (420, 500), (480, 430), (540, 460), (600, 402), (660, 372)]
    draw.line(chart, fill=PALETTE["teal"], width=5)
    for x, y in chart:
        draw.ellipse((x - 6, y - 6, x + 6, y + 6), fill=PALETTE["white"], outline=PALETTE["teal"], width=3)
    panel(draw, (720, 338, 1088, 604), "重点问题", "点击后直接进入问题详情。", "orange")
    for idx, txt in enumerate(["FS21-OPL-063 / 已超期 / 张鹏", "FS21-OPL-062 / 处理中 / Xuchu Liu", "FS21-OPL-061 / 待验证 / Jason Liu"]):
        rounded(draw, (742, 396 + idx * 62, 1066, 448 + idx * 62), 16, PALETTE["panel_alt"], outline=(229, 236, 242))
        draw_text(draw, (764, 422 + idx * 62), txt, 16, PALETTE["ink"], idx == 0, "lm")
    img.save(path)


def render_issue_library(path: Path):
    img, draw = mock_window_base("项目问题库 / Project OPL")
    panel(draw, (26, 96, 1094, 168), "FS21 项目问题库", "先选项目，再进入该项目的问题清单；支持在 Web 端直接录入新问题。", "orange")
    panel(draw, (26, 184, 1094, 278), "筛选区", "关键字 / 责任人 / 状态 / 优先级；顶部统计按整项目口径显示。", "cyan")
    for idx, title in enumerate(["全部 OPL 67", "处理中 58", "已超期 52", "已关闭 7"]):
        x = 26 + idx * 266
        draw_kpi(draw, (x, 294, x + 246, 388), title.split()[0], title.split()[-1], title.split()[1] if len(title.split()) > 2 else "项目统计")
    panel(draw, (26, 404, 720, 614), "问题清单", "支持从问题清单直接进入详情，或在右上角点击“新建问题”。", "navy")
    for idx, issue in enumerate(["ISS202603177835  123122312", "FS21-OPL-064  ST30 harting头未固定", "FS21-OPL-063  交付物需要提交"]):
        y = 462 + idx * 48
        rounded(draw, (48, y, 698, y + 38), 12, PALETTE["gray_bg"])
        draw_text(draw, (64, y + 19), issue, 17, PALETTE["ink"], idx == 0, "lm")
        draw_text(draw, (640, y + 19), ["待分派", "处理中", "已超期"][idx], 14, [PALETTE["orange"], PALETTE["teal"], PALETTE["red"]][idx], True, "rm")
    panel(draw, (738, 404, 1094, 614), "新建问题", "Web 端也可录入问题：标题、描述、优先级、责任人、图片/视频附件。", "orange")
    bullet_lines(draw, (760, 456), ["项目自动带入", "支持图片与视频上传", "单文件限制 5MB", "录入后立即进入项目问题库"], size=16)
    img.save(path)


def render_issue_detail(path: Path):
    img, draw = mock_window_base("问题详情 / Issue Detail")
    panel(draw, (26, 94, 706, 292), "FS21-OPL-064 / ST30 harting头还有未固定的需要固定", "当前阶段、责任人、下一步、录入时间集中展示；描述下方直接看附件缩略图。", "orange")
    for idx in range(3):
        x = 48 + idx * 208
        rounded(draw, (x, 180, x + 182, 270), 16, PALETTE["gray_bg"], outline=(229, 236, 242))
        if idx < 2:
            rounded(draw, (x + 12, 192, x + 170, 232), 10, [(233, 245, 250), (235, 240, 250)][idx])
            draw_text(draw, (x + 91, 212), ["图片预览", "视频预览"][idx], 15, PALETTE["muted"], True, "mm")
            if idx == 1:
                rounded(draw, (x + 64, 198, x + 118, 226), 12, PALETTE["navy"])
                draw_text(draw, (x + 91, 212), "PLAY", 12, PALETTE["white"], True, "mm")
        else:
            draw_text(draw, (x + 22, 206), "关闭证据", 16, PALETTE["muted"], True)
            draw_text(draw, (x + 22, 236), "支持图片/视频", 15, PALETTE["ink"])
    panel(draw, (724, 94, 1092, 354), "推进与关闭", "右侧可做责任分派、优先级调整、状态推进和关闭问题。", "teal")
    bullet_lines(draw, (748, 154), ["责任人只能从当前项目团队选择", "关闭时填写说明、关闭证据", "关闭人和关闭时间自动记录", "误录问题可执行逻辑删除"], size=16)
    panel(draw, (26, 312, 706, 614), "时间线 / 跟进", "系统日志与人工跟进统一进入时间线。", "navy")
    for idx, row in enumerate(["18:35  System  状态变更 IN_PROGRESS → PENDING_VERIFY", "18:48  Zhang Peng  上传现场照片 2 张", "19:05  PM  提交关闭说明和证据视频"]):
        rounded(draw, (48, 368 + idx * 70, 684, 422 + idx * 70), 16, PALETTE["gray_bg"])
        draw_text(draw, (66, 395 + idx * 70), row, 16, PALETTE["ink"], idx == 0, "lm")
    panel(draw, (724, 372, 1092, 614), "当前推进摘要", "责任人：Zhang Peng\n当前阶段：待验证\n下一步：确认关闭证据并关闭问题", "orange")
    img.save(path)


def render_project_team(path: Path):
    img, draw = mock_window_base("项目管理 / Project Management")
    panel(draw, (24, 94, 1094, 168), "项目信息", "项目编号、项目名称、客户、项目经理、开始/计划结束日期。", "orange")
    panel(draw, (24, 184, 670, 614), "项目团队", "项目成员用于限定责任人范围、项目岗位和问题权限。", "cyan")
    for idx, member in enumerate(["张鹏 / 机械工程师 / 可分派", "刘旭初 / 自动化工程师 / 可验证", "白海超 / 电气技术员 / 可关闭", "项目经理 / 可分派/验证/关闭"]):
        y = 248 + idx * 72
        rounded(draw, (48, y, 646, y + 54), 16, PALETTE["gray_bg"])
        draw_text(draw, (66, y + 27), member, 17, PALETTE["ink"], idx == 3, "lm")
    panel(draw, (694, 184, 1094, 614), "项目团队操作", "支持手工维护、批量导入、从现有 OPL 同步成员。", "teal")
    bullet_lines(draw, (720, 244), ["标记项目经理", "配置项目岗位", "控制可分派/验证/关闭", "同步 OPL 中已有责任人", "问题分派下拉只显示本项目成员"], size=16)
    img.save(path)


def render_user_mgmt(path: Path):
    img, draw = mock_window_base("用户管理 / User Management")
    for idx, title in enumerate(["用户总数", "已绑定微信", "待改密", "参与项目"]):
        x = 26 + idx * 266
        draw_kpi(draw, (x, 96, x + 246, 192), title, ["58", "32", "7", "49"][idx], "统一账号体系")
    panel(draw, (26, 212, 740, 614), "账号列表", "支持新建、编辑、批量导入、重置密码、解绑微信。", "navy")
    for idx, row in enumerate(["zhang.peng / 张鹏 / 机械工程师", "li.xuchu / 刘旭初 / 自动化工程师", "admin / 系统管理员 / ADMIN"]):
        y = 278 + idx * 74
        rounded(draw, (48, y, 718, y + 56), 16, PALETTE["gray_bg"])
        draw_text(draw, (66, y + 28), row, 17, PALETTE["ink"], idx == 2, "lm")
        draw_text(draw, (690, y + 28), ["已绑定", "已绑定", "未绑定"][idx], 15, [PALETTE["green"], PALETTE["green"], PALETTE["orange"]][idx], True, "rm")
    panel(draw, (760, 212, 1094, 614), "微信绑定流程", "企业账号为主身份，微信仅做登录入口和绑定关系。", "orange")
    bullet_lines(draw, (784, 280), ["管理员预先创建账号", "首次小程序登录时绑定企业账号", "若是初始密码则强制改密", "后续微信免密进入", "Web 端继续账号密码登录"], size=15)
    img.save(path)


def render_stats(path: Path):
    img, draw = mock_window_base("统计分析 / Analytics")
    panel(draw, (24, 94, 560, 614), "趋势分析", "新增 / 关闭趋势、超期风险和项目压力。", "cyan")
    bars = [320, 420, 360, 490, 540, 468]
    for idx, h in enumerate(bars):
        x = 62 + idx * 72
        draw.rounded_rectangle((x, 520 - h // 5, x + 38, 520), 10, fill=(71, 145, 223))
    panel(draw, (580, 94, 1094, 330), "状态结构", "可查看 NEW / IN_PROGRESS / PENDING_VERIFY / CLOSED 的分布。", "orange")
    for idx, (label, width) in enumerate([("处理中", 260), ("已关闭", 180), ("已超期", 140), ("待验证", 120)]):
        y = 154 + idx * 38
        draw_text(draw, (604, y), label, 16, PALETTE["ink"])
        rounded(draw, (690, y - 6, 690 + width, y + 18), 10, [(67, 162, 178), (82, 160, 96), (198, 88, 72), (210, 115, 56)][idx])
    panel(draw, (580, 348, 1094, 614), "字典标准化", "分类、来源、模块、部门等字典统一后，统计才有统一口径。", "navy")
    bullet_lines(draw, (604, 418), ["问题分类", "问题来源", "模块", "责任部门", "岗位与项目角色"], size=16)
    img.save(path)


def phone_base(title: str, size=(420, 860)):
    outer = Image.new("RGB", size, (36, 44, 56))
    d = ImageDraw.Draw(outer)
    rounded(d, (14, 14, size[0] - 14, size[1] - 14), 42, PALETTE["paper"])
    rounded(d, (120, 14, size[0] - 120, 48), 18, (15, 16, 18))
    screen = Image.new("RGB", (size[0] - 28, size[1] - 28), PALETTE["paper"])
    draw = ImageDraw.Draw(screen)
    draw_text(draw, (18, 18), "11:08", 18, PALETTE["ink"])
    draw_text(draw, (screen.size[0] // 2, 18), title, 19, PALETTE["ink"], True, "ma")
    rounded(draw, (screen.size[0] - 90, 12, screen.size[0] - 18, 44), 18, PALETTE["panel"], outline=(220, 230, 240), width=1)
    draw_text(draw, (screen.size[0] - 54, 28), "···", 18, PALETTE["ink"], True, "mm")
    return outer, screen, draw


def paste_screen(outer: Image.Image, screen: Image.Image):
    outer.paste(screen, (14, 14))
    return outer


def render_phone_home(path: Path):
    outer, screen, draw = phone_base("首页")
    panel(draw, (16, 70, 378, 278), "Issue Ops", "当前项目、录入入口、问题库和统计都放在首屏。", "orange")
    for idx, name in enumerate(["项目工作台", "OPL 录入", "项目问题库", "项目统计"]):
        x = 16 + (idx % 2) * 182
        y = 296 + (idx // 2) * 104
        rounded(draw, (x, y, x + 168, y + 88), 22, PALETTE["panel"], outline=(224, 233, 240), width=1)
        rounded(draw, (x + 16, y + 16, x + 72, y + 48), 16, PALETTE["navy_dark"])
        draw_text(draw, (x + 44, y + 32), ["OPS", "IN", "OPL", "STA"][idx], 15, PALETTE["white"], True, "mm")
        draw_text(draw, (x + 16, y + 64), name, 18, PALETTE["ink"], True)
    panel(draw, (16, 514, 378, 720), "个人工作区", "我的问题 / 待我处理 / 个人中心", "cyan")
    for idx, name in enumerate(["我的问题", "待我处理", "个人中心"]):
        x = 30 + idx * 118
        rounded(draw, (x, 570, x + 104, 664), 18, PALETTE["panel_alt"], outline=(226, 235, 242))
        draw_text(draw, (x + 52, 620), name, 16, PALETTE["ink"], True, "mm")
    paste_screen(outer, screen).save(path)


def render_phone_selector(path: Path):
    outer, screen, draw = phone_base("选择项目")
    panel(draw, (16, 70, 378, 190), "Project First", "先选项目，再进入工作台、录入页或项目问题库。", "orange")
    for idx, name in enumerate(["FS21 / ICE", "FS21 / Final Line", "A35 / Welding"]):
        y = 214 + idx * 128
        rounded(draw, (16, y, 378, y + 108), 22, PALETTE["panel"], outline=(223, 232, 239))
        draw_text(draw, (34, y + 26), name, 24, PALETTE["ink"], True)
        draw_text(draw, (34, y + 64), "项目经理 / 客户 / 日期", 15, PALETTE["muted"])
        rounded(draw, (276, y + 34, 360, y + 74), 18, PALETTE["navy"])
        draw_text(draw, (318, y + 54), "进入", 16, PALETTE["white"], True, "mm")
    paste_screen(outer, screen).save(path)


def render_phone_create(path: Path):
    outer, screen, draw = phone_base("新建问题")
    panel(draw, (16, 70, 378, 186), "Field Report", "录入人、录入时间、当前项目自动带入。", "orange")
    for idx, label in enumerate(["问题标题", "问题描述", "优先级 / 影响等级"]):
        y = 208 + idx * 108
        rounded(draw, (16, y, 378, y + (92 if idx == 1 else 74)), 18, PALETTE["panel"], outline=(223, 232, 239))
        draw_text(draw, (34, y + 18), label, 18, PALETTE["muted"])
    panel(draw, (16, 544, 378, 726), "现场媒体", "支持图片 / 视频；单文件 5MB，视频 15 秒。", "cyan")
    rounded(draw, (34, 606, 180, 690), 18, PALETTE["panel_alt"], outline=(223, 232, 239))
    rounded(draw, (214, 606, 360, 690), 18, PALETTE["panel_alt"], outline=(223, 232, 239))
    draw_text(draw, (107, 648), "上传图片", 17, PALETTE["ink"], True, "mm")
    draw_text(draw, (287, 648), "上传视频", 17, PALETTE["ink"], True, "mm")
    paste_screen(outer, screen).save(path)


def render_phone_library(path: Path):
    outer, screen, draw = phone_base("项目 OPL")
    panel(draw, (16, 70, 378, 166), "FS21 项目问题库", "支持整项目统计、状态筛选、责任人筛选。", "orange")
    for idx, (title, val) in enumerate([("全部", "67"), ("处理中", "58"), ("已超期", "52"), ("已关闭", "7")]):
        x = 16 + (idx % 2) * 182
        y = 190 + (idx // 2) * 94
        draw_kpi(draw, (x, y, x + 168, y + 82), title, val, "整项目")
    panel(draw, (16, 388, 378, 736), "问题清单", "点击即可进入问题详情。", "cyan")
    for idx, name in enumerate(["ISS202603187132 / 已关闭", "FS21-OPL-064 / 处理中", "FS21-OPL-063 / 已超期"]):
        y = 442 + idx * 82
        rounded(draw, (30, y, 364, y + 66), 16, PALETTE["panel_alt"], outline=(224, 233, 240))
        draw_text(draw, (46, y + 23), name, 16, PALETTE["ink"], idx == 0, "lm")
    paste_screen(outer, screen).save(path)


def render_phone_detail(path: Path):
    outer, screen, draw = phone_base("问题详情")
    panel(draw, (16, 70, 378, 220), "FS21-OPL-064", "描述下方直接预览图片 / 视频，并可点开中框查看。", "orange")
    rounded(draw, (34, 146, 170, 214), 16, PALETTE["panel_alt"], outline=(223, 232, 239))
    rounded(draw, (184, 146, 320, 214), 16, PALETTE["panel_alt"], outline=(223, 232, 239))
    draw_text(draw, (102, 181), "图片", 16, PALETTE["ink"], True, "mm")
    rounded(draw, (232, 160, 272, 200), 20, PALETTE["navy"])
    draw_text(draw, (252, 181), "PLAY", 12, PALETTE["white"], True, "mm")
    panel(draw, (16, 242, 378, 412), "当前推进", "责任人 / 当前阶段 / 下一步 / 关闭信息。", "cyan")
    panel(draw, (16, 430, 378, 736), "跟进与关闭证据", "跟进评论支持图片 / 视频；关闭时可提交关闭证据附件。", "navy")
    for idx, label in enumerate(["上传图片", "上传视频", "关闭问题"]):
        y = 516 + idx * 62
        rounded(draw, (36, y, 358, y + 46), 16, PALETTE["panel_alt"], outline=(223, 232, 239))
        draw_text(draw, (196, y + 23), label, 16, PALETTE["ink"], True, "mm")
    paste_screen(outer, screen).save(path)


def render_phone_profile(path: Path):
    outer, screen, draw = phone_base("个人中心")
    panel(draw, (16, 70, 378, 194), "个人中心", "查看账号、绑定状态、角色，并执行修改密码。", "orange")
    for idx, line in enumerate(["账号 / zhang.peng", "微信绑定 / 已绑定", "密码状态 / 已修改", "版本 / 1.0.0"]):
        y = 220 + idx * 76
        rounded(draw, (16, y, 378, y + 60), 16, PALETTE["panel"], outline=(224, 233, 240))
        draw_text(draw, (34, y + 30), line, 17, PALETTE["ink"], idx == 0, "lm")
    rounded(draw, (16, 570, 378, 628), 20, PALETTE["navy"])
    draw_text(draw, (197, 599), "修改密码", 18, PALETTE["white"], True, "mm")
    paste_screen(outer, screen).save(path)


def render_all_assets():
    ensure_assets()
    make_cover_bg(ASSET_DIR / "bg_cover.png")
    make_inner_bg(ASSET_DIR / "bg_inner.png")
    render_dashboard(ASSET_DIR / "web_dashboard.png")
    render_issue_library(ASSET_DIR / "web_issue_library.png")
    render_issue_detail(ASSET_DIR / "web_issue_detail.png")
    render_project_team(ASSET_DIR / "web_project_team.png")
    render_user_mgmt(ASSET_DIR / "web_user_mgmt.png")
    render_stats(ASSET_DIR / "web_stats.png")
    render_phone_home(ASSET_DIR / "mini_home.png")
    render_phone_selector(ASSET_DIR / "mini_selector.png")
    render_phone_create(ASSET_DIR / "mini_create.png")
    render_phone_library(ASSET_DIR / "mini_library.png")
    render_phone_detail(ASSET_DIR / "mini_detail.png")
    render_phone_profile(ASSET_DIR / "mini_profile.png")


def ppt_text(slide, left, top, width, height, text, size=18, color=(23, 48, 95), bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    return box


def ppt_bullets(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Microsoft YaHei"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(*PALETTE["teal"])
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(size)
        r2.font.color.rgb = RGBColor(*PALETTE["ink"])
    return box


def add_bg(slide, cover=False):
    bg = ASSET_DIR / ("bg_cover.png" if cover else "bg_inner.png")
    slide.shapes.add_picture(str(bg), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_brand(slide, page_num, light=False):
    color = PALETTE["white"] if light else PALETTE["ink"]
    sub = (198, 214, 232) if light else PALETTE["muted"]
    ppt_text(slide, Inches(10.92), Inches(0.32), Inches(1.85), Inches(0.28), "FIELD ISSUE OPS", 12, color, True, PP_ALIGN.RIGHT)
    ppt_text(slide, Inches(11.82), Inches(7.02), Inches(0.45), Inches(0.24), str(page_num), 10, sub, False, PP_ALIGN.RIGHT)


def add_header(slide, page_num, title, subtitle="", light=False):
    add_brand(slide, page_num, light)
    title_color = PALETTE["white"] if light else PALETTE["ink"]
    subtitle_color = (222, 231, 239) if light else PALETTE["muted"]
    deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(0.56), Inches(0.52), Inches(0.56), Inches(0.56))
    deco.fill.background()
    deco.line.color.rgb = RGBColor(*PALETTE["cyan"])
    deco.line.width = Pt(3)
    ppt_text(slide, Inches(0.95), Inches(0.52), Inches(8.8), Inches(0.5), title, 26, title_color, True)
    if subtitle:
        ppt_text(slide, Inches(0.98), Inches(0.98), Inches(9.1), Inches(0.42), subtitle, 12, subtitle_color)


def add_cover_text(slide):
    add_brand(slide, 1, True)
    ppt_text(slide, Inches(1.12), Inches(1.1), Inches(2.6), Inches(0.28), "Bridge to Execution", 16, PALETTE["cyan"], True)
    ppt_text(slide, Inches(1.12), Inches(2.32), Inches(8.6), Inches(0.9), "非标设备项目现场问题协同系统", 31, PALETTE["white"], True)
    ppt_text(slide, Inches(1.15), Inches(3.2), Inches(8.0), Inches(0.7), "企业风格增强版：项目优先、图像驱动、责任清晰、证据闭环", 16, (214, 224, 238))
    ppt_text(slide, Inches(1.15), Inches(6.22), Inches(2.5), Inches(0.3), "2026 / 03 / 25", 14, PALETTE["white"])


def add_image(slide, name, left, top, width=None, height=None):
    path = ASSET_DIR / name
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_panel(slide, left, top, width, height, fill=(255, 255, 255)):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*PALETTE["line"])
    shp.line.width = Pt(1)
    return shp


def build_slides():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_cover_text(slide)

    # 2 problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 2, "为什么要从 Excel 升级成项目化问题协同系统")
    add_panel(slide, Inches(0.72), Inches(1.55), Inches(12.0), Inches(5.55), fill=(255, 255, 255))
    ppt_bullets(slide, Inches(0.98), Inches(1.95), Inches(5.3), Inches(4.8), [
        "Excel 只有文字描述，现场图片、视频、责任归属和状态记录都不完整。",
        "管理层无法在 10 秒内看懂问题本质，需要反复追问“是什么问题、谁在处理、影响多大”。",
        "项目团队无法统一问题口径，责任和关闭验证缺少透明度。",
        "不同项目的问题容易混在一起，统计分析没有项目边界。",
        "现场与管理层之间信息割裂，关闭也常常没有证据。"
    ], 17)
    add_panel(slide, Inches(6.38), Inches(1.88), Inches(5.9), Inches(4.75), fill=(244, 249, 252))
    ppt_text(slide, Inches(6.68), Inches(2.08), Inches(5.2), Inches(0.35), "升级后的目标", 21, PALETTE["ink"], True)
    ppt_bullets(slide, Inches(6.64), Inches(2.46), Inches(5.1), Inches(3.8), [
        "小程序负责快速提报和现场跟进，尽量少打字，多媒体先行。",
        "Web 负责项目问题库、责任分派、状态推进和管理分析。",
        "每条问题必须绑定项目、责任人、时间线和关闭证据。",
        "整套系统围绕“项目优先、闭环可追溯”设计。",
    ], 16)

    # 3 architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 3, "总体架构与协同方式", "微信小程序 + Web 管理端 + 统一后端 + MySQL/Redis/MinIO")
    for left, title, desc in [
        (0.85, "微信小程序", "现场录入\n项目工作台\n项目问题库\n跟进 / 关闭"),
        (3.55, "Web 管理端", "驾驶舱\n项目问题库\n问题详情\n用户 / 统计"),
        (6.25, "统一后端", "认证与绑定\nRBAC 权限\n状态机\n附件 / 统计"),
        (9.15, "数据层", "MySQL\nRedis\nMinIO\nDocker Compose"),
    ]:
        add_panel(slide, Inches(left), Inches(2.15), Inches(2.1), Inches(2.2), fill=(255, 255, 255))
        ppt_text(slide, Inches(left + 0.18), Inches(2.42), Inches(1.7), Inches(0.32), title, 20, PALETTE["ink"], True, PP_ALIGN.CENTER)
        ppt_text(slide, Inches(left + 0.18), Inches(2.86), Inches(1.7), Inches(1.1), desc, 15, PALETTE["muted"], False, PP_ALIGN.CENTER)
    ppt_text(slide, Inches(0.92), Inches(5.08), Inches(11.4), Inches(1.05), "关键原则：状态流转、权限校验、编号生成、统计口径全部收口在后端；前端只负责录入、查看和推进。", 18, PALETTE["ink"], True, PP_ALIGN.CENTER)

    # 4 web map
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 4, "Web 管理端：页面总览与定位")
    add_image(slide, "web_dashboard.png", Inches(0.78), Inches(1.62), width=Inches(7.25))
    add_panel(slide, Inches(8.25), Inches(1.62), Inches(4.1), Inches(4.95), fill=(255, 255, 255))
    ppt_text(slide, Inches(8.5), Inches(1.88), Inches(3.4), Inches(0.3), "核心页面", 21, PALETTE["ink"], True)
    ppt_bullets(slide, Inches(8.46), Inches(2.22), Inches(3.45), Inches(4.0), [
        "登录页：账号密码登录后台。",
        "驾驶舱：按项目查看总览，未选项目时自动轮巡。",
        "项目选择器：先选项目，再进项目问题库。",
        "项目问题库：筛选、录入、查看问题。",
        "问题详情：附件预览、分派、推进、关闭。",
        "项目管理：项目主数据与项目团队。",
        "用户管理：企业账号、微信绑定、导入。",
        "统计分析与字典管理。",
    ], 15)

    # 5 dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 5, "Web 驾驶舱：管理层的一号页面")
    add_image(slide, "web_dashboard.png", Inches(0.72), Inches(1.7), width=Inches(7.8))
    add_panel(slide, Inches(8.7), Inches(1.7), Inches(3.65), Inches(4.9))
    ppt_bullets(slide, Inches(8.92), Inches(1.98), Inches(3.15), Inches(4.3), [
        "顶部项目选择器：手动切换到单项目视图。",
        "若不选择项目，系统自动轮巡各项目，适合大屏场景。",
        "KPI 卡显示问题总数、未关闭、高优先级、超期。",
        "趋势图和重点问题帮助管理层快速找到要先处理的项目与问题。",
        "点击项目或重点问题即可继续深入。"
    ], 15)

    # 6 issue library
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 6, "Web 项目问题库：项目优先 + 直接录入")
    add_image(slide, "web_issue_library.png", Inches(0.72), Inches(1.7), width=Inches(7.8))
    add_panel(slide, Inches(8.7), Inches(1.7), Inches(3.65), Inches(4.9))
    ppt_bullets(slide, Inches(8.92), Inches(1.98), Inches(3.15), Inches(4.25), [
        "所有问题都先归属项目，避免不同项目混池。",
        "顶部统计显示整项目口径，不受分页影响。",
        "支持关键字、责任人、状态、优先级等筛选。",
        "Web 端现在也能录入新问题，不必只靠小程序。",
        "图片和视频附件可直接随问题一起上传。"
    ], 15)

    # 7 issue detail
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 7, "Web 问题详情：责任、证据、进度在一页完成")
    add_image(slide, "web_issue_detail.png", Inches(0.7), Inches(1.7), width=Inches(7.85))
    add_panel(slide, Inches(8.7), Inches(1.7), Inches(3.65), Inches(4.9))
    ppt_bullets(slide, Inches(8.92), Inches(1.98), Inches(3.15), Inches(4.25), [
        "打开问题后先看：当前阶段、责任人、下一步、最新动作。",
        "附件在描述区下方直接预览，图片和视频都能点开查看。",
        "右侧直接做分派、优先级调整和状态推进。",
        "关闭时自动记录关闭人、关闭时间，并支持关闭证据。",
        "时间线集中展示系统日志和人工跟进。"
    ], 15)

    # 8 project team
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 8, "项目管理与项目团队：岗位不等于系统权限")
    add_image(slide, "web_project_team.png", Inches(0.76), Inches(1.7), width=Inches(7.8))
    add_panel(slide, Inches(8.72), Inches(1.7), Inches(3.6), Inches(4.9))
    ppt_bullets(slide, Inches(8.94), Inches(1.98), Inches(3.05), Inches(4.25), [
        "项目页维护项目编号、名称、客户、项目经理和计划日期。",
        "项目团队决定谁属于这个项目，谁能被分派问题。",
        "每个成员可配置项目岗位，以及可分派/可验证/可关闭权限。",
        "支持从现有 OPL 同步成员，也支持批量导入。",
        "这样责任人下拉不会出现不属于该项目的人。"
    ], 15)

    # 9 user management
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 9, "用户体系：企业账号为主，微信只做绑定和免密入口")
    add_image(slide, "web_user_mgmt.png", Inches(0.72), Inches(1.7), width=Inches(7.85))
    add_panel(slide, Inches(8.74), Inches(1.7), Inches(3.6), Inches(4.9))
    ppt_bullets(slide, Inches(8.96), Inches(1.98), Inches(3.02), Inches(4.25), [
        "管理员先创建企业账号、角色和部门。",
        "首次小程序登录时进行企业账号绑定。",
        "若仍是初始密码，系统强制改密。",
        "后续小程序可微信免密进入，Web 端继续账号密码登录。",
        "避免纯微信昵称登录带来的身份混乱。"
    ], 15)

    # 10 stats
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 10, "统计分析与字典管理：把项目复盘和标准口径做实")
    add_image(slide, "web_stats.png", Inches(0.72), Inches(1.72), width=Inches(7.85))
    add_panel(slide, Inches(8.72), Inches(1.72), Inches(3.62), Inches(4.9))
    ppt_bullets(slide, Inches(8.94), Inches(2.0), Inches(3.06), Inches(4.2), [
        "统计页看趋势、状态结构、项目压力和重点问题。",
        "驾驶舱偏总览，统计页偏分析与复盘。",
        "字典管理负责问题分类、来源、模块、部门等标准化。",
        "没有统一字典，后续筛选和统计会越来越失真。"
    ], 15)

    # 11 mini map
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 11, "小程序：现场轻操作的完整页面体系")
    add_image(slide, "mini_home.png", Inches(0.72), Inches(1.65), width=Inches(2.15))
    add_image(slide, "mini_selector.png", Inches(3.0), Inches(1.65), width=Inches(2.15))
    add_image(slide, "mini_create.png", Inches(5.28), Inches(1.65), width=Inches(2.15))
    add_image(slide, "mini_library.png", Inches(7.56), Inches(1.65), width=Inches(2.15))
    add_image(slide, "mini_detail.png", Inches(9.84), Inches(1.65), width=Inches(2.15))
    add_panel(slide, Inches(0.72), Inches(6.05), Inches(11.3), Inches(0.7), fill=(255, 255, 255))
    ppt_text(slide, Inches(1.0), Inches(6.24), Inches(10.5), Inches(0.24), "登录、首页、项目选择、项目工作台、项目统计、新建问题、项目问题库、问题详情、我的问题、待我处理、个人中心。", 15, PALETTE["ink"], True, PP_ALIGN.CENTER)

    # 12 mini project flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 12, "小程序主链路：先选项目，再录入 / 查看 / 跟进")
    add_image(slide, "mini_selector.png", Inches(0.84), Inches(1.72), width=Inches(2.2))
    add_image(slide, "mini_home.png", Inches(3.3), Inches(1.72), width=Inches(2.2))
    add_image(slide, "mini_library.png", Inches(5.76), Inches(1.72), width=Inches(2.2))
    add_image(slide, "mini_detail.png", Inches(8.22), Inches(1.72), width=Inches(2.2))
    ppt_text(slide, Inches(11.05), Inches(2.05), Inches(1.15), Inches(0.4), "1", 28, PALETTE["ink"], True, PP_ALIGN.CENTER)
    ppt_text(slide, Inches(11.05), Inches(2.72), Inches(1.15), Inches(0.4), "2", 28, PALETTE["ink"], True, PP_ALIGN.CENTER)
    ppt_text(slide, Inches(11.05), Inches(3.39), Inches(1.15), Inches(0.4), "3", 28, PALETTE["ink"], True, PP_ALIGN.CENTER)
    ppt_text(slide, Inches(11.05), Inches(4.06), Inches(1.15), Inches(0.4), "4", 28, PALETTE["ink"], True, PP_ALIGN.CENTER)
    ppt_bullets(slide, Inches(10.45), Inches(2.02), Inches(1.6), Inches(3.2), [
        "选项目",
        "进工作台",
        "看项目问题库",
        "进入详情做跟进/关闭"
    ], 14)

    # 13 mini create
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 13, "小程序录入新问题：现场员工的核心动作")
    add_image(slide, "mini_create.png", Inches(0.82), Inches(1.78), width=Inches(2.45))
    add_panel(slide, Inches(3.62), Inches(1.78), Inches(8.65), Inches(4.8))
    ppt_bullets(slide, Inches(3.9), Inches(2.08), Inches(8.0), Inches(4.15), [
        "当前项目、录入人、录入时间自动带入。",
        "录入时填写标题、描述、优先级、影响等级，尽量减少现场自由发挥。",
        "支持照片和视频；单文件 5MB、视频 15 秒，可压缩后上传。",
        "图片和视频会直接成为问题附件，管理层打开问题即可看到现场。",
        "这页的目标是让现场员工随时拿起手机就能报问题，而不是回办公室补 Excel。"
    ], 16)

    # 14 mini library
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 14, "小程序项目问题库：管理层也能在手机上看清项目 OPL")
    add_image(slide, "mini_library.png", Inches(0.82), Inches(1.78), width=Inches(2.45))
    add_panel(slide, Inches(3.62), Inches(1.78), Inches(8.65), Inches(4.8))
    ppt_bullets(slide, Inches(3.9), Inches(2.08), Inches(8.0), Inches(4.15), [
        "顶部显示整项目口径统计，不会随着翻页而改变。",
        "状态卡已经压成双列紧凑布局，适合现场手机查看。",
        "支持状态筛选、责任人筛选和关键字搜索。",
        "重点信息不是堆字段，而是看责任人、当前进度和下一步。",
        "管理层在现场也能快速看某个项目的全部 OPL。"
    ], 16)

    # 15 mini detail
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 15, "小程序问题详情：附件预览、跟进提交、关闭证据")
    add_image(slide, "mini_detail.png", Inches(0.82), Inches(1.78), width=Inches(2.45))
    add_image(slide, "mini_profile.png", Inches(9.72), Inches(1.78), width=Inches(2.2))
    add_panel(slide, Inches(3.6), Inches(1.78), Inches(5.78), Inches(4.8))
    ppt_bullets(slide, Inches(3.86), Inches(2.08), Inches(5.2), Inches(4.15), [
        "问题描述下方直接预览图片和视频，小图点开可中框预览。",
        "跟进记录支持文本 + 图片 + 视频。",
        "如果当前动作是关闭问题，可以上传关闭证据附件。",
        "关闭后自动显示关闭人、关闭时间、关闭说明和关闭证据。",
        "这页既是现场跟进页，也是移动端证据页。"
    ], 16)

    # 16 usage by role
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 16, "怎么使用：按角色看最直接")
    for idx, data in enumerate([
        ("现场人员", ["选项目", "录新问题", "上传图片/视频", "后续补跟进"]),
        ("责任工程师 / 项目经理", ["看待我处理", "分派责任", "推进状态", "关闭并留证据"]),
        ("管理层 / 管理员", ["看驾驶舱", "看项目统计", "维护团队和用户", "复盘项目问题结构"]),
    ]):
        left = Inches(0.78 + idx * 4.15)
        add_panel(slide, left, Inches(1.9), Inches(3.85), Inches(4.9))
        ppt_text(slide, left + Inches(0.2), Inches(2.16), Inches(3.3), Inches(0.3), data[0], 22, PALETTE["ink"], True)
        ppt_bullets(slide, left + Inches(0.18), Inches(2.62), Inches(3.25), Inches(3.6), data[1], 16)

    # 17 deployment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 17, "部署与上线注意事项")
    add_panel(slide, Inches(0.82), Inches(1.74), Inches(5.7), Inches(4.9))
    ppt_text(slide, Inches(1.08), Inches(2.02), Inches(4.8), Inches(0.3), "部署方式", 22, PALETTE["ink"], True)
    ppt_bullets(slide, Inches(1.02), Inches(2.42), Inches(4.9), Inches(3.7), [
        "Docker Compose 一键启动：MySQL / Redis / MinIO / Backend / Web。",
        "所有服务统一 UTF-8、Asia/Shanghai 和 ISO 8601。",
        "MinIO 后续可平滑替换为 OSS / COS。",
        "支持本地环境、内网环境和私有化迁移。"
    ], 16)
    add_panel(slide, Inches(6.72), Inches(1.74), Inches(5.6), Inches(4.9))
    ppt_text(slide, Inches(6.98), Inches(2.02), Inches(4.8), Inches(0.3), "小程序真机前要做的事", 22, PALETTE["ink"], True)
    ppt_bullets(slide, Inches(6.92), Inches(2.42), Inches(4.9), Inches(3.7), [
        "接口地址必须切换到手机可访问的 HTTPS 域名。",
        "微信公众平台配置 request / upload / download 合法域名。",
        "正式环境关闭 mock，配置 AppID / Secret。",
        "附件预览和上传链路需按正式域名验证。"
    ], 16)

    # 18 close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_brand(slide, 18, True)
    ppt_text(slide, Inches(1.1), Inches(1.55), Inches(8.6), Inches(0.7), "最终价值", 18, PALETTE["cyan"], True)
    ppt_text(slide, Inches(1.1), Inches(2.18), Inches(9.4), Inches(1.2), "把 Excel 问题表，变成一个真正可协同、可追责、可关闭、可复盘的企业现场问题系统", 28, PALETTE["white"], True)
    ppt_bullets(slide, Inches(1.16), Inches(3.7), Inches(8.2), Inches(2.1), [
        "项目内问题分库，避免跨项目混乱。",
        "现场图像化录入，管理层打开问题即可理解现场。",
        "责任清晰、状态明确、关闭有证据。",
        "Web 和小程序共用一套后端与数据，适合持续迭代。"
    ], 17)

    return prs


def main():
    render_all_assets()
    prs = build_slides()
    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    main()
