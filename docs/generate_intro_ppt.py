from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_FILE = Path(__file__).resolve().parent / "非标设备项目现场问题协同系统-系统介绍与使用说明.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR = {
    "ink": RGBColor(24, 34, 43),
    "muted": RGBColor(96, 107, 120),
    "subtle": RGBColor(139, 150, 162),
    "paper": RGBColor(249, 246, 241),
    "panel": RGBColor(255, 252, 248),
    "panel_alt": RGBColor(245, 249, 248),
    "line": RGBColor(225, 220, 212),
    "orange": RGBColor(203, 93, 31),
    "orange_soft": RGBColor(251, 233, 218),
    "teal": RGBColor(22, 123, 120),
    "teal_soft": RGBColor(224, 243, 241),
    "blue": RGBColor(57, 93, 132),
    "blue_soft": RGBColor(230, 238, 247),
    "red": RGBColor(176, 78, 62),
    "red_soft": RGBColor(250, 230, 226),
    "dark": RGBColor(18, 25, 34),
    "dark_2": RGBColor(36, 46, 58),
    "white": RGBColor(255, 255, 255),
}

FONT_HEAD = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Bahnschrift"


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_full_bg(slide, color) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_blob(slide, left, top, width, height, color, transparency=0.78):
    blob = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
    blob.fill.solid()
    blob.fill.fore_color.rgb = color
    blob.fill.transparency = transparency
    blob.line.fill.background()
    return blob


def add_line(slide, left, top, width, height, color, transparency=0):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.fill.transparency = transparency
    line.line.fill.background()
    return line


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=18,
    color=COLOR["ink"],
    bold=False,
    font_name=FONT_BODY,
    align=PP_ALIGN.LEFT,
    vertical=MSO_VERTICAL_ANCHOR.TOP,
    margins=(6, 6, 4, 4),
):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Pt(margins[0])
    text_frame.margin_right = Pt(margins[1])
    text_frame.margin_top = Pt(margins[2])
    text_frame.margin_bottom = Pt(margins[3])
    text_frame.vertical_anchor = vertical
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(shape, items, font_size=15, color=COLOR["ink"], bullet_color=None):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(8)
    text_frame.margin_right = Pt(8)
    text_frame.margin_top = Pt(8)
    text_frame.margin_bottom = Pt(8)
    bullet_color = bullet_color or color
    for idx, item in enumerate(items):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run1 = p.add_run()
        run1.text = "• "
        run1.font.name = FONT_BODY
        run1.font.size = Pt(font_size)
        run1.font.bold = True
        run1.font.color.rgb = bullet_color
        run2 = p.add_run()
        run2.text = item
        run2.font.name = FONT_BODY
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color


def add_title_block(slide, kicker, title, subtitle="", dark=False):
    kicker_color = COLOR["orange"] if not dark else RGBColor(255, 194, 159)
    title_color = COLOR["ink"] if not dark else COLOR["white"]
    subtitle_color = COLOR["muted"] if not dark else RGBColor(212, 220, 229)
    add_textbox(slide, Inches(0.65), Inches(0.45), Inches(5.9), Inches(0.35), kicker.upper(), 11, kicker_color, False, FONT_MONO)
    add_textbox(slide, Inches(0.62), Inches(0.78), Inches(8.5), Inches(0.72), title, 28, title_color, True)
    if subtitle:
        add_textbox(slide, Inches(0.64), Inches(1.35), Inches(8.8), Inches(0.55), subtitle, 12.5, subtitle_color, False)


def add_footer(slide, page_num, label="WeCHAT OPL Platform"):
    add_line(slide, Inches(0.65), Inches(7.07), Inches(12.0), Pt(1.2), COLOR["line"])
    add_textbox(slide, Inches(0.65), Inches(7.08), Inches(4.0), Inches(0.22), label, 9.5, COLOR["subtle"], False, FONT_MONO)
    add_textbox(slide, Inches(12.1), Inches(7.02), Inches(0.45), Inches(0.25), str(page_num), 9.5, COLOR["subtle"], False, FONT_MONO, PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body="", accent="orange", fill=None):
    fill = fill or COLOR["panel"]
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = COLOR["line"]
    card.line.width = Pt(1)
    accent_map = {
        "orange": COLOR["orange"],
        "teal": COLOR["teal"],
        "blue": COLOR["blue"],
        "red": COLOR["red"],
    }
    accent_color = accent_map.get(accent, COLOR["orange"])
    add_line(slide, left + Pt(6), top + Pt(8), Pt(5), height - Pt(16), accent_color, 0.08)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.26), Inches(0.32), title, 18, COLOR["ink"], True)
    if body:
        add_textbox(slide, left + Inches(0.16), top + Inches(0.48), width - Inches(0.28), height - Inches(0.58), body, 12.5, COLOR["muted"])
    return card


def add_tag(slide, left, top, text, tone="orange", width=None):
    tone_map = {
        "orange": (COLOR["orange_soft"], COLOR["orange"]),
        "teal": (COLOR["teal_soft"], COLOR["teal"]),
        "blue": (COLOR["blue_soft"], COLOR["blue"]),
        "red": (COLOR["red_soft"], COLOR["red"]),
        "dark": (COLOR["dark_2"], COLOR["white"]),
    }
    bg, fg = tone_map[tone]
    width = width or Inches(1.2)
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34))
    tag.fill.solid()
    tag.fill.fore_color.rgb = bg
    tag.line.fill.background()
    add_textbox(slide, left, top + Inches(0.01), width, Inches(0.24), text, 10.5, fg, True, FONT_HEAD, PP_ALIGN.CENTER)


def add_three_kpis(slide, cards):
    left = Inches(0.72)
    top = Inches(1.95)
    gap = Inches(0.18)
    width = Inches(4.05)
    height = Inches(1.3)
    for idx, item in enumerate(cards):
        add_card(slide, left + idx * (width + gap), top, width, height, item["title"], item["body"], item.get("accent", "orange"))


def add_two_column_section(slide, left_title, left_items, right_title, right_items, top=Inches(1.95)):
    left_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.68), top, Inches(6.05), Inches(4.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR["panel"]
    left_card.line.color.rgb = COLOR["line"]
    right_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), top, Inches(5.75), Inches(4.7))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR["panel_alt"]
    right_card.line.color.rgb = COLOR["line"]
    add_textbox(slide, Inches(0.88), top + Inches(0.18), Inches(5.6), Inches(0.3), left_title, 19, COLOR["ink"], True)
    add_textbox(slide, Inches(7.1), top + Inches(0.18), Inches(5.2), Inches(0.3), right_title, 19, COLOR["ink"], True)
    left_box = slide.shapes.add_textbox(Inches(0.86), top + Inches(0.52), Inches(5.62), Inches(3.88))
    add_bullets(left_box, left_items, 15)
    right_box = slide.shapes.add_textbox(Inches(7.08), top + Inches(0.52), Inches(5.18), Inches(3.88))
    add_bullets(right_box, right_items, 15)


def add_page_grid(slide, items, cols=3, top=Inches(1.95), card_height=Inches(1.28)):
    left = Inches(0.68)
    gap_x = Inches(0.16)
    gap_y = Inches(0.16)
    total_width = Inches(12.0)
    card_width = (total_width - gap_x * (cols - 1)) / cols
    for idx, item in enumerate(items):
        row = idx // cols
        col = idx % cols
        card_left = left + col * (card_width + gap_x)
        card_top = top + row * (card_height + gap_y)
        fill = COLOR["panel"] if row % 2 == 0 else COLOR["panel_alt"]
        add_card(slide, card_left, card_top, card_width, card_height, item["title"], item["body"], item.get("accent", "orange"), fill)


def add_flow(slide, steps, top=Inches(2.1)):
    left = Inches(0.72)
    box_w = Inches(2.0)
    box_h = Inches(0.95)
    gap = Inches(0.22)
    for idx, step in enumerate(steps):
        x = left + idx * (box_w + gap)
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR["panel"]
        box.line.color.rgb = step.get("line", COLOR["line"])
        box.line.width = Pt(1.25)
        add_textbox(slide, x + Inches(0.08), top + Inches(0.16), box_w - Inches(0.16), Inches(0.24), step["title"], 17, COLOR["ink"], True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.08), top + Inches(0.46), box_w - Inches(0.16), Inches(0.24), step["desc"], 10.5, COLOR["muted"], False, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, arrow_x + Inches(0.02), top + Inches(0.31), Inches(0.16), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR["orange"]
            arrow.line.fill.background()


def add_role_band(slide, left, top, width, title, desc, tone, bullets):
    add_card(slide, left, top, width, Inches(4.3), title, desc, tone)
    box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(1.2), width - Inches(0.28), Inches(2.8))
    add_bullets(box, bullets, 13.5)


def add_arch_box(slide, left, top, width, height, title, subtitle, tone="orange"):
    fill_map = {
        "orange": COLOR["orange_soft"],
        "teal": COLOR["teal_soft"],
        "blue": COLOR["blue_soft"],
        "red": COLOR["red_soft"],
    }
    line_map = {
        "orange": COLOR["orange"],
        "teal": COLOR["teal"],
        "blue": COLOR["blue"],
        "red": COLOR["red"],
    }
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_map[tone]
    box.line.color.rgb = line_map[tone]
    box.line.width = Pt(1.4)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.13), width - Inches(0.2), Inches(0.24), title, 16, COLOR["ink"], True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.42), width - Inches(0.2), height - Inches(0.48), subtitle, 11.2, COLOR["muted"], False, align=PP_ALIGN.CENTER)


def add_step_column(slide, left, top, width, title, steps, tone="orange"):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(4.7))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR["panel"]
    card.line.color.rgb = COLOR["line"]
    add_textbox(slide, left + Inches(0.16), top + Inches(0.15), width - Inches(0.32), Inches(0.3), title, 19, COLOR["ink"], True)
    start_y = top + Inches(0.58)
    for idx, step in enumerate(steps, start=1):
        y = start_y + Inches(0.84) * (idx - 1)
        add_tag(slide, left + Inches(0.16), y, f"STEP {idx}", tone, Inches(0.9))
        add_textbox(slide, left + Inches(1.12), y, width - Inches(1.28), Inches(0.55), step, 13.4, COLOR["muted"])


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["dark"])
    add_blob(slide, Inches(8.6), Inches(-0.8), Inches(4.4), Inches(4.4), COLOR["teal"], 0.78)
    add_blob(slide, Inches(-0.9), Inches(4.5), Inches(3.8), Inches(3.8), COLOR["orange"], 0.84)
    add_line(slide, Inches(0.74), Inches(0.86), Inches(1.6), Pt(3), COLOR["orange"])
    add_textbox(slide, Inches(0.78), Inches(1.12), Inches(4.2), Inches(0.34), "FIELD ISSUE OPS / SYSTEM PRESENTATION", 11.5, RGBColor(255, 199, 168), False, FONT_MONO)
    add_textbox(slide, Inches(0.78), Inches(1.62), Inches(9.4), Inches(1.25), "非标设备项目现场问题协同系统", 30, COLOR["white"], True)
    add_textbox(slide, Inches(0.82), Inches(2.72), Inches(7.8), Inches(0.8), "从 Excel 问题表，升级为项目优先、图像驱动、责任清晰、闭环可追溯的企业级问题协同平台", 15, RGBColor(218, 226, 234))
    add_tag(slide, Inches(0.82), Inches(3.68), "项目优先", "orange", Inches(1.05))
    add_tag(slide, Inches(1.98), Inches(3.68), "现场录入", "teal", Inches(1.05))
    add_tag(slide, Inches(3.14), Inches(3.68), "闭环管理", "blue", Inches(1.05))
    add_tag(slide, Inches(4.3), Inches(3.68), "附件证据", "red", Inches(1.05))
    add_arch_box(slide, Inches(8.2), Inches(1.55), Inches(1.9), Inches(1.15), "微信小程序", "现场录入\n跟进 / 关闭", "orange")
    add_arch_box(slide, Inches(10.3), Inches(1.55), Inches(1.9), Inches(1.15), "Web 管理端", "项目 / 统计\n分派 / 分析", "teal")
    add_arch_box(slide, Inches(8.2), Inches(2.95), Inches(4.0), Inches(1.05), "统一后端 + MySQL / Redis / MinIO", "统一权限、状态机、附件、编号、统计口径", "blue")
    add_textbox(slide, Inches(0.82), Inches(6.68), Inches(5.2), Inches(0.26), "版本范围：当前可运行系统说明稿", 10.5, RGBColor(180, 190, 200), False, FONT_MONO)
    add_textbox(slide, Inches(10.48), Inches(6.68), Inches(1.7), Inches(0.26), "2026.03", 10.5, RGBColor(180, 190, 200), False, FONT_MONO, PP_ALIGN.RIGHT)


def build_problem(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_blob(slide, Inches(8.9), Inches(-0.7), Inches(4.5), Inches(4.5), COLOR["teal"], 0.9)
    add_blob(slide, Inches(-0.6), Inches(5.3), Inches(3.0), Inches(3.0), COLOR["orange"], 0.92)
    add_title_block(slide, "Why Change", "为什么必须从 Excel 升级", "这不是换一个录入界面，而是重构现场问题的传递方式。")
    add_three_kpis(slide, [
        {"title": "Excel 的真实问题", "body": "文字描述不足、现场情况还要反复追问、状态和责任缺少统一口径。", "accent": "orange"},
        {"title": "管理层的真实痛点", "body": "打开问题后不能在 10 秒内理解本质，不知道影响、责任人、下一步。", "accent": "teal"},
        {"title": "项目团队的真实成本", "body": "责任归属不清、反复催办、超期问题不可视，统计也无法沉淀。", "accent": "blue"},
    ])
    add_two_column_section(
        slide,
        "本系统要解决什么",
        [
            "把每个问题从“文字行”变成“现场问题对象”：有项目、有责任、有附件、有时间线。",
            "让现场人员以拍照、拍视频、选项目、简单跟进为主，尽量减少纯文字输入。",
            "让项目经理和管理层按项目查看问题池，快速识别高优先级、超期、未分派问题。",
            "让问题关闭不再靠口头确认，而是保留关闭人、关闭时间、关闭说明和证据附件。",
        ],
        "交付给企业的直接结果",
        [
            "现场与管理层使用同一套数据，不再重复维护 Excel 和聊天记录。",
            "管理端可以做项目团队、责任人、优先级、统计分析和项目轮巡。",
            "小程序与 Web 共用账号体系、同一后端和数据库，便于后续私有化部署。",
            "整个问题链路具备审计能力，适合企业现场管理、追溯和复盘。",
        ],
        top=Inches(3.45),
    )
    add_footer(slide, num)


def build_roles(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Role Design", "角色分工不是一堆账号，而是现场协作关系")
    add_role_band(slide, Inches(0.72), Inches(1.95), Inches(3.9), "现场人员 / 工厂", "小程序轻操作，核心目标是随时记录真实现场。", "orange", [
        "选择项目后快速录入问题。",
        "上传照片、视频，自动带录入人和录入时间。",
        "查看我的问题、待我处理、当前项目 OPL。",
        "在问题详情里提交跟进记录和关闭证据。",
    ])
    add_role_band(slide, Inches(4.74), Inches(1.95), Inches(3.9), "责任工程师 / 项目经理", "Web 做重管理，小程序也能随时推进状态。", "teal", [
        "只在当前项目问题库内分派责任人。",
        "调整优先级、推进状态、补充处理说明。",
        "维护项目团队，限定谁可以分派、验证、关闭问题。",
        "按项目视角看清重点问题、下一步和风险。",
    ])
    add_role_band(slide, Inches(8.76), Inches(1.95), Inches(3.86), "管理层 / 管理员", "一个做运营看板，一个做基础账号和权限。", "blue", [
        "驾驶舱按项目轮巡或手动选择项目查看整体状态。",
        "统计页看趋势、结构、超期和高优先级压力。",
        "用户管理支持预创建账号、微信绑定、导入、重置密码。",
        "系统角色负责权限，项目团队负责项目内岗位和职责。",
    ])
    add_textbox(slide, Inches(0.78), Inches(6.55), Inches(12.0), Inches(0.32), "关键原则：权限角色控制“能做什么”，项目团队控制“在哪个项目做什么”。", 12.2, COLOR["muted"])
    add_footer(slide, num)


def build_architecture(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Architecture", "系统总体架构", "一套后端、一套数据库、两套前端入口，所有规则统一收口。")
    add_arch_box(slide, Inches(0.9), Inches(2.2), Inches(2.0), Inches(1.2), "微信小程序", "登录绑定\n项目选择\n问题录入\n跟进与关闭", "orange")
    add_arch_box(slide, Inches(0.9), Inches(4.0), Inches(2.0), Inches(1.2), "Web 管理端", "驾驶舱\n项目问题库\n问题详情\n用户与统计", "teal")
    add_arch_box(slide, Inches(4.1), Inches(2.15), Inches(4.1), Inches(3.15), "Spring Boot 后端", "认证与微信绑定\nRBAC 权限\nIssue 状态机\n附件服务\n项目团队\n统计计算", "blue")
    add_arch_box(slide, Inches(9.2), Inches(1.7), Inches(1.7), Inches(1.0), "MySQL", "业务主库", "orange")
    add_arch_box(slide, Inches(11.0), Inches(1.7), Inches(1.7), Inches(1.0), "Redis", "缓存 / 登录辅助", "teal")
    add_arch_box(slide, Inches(9.2), Inches(3.0), Inches(3.5), Inches(1.1), "MinIO", "图片 / 视频 / 证据附件", "blue")
    add_arch_box(slide, Inches(9.2), Inches(4.45), Inches(3.5), Inches(0.95), "Docker Compose", "MySQL + Redis + MinIO + Backend + Web", "red")
    add_line(slide, Inches(2.95), Inches(2.72), Inches(1.0), Pt(2), COLOR["orange"])
    add_line(slide, Inches(2.95), Inches(4.54), Inches(1.0), Pt(2), COLOR["teal"])
    add_line(slide, Inches(8.25), Inches(2.25), Inches(0.88), Pt(2), COLOR["blue"])
    add_line(slide, Inches(8.25), Inches(3.55), Inches(0.88), Pt(2), COLOR["blue"])
    add_line(slide, Inches(8.25), Inches(4.93), Inches(0.88), Pt(2), COLOR["blue"])
    add_textbox(slide, Inches(0.82), Inches(6.38), Inches(12.0), Inches(0.45), "为什么这套架构稳：前后端职责清晰、状态规则只在后端、附件存储可替换、部署可本地可私有化。", 12.5, COLOR["muted"])
    add_footer(slide, num)


def build_lifecycle(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Lifecycle", "问题闭环主链路", "系统的核心不是 CRUD，而是让问题形成可执行、可追溯、可关闭的流程。")
    add_flow(slide, [
        {"title": "新建 NEW", "desc": "项目 + 标题 + 描述 + 媒体"},
        {"title": "已受理", "desc": "明确责任人和处理入口"},
        {"title": "处理中", "desc": "跟进记录、附件、说明"},
        {"title": "待验证", "desc": "处理完成，等待确认"},
        {"title": "已关闭", "desc": "关闭人 / 时间 / 证据"},
    ], top=Inches(2.2))
    add_tag(slide, Inches(9.88), Inches(3.42), "可挂起", "blue", Inches(0.82))
    add_tag(slide, Inches(10.82), Inches(3.42), "可取消", "red", Inches(0.82))
    add_two_column_section(
        slide,
        "每次推进都记录什么",
        ["状态从哪里变到哪里。", "谁执行了这次操作。", "操作时间和备注。", "是否上传了处理附件、跟进附件或关闭证据。"],
        "关闭问题时必须关注什么",
        ["谁关闭了这条问题、关闭时间是什么时候。", "关闭说明可以补充处理结论。", "关闭证据可上传图片和视频，直接作为验收依据。", "关闭后仍保留时间线和附件审计，不丢失过程记录。"],
        top=Inches(3.95),
    )
    add_footer(slide, num)


def build_web_map(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web Pages", "Web 管理端页面总览", "这部分面向项目经理、管理层和管理员，强调重管理、重分析、重追踪。")
    items = [
        {"title": "登录页", "body": "账号密码登录，进入统一后台。", "accent": "orange"},
        {"title": "驾驶舱", "body": "全局或单项目 KPI、趋势、轮巡。", "accent": "teal"},
        {"title": "项目选择器", "body": "先选项目，再进入项目问题库。", "accent": "blue"},
        {"title": "项目问题库", "body": "筛选、分页、新建问题、整项目统计。", "accent": "orange"},
        {"title": "问题详情页", "body": "附件预览、分派、优先级、状态推进、时间线。", "accent": "teal"},
        {"title": "项目管理", "body": "项目主数据、项目团队、成员同步与导入。", "accent": "blue"},
        {"title": "统计分析", "body": "趋势、状态结构、项目压力和重点问题。", "accent": "red"},
        {"title": "用户管理", "body": "创建/导入账号、微信绑定、角色、参与项目。", "accent": "orange"},
        {"title": "字典管理", "body": "分类、来源、部门、模块等标准口径。", "accent": "teal"},
    ]
    add_page_grid(slide, items, cols=3, top=Inches(1.95), card_height=Inches(1.35))
    add_footer(slide, num)


def build_web_dashboard(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - Dashboard", "驾驶舱：管理层首先看到的页面")
    add_two_column_section(
        slide,
        "页面功能",
        ["顶部项目选择器：手动切换项目查看单项目数据。", "未选择项目时可自动轮巡，适合会议室大屏或管理巡检。", "展示总数、未关闭、高优先级、超期、趋势和重点问题。", "把“要先看什么项目、哪类问题最急”收敛成一屏信息。"],
        "怎么使用",
        ["1. 登录后直接进入驾驶舱。", "2. 选择具体项目，查看该项目的 KPI 和重点问题。", "3. 不选择项目时，系统自动轮巡各项目。", "4. 点重点问题或某个项目入口，进入项目问题库或问题详情。"],
        top=Inches(1.95),
    )
    add_card(slide, Inches(0.72), Inches(5.05), Inches(12.0), Inches(1.2), "适合谁", "管理层看节奏，项目经理看风险，会议场景看轮巡。", "orange")
    add_footer(slide, num)


def build_web_issue_library(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - Issue Library", "项目选择器 + 项目问题库", "从 Web 端录入、筛选、查看问题时，永远先锁定项目。")
    add_step_column(slide, Inches(0.72), Inches(1.95), Inches(3.86), "先选项目，再进问题库", [
        "通过项目选择器按项目编号、名称、客户筛选。",
        "进入项目后，顶部统计按整项目口径展示，不受分页影响。",
        "支持关键字、责任人、状态等条件筛选。",
        "列表支持直接进入问题详情。",
    ], "orange")
    add_step_column(slide, Inches(4.73), Inches(1.95), Inches(3.86), "在 Web 端录入新问题", [
        "问题库页面直接点击“新建问题”。",
        "当前项目自动带入，避免录错项目。",
        "可填写标题、描述、分类、优先级、设备/模块、责任人。",
        "支持上传图片和视频附件，作为问题现场证据。",
    ], "teal")
    add_step_column(slide, Inches(8.74), Inches(1.95), Inches(3.86), "这页解决了什么", [
        "把项目问题池和全局问题池彻底区分开。",
        "让问题筛选、统计、责任判断都围绕项目展开。",
        "让 Web 端不只是查看，也能补录问题。",
        "为项目经理提供最常用的管理入口。",
    ], "blue")
    add_footer(slide, num)


def build_web_issue_detail(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - Issue Detail", "问题详情页：最重要的一页", "打开后应在 10 秒内知道当前阶段、责任人、下一步和最新动作。")
    add_two_column_section(
        slide,
        "页面内能做什么",
        ["顶部直接看标题、优先级、当前阶段与描述。", "附件在描述区下方直接预览，图片和视频都能查看。", "右侧做责任分派、优先级调整和状态推进。", "时间线查看系统记录、人工跟进和附件评论。", "关闭时填写关闭说明、关闭证据，并记录关闭人和关闭时间。"],
        "如何实际使用",
        ["1. 先判断当前阶段、责任人和下一步是否正确。", "2. 若需改责任人，在此页直接分派，且只能分给当前项目团队成员。", "3. 若问题进入收尾，切到“关闭问题”，补说明和证据附件。", "4. 所有状态变更、评论和附件都会沉淀在时间线中。"],
        top=Inches(1.95),
    )
    add_card(slide, Inches(0.72), Inches(5.15), Inches(12.0), Inches(1.07), "为什么这页重要", "它把“描述、证据、责任、进度、关闭”五件事收敛到一处，是全系统的协同核心。", "red")
    add_footer(slide, num)


def build_web_project_team(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - Project Management", "项目管理与项目团队", "项目不是只有项目经理一个字段，而是一支有岗位、有权限的团队。")
    add_two_column_section(
        slide,
        "项目管理页包含什么",
        ["维护项目编号、名称、客户、项目经理、起止日期、项目状态。", "从项目页直接进入该项目的问题库。", "查看并维护“项目团队”。", "支持从现有 OPL 同步项目成员，减少手工维护成本。", "支持项目成员批量导入。"],
        "项目团队能表达什么",
        ["同一个人在不同项目内可以承担不同岗位。", "项目成员可以标记是否为项目经理、是否可分派、是否可验证、是否可关闭问题。", "问题责任人下拉只显示当前项目团队成员。", "岗位与系统权限分离，避免 RBAC 与现场职责混乱。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_web_user_mgmt(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - User Management", "用户管理、微信绑定与账号体系", "企业账号是主身份，微信只是登录入口和绑定关系。")
    add_step_column(slide, Inches(0.72), Inches(1.95), Inches(3.86), "管理员怎么做", [
        "预先创建企业账号，配置姓名、部门、角色。",
        "支持新建用户、编辑、批量导入、重置密码、解绑微信。",
        "查看用户参与了哪些项目，并可跳转项目团队维护。",
        "角色已扩展到项目经理、质量经理、机械工程师等岗位。",
    ], "orange")
    add_step_column(slide, Inches(4.73), Inches(1.95), Inches(3.86), "小程序如何登录", [
        "首次进入：先微信登录，再绑定一次企业账号和初始密码。",
        "如果还是初始密码，系统强制修改密码。",
        "绑定成功后，后续可直接微信免密进入。",
        "Web 端继续使用账号密码登录，适合后台管理场景。",
    ], "teal")
    add_step_column(slide, Inches(8.74), Inches(1.95), Inches(3.86), "这套设计的价值", [
        "避免纯微信昵称登录带来的重名和混乱。",
        "审计记录始终能追到企业账号和真实人员。",
        "离职、换岗、外协账号都能在后台统一管理。",
        "后续接 LDAP / AD / 企业微信也有空间。",
    ], "blue")
    add_footer(slide, num)


def build_web_stats_dict(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Web - Analytics", "统计分析与字典标准化")
    add_two_column_section(
        slide,
        "统计分析页",
        ["展示新增/关闭趋势、状态结构、项目压力和风险分布。", "适合管理层定期 review 项目健康度。", "和驾驶舱相比，这里更偏分析与结构化展示。", "可以快速识别哪个项目、哪个部门、哪类问题最集中。"],
        "字典管理页",
        ["统一问题分类、问题来源、模块、部门等标准口径。", "没有字典标准化，筛选和统计很快就会失真。", "适合在正式上线前由管理员和项目管理层共同维护。", "为后续报表、跨项目分析和导出做准备。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_mini_map(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp Pages", "小程序页面总览", "小程序负责轻操作：选项目、录问题、看项目 OPL、做跟进。")
    items = [
        {"title": "登录", "body": "微信登录、企业账号绑定、首次改密。", "accent": "orange"},
        {"title": "首页", "body": "主入口：项目工作台、录入、问题库、统计。", "accent": "teal"},
        {"title": "项目选择", "body": "所有项目级动作前先选项目。", "accent": "blue"},
        {"title": "项目工作台", "body": "紧凑入口 + 重点问题 + 项目概览。", "accent": "orange"},
        {"title": "项目统计", "body": "项目状态分布、优先级和重点问题。", "accent": "teal"},
        {"title": "新建问题", "body": "现场录入问题并上传照片/视频。", "accent": "red"},
        {"title": "项目问题库", "body": "按项目查看所有 OPL，并筛选责任人/状态。", "accent": "orange"},
        {"title": "问题详情", "body": "查看附件、跟进记录、关闭证据、推进状态。", "accent": "teal"},
        {"title": "我的问题", "body": "我提报的问题。", "accent": "blue"},
        {"title": "待我处理", "body": "分配给我的问题。", "accent": "orange"},
        {"title": "个人中心", "body": "绑定状态、版本、改密、退出登录。", "accent": "teal"},
    ]
    add_page_grid(slide, items, cols=4, top=Inches(1.95), card_height=Inches(1.18))
    add_footer(slide, num)


def build_mini_login_home(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Login & Home", "登录页与首页：让现场人员少思考，直接进入工作")
    add_two_column_section(
        slide,
        "登录页",
        ["Step 1：微信登录并识别当前微信身份。", "Step 2：若未绑定，则输入企业账号和初始密码完成绑定。", "Step 3：若仍使用初始密码，则强制修改密码后进入系统。", "整个过程围绕“企业身份稳定、微信体验轻量”设计。"],
        "首页",
        ["首页只保留高频入口，不堆太多说明文字。", "主入口是项目工作台、OPL 录入、项目问题库、项目统计。", "个人区包含我的问题、待我处理、个人中心。", "管理层账号可在首页看到简要项目快照。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_mini_project_hub(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Project Flow", "项目选择、项目工作台、项目统计", "项目优先是小程序最重要的交互原则。")
    add_step_column(slide, Inches(0.72), Inches(1.95), Inches(3.86), "1. 项目选择", [
        "按项目编号、名称、客户搜索。",
        "先锁定项目，再进入下一级菜单。",
        "避免跨项目问题混在一起。",
        "适合一个人同时参与多个项目的场景。",
    ], "orange")
    add_step_column(slide, Inches(4.73), Inches(1.95), Inches(3.86), "2. 项目工作台", [
        "首屏以紧凑布局展示项目 KPI 和四个主入口。",
        "入口包括 OPL 录入、项目问题库、项目统计、待我处理。",
        "下面附带重点问题，便于直接进入详情。",
        "适合项目经理和现场负责人随手查看。",
    ], "teal")
    add_step_column(slide, Inches(8.74), Inches(1.95), Inches(3.86), "3. 项目统计", [
        "管理层在小程序也能看项目状态分布和重点问题。",
        "包含全部 OPL、处理中、已关闭、已超期等指标。",
        "直接查看重点问题的责任人、截止时间和下一步。",
        "适合现场会议或巡检时快速查看项目健康度。",
    ], "blue")
    add_footer(slide, num)


def build_mini_create(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Create Issue", "新建问题页：现场录入主链路", "目标是现场员工可随时拍照、拍视频、快速提报。")
    add_two_column_section(
        slide,
        "录入时会自动带上的信息",
        ["当前项目：必须先选项目，避免问题挂错项目。", "录入人：来自当前登录用户。", "录入时间：系统自动记录。", "标题、描述、优先级、影响等级可手动填写。", "影响出货 / 影响联调 / 需要停机等开关可快速标记。"],
        "媒体上传规则",
        ["支持图片和视频。", "单文件限制 5MB，视频限制 15 秒。", "上传前可压缩，优先满足快速提报而不是高画质。", "附件提交后会和问题一起落库，作为现场证据。", "Web 端和小程序详情页都可以预览这些附件。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_mini_issue_library(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Project OPL", "项目问题库：小程序里的整项目视图")
    add_two_column_section(
        slide,
        "页面功能",
        ["顶部显示整项目口径的统计，不随分页变化。", "支持状态筛选、责任人筛选、关键字搜索。", "统计卡已经压成双列紧凑布局，减少页面高度。", "列表重点显示责任人、当前进度和下一步。"],
        "怎么使用",
        ["1. 进入项目工作台后点击“项目问题库”。", "2. 需要看全部就用默认状态；要看关闭、超期、处理中可直接切换筛选。", "3. 如需看某个人的处理情况，用责任人筛选。", "4. 点任意问题进入详情做跟进、查看附件或关闭。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_mini_detail(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Issue Detail", "问题详情：附件预览、跟进、关闭证据都在这里")
    add_two_column_section(
        slide,
        "这页能做什么",
        ["问题描述下方直接预览图片和视频缩略图。", "点击后弹出中框预览，大图查看或直接播放视频。", "时间线支持提交文字跟进，也支持图片和视频跟进。", "关闭问题时，可上传关闭证据图片/视频，并单独回显为关闭证据附件。", "页面同时显示当前阶段、责任人、下一步和关闭信息。"],
        "建议使用方式",
        ["现场处理过程中，多次补充跟进记录和附件。", "如果问题处理完成且进入关闭，先补关闭说明和关闭证据。", "管理层可直接在此页看证据，而不是再去聊天记录里找。", "这页既是移动端详情页，也是现场协同的证据页。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_mini_personal(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Miniapp - Personal Workspace", "我的问题、待我处理、个人中心")
    add_three_kpis(slide, [
        {"title": "我的问题", "body": "查看我提报的问题现在到哪一步，是否已经分派、是否超期。", "accent": "orange"},
        {"title": "待我处理", "body": "聚焦分配给我的问题，快速进入详情推进。", "accent": "teal"},
        {"title": "个人中心", "body": "查看账号、绑定状态、角色、版本，并修改密码。", "accent": "blue"},
    ])
    add_card(slide, Inches(0.72), Inches(3.65), Inches(12.0), Inches(2.05), "为什么这三页重要", "它们把小程序从“只会录问题”变成“有个人工作区”的现场工具：我提报的、分给我的、我的账号状态，都有独立入口。", "orange")
    add_footer(slide, num)


def build_rules(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Business Rules", "关键业务规则", "这套系统稳定运行，靠的不是页面多，而是规则统一。")
    items = [
        {"title": "项目优先", "body": "录入、查看、分派都先选项目，不做全局混用问题池。", "accent": "orange"},
        {"title": "项目团队约束", "body": "责任人只能从当前项目团队里选择。", "accent": "teal"},
        {"title": "权限角色 + 岗位角色", "body": "系统权限与岗位分离，避免角色体系混乱。", "accent": "blue"},
        {"title": "媒体规则", "body": "图片 / 视频单文件 5MB，视频 15 秒，支持压缩。", "accent": "red"},
        {"title": "关闭必须可追溯", "body": "关闭人、关闭时间、关闭说明、关闭证据全部保留。", "accent": "orange"},
        {"title": "整项目统计", "body": "统计看的是整个项目，不随分页变化。", "accent": "teal"},
        {"title": "状态机统一收口", "body": "NEW → ACCEPTED → IN_PROGRESS → PENDING_VERIFY → CLOSED。", "accent": "blue"},
        {"title": "附件统一落库", "body": "问题附件、跟进附件、关闭证据都可跨端查看。", "accent": "red"},
    ]
    add_page_grid(slide, items, cols=4, top=Inches(1.95), card_height=Inches(1.28))
    add_footer(slide, num)


def build_deploy(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Deployment", "部署与运维", "当前系统已按企业内网、本地部署、私有化迁移来设计。")
    add_two_column_section(
        slide,
        "部署组成",
        ["Docker Compose 一键启动：mysql / redis / minio / backend / web。", "所有服务统一 UTF-8、Asia/Shanghai、ISO 8601 时间格式。", "MinIO 可后续切换为 OSS / COS，不影响业务接口层。", "Web 与小程序共用后端与数据库。"],
        "小程序真机前还要做什么",
        ["把接口地址切换成可从手机访问的 HTTPS 域名。", "在微信公众平台配置 request / upload / download 合法域名。", "若需要真实微信登录，配置 AppID / Secret 并关闭本地 mock。", "开发者工具可本地调试，真机必须遵守域名和证书规则。"],
        top=Inches(1.95),
    )
    add_footer(slide, num)


def build_sop(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["paper"])
    add_title_block(slide, "Usage SOP", "不同角色怎么使用这套系统", "这页可以直接作为培训页使用。")
    add_page_grid(slide, [
        {"title": "现场人员", "body": "选项目 → 新建问题 → 上传媒体 → 后续补跟进。", "accent": "orange"},
        {"title": "责任工程师", "body": "看待我处理 → 处理问题 → 提交验证 → 补证据。", "accent": "teal"},
        {"title": "项目经理", "body": "进项目问题库 → 分派责任人 → 调优先级 → 跟进超期。", "accent": "blue"},
        {"title": "管理层", "body": "先看驾驶舱 / 项目统计 → 再看重点问题详情。", "accent": "red"},
        {"title": "管理员", "body": "建项目、建用户、导入项目成员、维护角色和字典。", "accent": "orange"},
        {"title": "关闭问题", "body": "补关闭说明和关闭证据，不建议直接口头关闭。", "accent": "teal"},
    ], cols=3, top=Inches(2.05), card_height=Inches(1.35))
    add_card(slide, Inches(0.72), Inches(5.22), Inches(12.0), Inches(1.0), "培训建议", "先培训“项目优先”和“关闭证据”这两件事，再培训高级统计和项目团队维护。", "blue")
    add_footer(slide, num)


def build_value(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLOR["dark"])
    add_blob(slide, Inches(-0.8), Inches(4.6), Inches(3.6), Inches(3.6), COLOR["orange"], 0.85)
    add_blob(slide, Inches(9.5), Inches(-0.6), Inches(4.0), Inches(4.0), COLOR["teal"], 0.84)
    add_title_block(slide, "Outcome", "这套系统最终给企业带来的不是“一个软件”，而是一种现场协同方式", "项目内问题分库、现场图像化录入、管理层快速理解、责任清晰、证据闭环。", dark=True)
    add_page_grid(slide, [
        {"title": "现场端更快", "body": "少打字，多拍图，多选项，随时提报。", "accent": "orange"},
        {"title": "管理端更清楚", "body": "看项目、看趋势、看重点问题，不再盯 Excel。", "accent": "teal"},
        {"title": "责任更明确", "body": "谁录入、谁处理、谁关闭，全链路可追溯。", "accent": "blue"},
        {"title": "后续更容易扩展", "body": "可继续做报表、导出、AI 分析、企业集成。", "accent": "red"},
    ], cols=2, top=Inches(2.35), card_height=Inches(1.45))
    add_textbox(slide, Inches(0.8), Inches(6.45), Inches(7.0), Inches(0.34), "建议汇报顺序：先讲痛点和闭环，再演示 Web 驾驶舱与小程序录入主链路。", 12, RGBColor(214, 221, 230))
    add_footer(slide, num, "Presentation Generated by Codex")


def build_presentation():
    prs = Presentation()
    set_slide_size(prs)
    build_cover(prs)
    build_problem(prs, 2)
    build_roles(prs, 3)
    build_architecture(prs, 4)
    build_lifecycle(prs, 5)
    build_web_map(prs, 6)
    build_web_dashboard(prs, 7)
    build_web_issue_library(prs, 8)
    build_web_issue_detail(prs, 9)
    build_web_project_team(prs, 10)
    build_web_user_mgmt(prs, 11)
    build_web_stats_dict(prs, 12)
    build_mini_map(prs, 13)
    build_mini_login_home(prs, 14)
    build_mini_project_hub(prs, 15)
    build_mini_create(prs, 16)
    build_mini_issue_library(prs, 17)
    build_mini_detail(prs, 18)
    build_mini_personal(prs, 19)
    build_rules(prs, 20)
    build_deploy(prs, 21)
    build_sop(prs, 22)
    build_value(prs, 23)
    return prs


def main():
    prs = build_presentation()
    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    main()
