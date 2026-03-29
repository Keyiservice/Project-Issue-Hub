from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
BASE_PPT = next(ROOT.glob("*第一版补图增强版-v2.pptx"))
OUT_PPT = ROOT / "非标设备项目现场问题协同系统-第一版补图增强版-v3.pptx"
REAL = ROOT / "ppt_real_assets"

INK = RGBColor(24, 34, 43)
MUTED = RGBColor(96, 107, 120)
SUBTLE = RGBColor(139, 150, 162)
PAPER = RGBColor(249, 246, 241)
PANEL = RGBColor(255, 252, 248)
LINE = RGBColor(225, 220, 212)
ORANGE = RGBColor(203, 93, 31)
TEAL = RGBColor(22, 123, 120)


def add_bg(slide, width, height):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.8), Inches(5.1), Inches(3.4), Inches(3.4))
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = ORANGE
    glow_left.fill.transparency = 0.93
    glow_left.line.fill.background()

    glow_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.1), Inches(-0.7), Inches(4.2), Inches(4.2))
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = TEAL
    glow_right.fill.transparency = 0.91
    glow_right.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        lead = p.add_run()
        lead.text = "• "
        lead.font.name = "Microsoft YaHei"
        lead.font.size = Pt(size)
        lead.font.bold = True
        lead.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = INK


def add_header(slide, page_num, title, subtitle):
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(0.56), Inches(0.5), Inches(0.56), Inches(0.56))
    arc.fill.background()
    arc.line.color.rgb = TEAL
    arc.line.width = Pt(3)

    add_text(slide, Inches(0.95), Inches(0.5), Inches(9.4), Inches(0.42), title, 26, INK, True)
    add_text(slide, Inches(0.98), Inches(0.94), Inches(9.8), Inches(0.35), subtitle, 12, MUTED)
    add_text(slide, Inches(10.78), Inches(0.32), Inches(1.65), Inches(0.22), "REAL SCREENS v3", 10.5, SUBTLE, False, PP_ALIGN.RIGHT)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(7.07), Inches(12.0), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    add_text(slide, Inches(0.65), Inches(7.08), Inches(5.0), Inches(0.22), "First Edition Enhanced v3", 9.5, SUBTLE)
    add_text(slide, Inches(12.0), Inches(7.02), Inches(0.4), Inches(0.22), str(page_num), 9.5, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    panel.line.width = Pt(1)
    return panel


def add_image_panel(slide, img_name, left, top, width, height):
    add_panel(slide, left, top, width, height)
    path = REAL / img_name
    if path.exists():
        slide.shapes.add_picture(str(path), left + Inches(0.08), top + Inches(0.08), width=width - Inches(0.16), height=height - Inches(0.16))


def add_single_image_slide(prs, page_num, title, subtitle, img_name, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle)
    add_image_panel(slide, img_name, Inches(0.72), Inches(1.55), Inches(7.45), Inches(5.18))
    add_panel(slide, Inches(8.35), Inches(1.55), Inches(3.98), Inches(5.18))
    add_bullets(slide, Inches(8.57), Inches(1.84), Inches(3.56), Inches(4.62), bullets, 14.6)


def add_dual_image_slide(prs, page_num, title, subtitle, left_img, right_img, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle)
    add_image_panel(slide, left_img, Inches(0.72), Inches(1.55), Inches(5.65), Inches(4.8))
    add_image_panel(slide, right_img, Inches(6.52), Inches(1.55), Inches(5.8), Inches(4.8))
    add_panel(slide, Inches(0.72), Inches(6.5), Inches(11.6), Inches(0.55))
    add_bullets(slide, Inches(0.92), Inches(6.61), Inches(11.2), Inches(0.32), bullets, 13.4)


def add_three_panel_slide(prs, page_num, title, subtitle, main_img, sub_img1, sub_img2, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle)
    add_image_panel(slide, main_img, Inches(0.72), Inches(1.55), Inches(4.2), Inches(5.15))
    add_image_panel(slide, sub_img1, Inches(5.1), Inches(1.55), Inches(3.12), Inches(2.48))
    add_image_panel(slide, sub_img2, Inches(5.1), Inches(4.22), Inches(3.12), Inches(2.48))
    add_panel(slide, Inches(8.42), Inches(1.55), Inches(3.9), Inches(5.15))
    add_bullets(slide, Inches(8.63), Inches(1.85), Inches(3.5), Inches(4.55), bullets, 14.2)


def main():
    prs = Presentation(str(BASE_PPT))
    start_page = len(prs.slides) + 1

    add_dual_image_slide(
        prs,
        start_page,
        "Web 真实补图：项目入口到驾驶舱",
        "管理层和项目经理最常走的使用路径。",
        "issue_projects_real.png",
        "dashboard_real.png",
        [
            "先选项目，再进入该项目自己的问题库或驾驶舱，这是系统当前明确收口的使用路径。",
            "项目选择页解决了“所有项目问题混在一起”的阅读问题，驾驶舱则负责看当前项目的节奏和风险。",
            "这两张图连起来以后，PPT 里“怎么进入系统主工作区”这件事会更清楚。"
        ]
    )

    add_single_image_slide(
        prs,
        start_page + 1,
        "小程序真实补图：运行中的开发现场",
        "不是示意图，是开发者工具里真实运行的小程序。",
        "miniapp_after_seed_reload.png",
        [
            "这页保留了真实调试场景，能说明小程序当前确实已经具备首页、项目入口、个人工作区和快照区。",
            "相比只放单张手机裁切图，这种页面更适合对内汇报时证明系统已经进入实际联调阶段。",
            "如果后续补真机截图，这一页可以直接替换成真机页而不用改文案逻辑。"
        ]
    )

    add_three_panel_slide(
        prs,
        start_page + 2,
        "小程序真实补图：首页入口拆解",
        "把高频使用区域拆开看，便于培训和介绍。",
        "miniapp_home_real.png",
        "miniapp_home_actions_real.png",
        "miniapp_home_lower_real.png",
        [
            "首页上半区承载项目级入口，用户先进入项目，再去录入 OPL、看问题库和项目统计。",
            "中下区承载个人工作区和快照区，方便现场人员和管理层各自看自己的入口。",
            "这组图不改颜色、不重绘，只是把真实首页拆解成更易讲解的几个区域。"
        ]
    )

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
