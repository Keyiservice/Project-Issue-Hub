from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET = ROOT / "ppt_corporate_assets"
REAL = ROOT / "ppt_real_assets"
OUT = ROOT / "非标设备项目现场问题协同系统-原图截图版.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(23, 48, 95)
MUTED = RGBColor(86, 103, 128)
TEAL = RGBColor(37, 171, 179)
LINE = RGBColor(220, 231, 243)
WHITE = RGBColor(255, 255, 255)


def crop_mini_login():
    src = REAL / "miniapp_devtools_foreground.png"
    out = REAL / "miniapp_login_crop.png"
    if src.exists() and not out.exists():
        img = Image.open(src)
        crop = img.crop((42, 92, 355, 650))
        crop.save(out)
    return out


def add_bg(slide, cover=False):
    name = "bg_cover.png" if cover else "bg_inner.png"
    slide.shapes.add_picture(str(ASSET / name), 0, 0, width=SLIDE_W, height=SLIDE_H)


def text(slide, left, top, width, height, value, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Microsoft YaHei"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = TEAL
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK


def panel(slide, left, top, width, height, fill=WHITE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def header(slide, page, title, subtitle=""):
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(0.56), Inches(0.5), Inches(0.56), Inches(0.56))
    arc.fill.background()
    arc.line.color.rgb = TEAL
    arc.line.width = Pt(3)
    text(slide, Inches(0.95), Inches(0.5), Inches(9.0), Inches(0.42), title, 26, INK, True)
    if subtitle:
        text(slide, Inches(0.98), Inches(0.95), Inches(9.3), Inches(0.35), subtitle, 12, MUTED)
    text(slide, Inches(10.9), Inches(0.32), Inches(1.8), Inches(0.24), "FIELD ISSUE OPS", 11, INK, True, PP_ALIGN.RIGHT)
    text(slide, Inches(12.0), Inches(7.02), Inches(0.3), Inches(0.2), str(page), 10, MUTED, False, PP_ALIGN.RIGHT)


def pic(slide, name, left, top, width=None, height=None):
    path = REAL / name
    if path.exists():
        kwargs = {}
        if width is not None:
            kwargs["width"] = width
        if height is not None:
            kwargs["height"] = height
        slide.shapes.add_picture(str(path), left, top, **kwargs)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    crop_mini_login()

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, True)
    text(s, Inches(1.08), Inches(1.1), Inches(3.0), Inches(0.3), "Bridge to Execution", 16, RGBColor(38, 214, 234), True)
    text(s, Inches(1.08), Inches(2.2), Inches(8.8), Inches(1.0), "非标设备项目现场问题协同系统", 31, WHITE, True)
    text(s, Inches(1.1), Inches(3.08), Inches(8.0), Inches(0.6), "原图截图版：保留真实系统界面，不对截图颜色做任何处理。", 16, RGBColor(220, 229, 238))
    text(s, Inches(1.1), Inches(6.22), Inches(2.2), Inches(0.26), "2026 / 03 / 25", 13, WHITE)
    text(s, Inches(10.9), Inches(0.32), Inches(1.8), Inches(0.24), "FIELD ISSUE OPS", 11, WHITE, True, PP_ALIGN.RIGHT)

    # 2 why
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 2, "为什么要从 Excel 升级成项目化问题协同系统")
    panel(s, Inches(0.78), Inches(1.55), Inches(12.0), Inches(5.5))
    bullets(s, Inches(1.02), Inches(1.92), Inches(5.4), Inches(4.8), [
        "Excel 只有文本，没有现场画面、责任状态和关闭证据。",
        "管理层打开问题后，往往还要追问“这是什么、影响多大、谁在处理”。",
        "不同项目的问题容易混在一起，责任归属和统计口径都会变形。",
        "关闭动作没有证据沉淀，后续复盘难以还原现场。",
        "系统升级的目标不是替代表格，而是建立统一的问题闭环。"
    ], 17)
    panel(s, Inches(6.5), Inches(1.86), Inches(5.94), Inches(4.8), RGBColor(244, 249, 252))
    text(s, Inches(6.8), Inches(2.08), Inches(5.0), Inches(0.3), "系统升级后的结果", 21, INK, True)
    bullets(s, Inches(6.76), Inches(2.46), Inches(5.0), Inches(4.0), [
        "问题先归属项目，再做录入、分派、跟进、关闭。",
        "现场以图片和视频为主，管理层打开问题即可看懂。",
        "Web 与小程序共用数据和规则，责任和状态统一。",
        "关闭问题必须留痕：关闭人、关闭时间、关闭证据。"
    ], 16)

    # 3 architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 3, "总体架构", "微信小程序 + Web 管理端 + 统一后端 + MySQL/Redis/MinIO")
    for idx, data in enumerate([
        ("微信小程序", "项目选择\n现场录入\n问题库\n跟进/关闭"),
        ("Web 管理端", "驾驶舱\n项目问题库\n问题详情\n用户/统计"),
        ("统一后端", "认证与绑定\nRBAC\n状态机\n附件/统计"),
        ("数据层", "MySQL\nRedis\nMinIO\nDocker"),
    ]):
        left = Inches(0.88 + idx * 2.95)
        panel(s, left, Inches(2.18), Inches(2.18), Inches(2.15))
        text(s, left + Inches(0.12), Inches(2.46), Inches(1.92), Inches(0.3), data[0], 20, INK, True, PP_ALIGN.CENTER)
        text(s, left + Inches(0.14), Inches(2.92), Inches(1.88), Inches(1.0), data[1], 15, MUTED, False, PP_ALIGN.CENTER)
    text(s, Inches(0.92), Inches(5.1), Inches(11.2), Inches(0.8), "关键原则：前端负责录入与查看，业务规则全部统一收口在后端，状态流转和权限不分散到页面里。", 18, INK, True, PP_ALIGN.CENTER)

    # 4 dashboard
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 4, "Web 驾驶舱：真实系统截图", "以下开始全部使用当前系统原始截图，不改颜色。")
    pic(s, "dashboard_real.png", Inches(0.72), Inches(1.55), width=Inches(8.15))
    panel(s, Inches(9.02), Inches(1.55), Inches(3.3), Inches(5.15))
    bullets(s, Inches(9.2), Inches(1.88), Inches(2.9), Inches(4.65), [
        "登录后默认进入驾驶舱。",
        "支持顶部项目选择器，也支持未选项目时自动轮巡。",
        "集中展示 KPI、趋势、重点问题和当前项目状态。",
        "适合管理层例会、项目经理巡检和大屏展示。"
    ], 15)

    # 5 issue library
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 5, "Web 项目问题库：真实系统截图")
    pic(s, "issues_real.png", Inches(0.72), Inches(1.55), width=Inches(8.15))
    panel(s, Inches(9.02), Inches(1.55), Inches(3.3), Inches(5.15))
    bullets(s, Inches(9.2), Inches(1.88), Inches(2.9), Inches(4.65), [
        "先选项目，再进入项目问题库。",
        "顶部统计是整项目口径，不跟翻页变化。",
        "支持关键字、责任人、状态、优先级等筛选。",
        "项目问题库页可直接录入新问题。"
    ], 15)

    # 6 issue detail
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 6, "Web 问题详情：真实系统截图")
    pic(s, "issue_detail_real.png", Inches(0.58), Inches(1.48), width=Inches(8.5))
    panel(s, Inches(9.2), Inches(1.55), Inches(3.1), Inches(5.15))
    bullets(s, Inches(9.36), Inches(1.88), Inches(2.78), Inches(4.65), [
        "打开问题后先看当前阶段、责任人、下一步。",
        "附件直接在描述区预览，支持图片和视频。",
        "右侧可以做分派、优先级调整和状态推进。",
        "关闭问题时填写关闭说明和关闭证据。"
    ], 15)

    # 7 projects
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 7, "Web 项目管理：真实系统截图")
    pic(s, "projects_real.png", Inches(0.72), Inches(1.55), width=Inches(8.15))
    panel(s, Inches(9.02), Inches(1.55), Inches(3.3), Inches(5.15))
    bullets(s, Inches(9.2), Inches(1.88), Inches(2.9), Inches(4.65), [
        "维护项目编号、名称、客户、项目经理和计划日期。",
        "项目团队用于表达项目内岗位和责任范围。",
        "问题责任人下拉只显示当前项目团队成员。",
        "支持同步 OPL 成员和团队批量导入。"
    ], 15)

    # 8 users
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 8, "Web 用户管理：真实系统截图")
    pic(s, "users_real.png", Inches(0.72), Inches(1.55), width=Inches(8.15))
    panel(s, Inches(9.02), Inches(1.55), Inches(3.3), Inches(5.15))
    bullets(s, Inches(9.2), Inches(1.88), Inches(2.9), Inches(4.65), [
        "管理员预创建企业账号与角色。",
        "支持新建、编辑、导入、重置密码、解绑微信。",
        "小程序首次登录时绑定企业账号，而不是直接用微信昵称。",
        "企业账号是主身份，微信只做登录入口。"
    ], 15)

    # 9 stats
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 9, "Web 统计分析：真实系统截图")
    pic(s, "stats_real.png", Inches(0.72), Inches(1.55), width=Inches(8.15))
    panel(s, Inches(9.02), Inches(1.55), Inches(3.3), Inches(5.15))
    bullets(s, Inches(9.2), Inches(1.88), Inches(2.9), Inches(4.65), [
        "管理层可按项目看趋势、结构和风险。",
        "统计页偏分析和复盘，驾驶舱偏总览和巡检。",
        "后续可继续扩展导出报表和跨项目对比。"
    ], 15)

    # 10 miniapp real
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 10, "小程序：当前已抓到的原始截图", "这里使用开发者工具中的真实小程序界面，不做颜色改动。")
    pic(s, "miniapp_login_crop.png", Inches(0.9), Inches(1.52), height=Inches(5.55))
    panel(s, Inches(4.15), Inches(1.55), Inches(8.1), Inches(5.15))
    bullets(s, Inches(4.36), Inches(1.9), Inches(7.7), Inches(4.5), [
        "当前已抓到的是小程序真实登录页原图。",
        "小程序登录方式是：微信登录 → 企业账号绑定 → 首次强制改密。",
        "这页也能看到当前开发环境的提示：真机前需切换到可访问的 HTTPS 域名。",
        "我没有再对小程序截图调色或重新绘制。"
    ], 16)

    # 11 miniapp usage
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 11, "小程序主链路怎么用", "先选项目，再录入问题、查看项目问题库、提交跟进和关闭证据。")
    panel(s, Inches(0.82), Inches(1.7), Inches(3.65), Inches(4.9))
    text(s, Inches(1.02), Inches(2.0), Inches(3.1), Inches(0.3), "现场员工", 22, INK, True)
    bullets(s, Inches(0.98), Inches(2.42), Inches(3.1), Inches(3.8), [
        "选择项目",
        "录入新问题",
        "上传现场图片/视频",
        "后续补跟进记录"
    ], 16)
    panel(s, Inches(4.82), Inches(1.7), Inches(3.65), Inches(4.9))
    text(s, Inches(5.02), Inches(2.0), Inches(3.1), Inches(0.3), "责任工程师 / 项目经理", 22, INK, True)
    bullets(s, Inches(4.98), Inches(2.42), Inches(3.1), Inches(3.8), [
        "查看项目问题库",
        "看待我处理",
        "提交处理进展",
        "关闭前补充证据"
    ], 16)
    panel(s, Inches(8.82), Inches(1.7), Inches(3.45), Inches(4.9))
    text(s, Inches(9.02), Inches(2.0), Inches(2.9), Inches(0.3), "管理层", 22, INK, True)
    bullets(s, Inches(8.98), Inches(2.42), Inches(2.9), Inches(3.8), [
        "看项目统计",
        "看项目问题库",
        "看重点问题和责任人",
        "看关闭证据"
    ], 16)

    # 12 business rules
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 12, "关键业务规则")
    panel(s, Inches(0.82), Inches(1.68), Inches(11.45), Inches(5.25))
    bullets(s, Inches(1.06), Inches(2.0), Inches(10.8), Inches(4.6), [
        "所有问题先归属项目，不允许全局问题池混用不同项目数据。",
        "责任人只能从当前项目团队中选择。",
        "关闭问题要保留关闭人、关闭时间、关闭说明和关闭证据。",
        "图片和视频是问题证据的一部分，不是附属备注。",
        "系统角色负责权限，项目团队负责项目内岗位与职责。",
        "统计看整项目口径，不跟分页变化。"
    ], 17)

    # 13 deploy
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 13, "部署与上线说明")
    panel(s, Inches(0.82), Inches(1.75), Inches(5.72), Inches(4.9))
    text(s, Inches(1.04), Inches(2.04), Inches(4.9), Inches(0.3), "部署组成", 22, INK, True)
    bullets(s, Inches(1.0), Inches(2.42), Inches(4.9), Inches(3.8), [
        "Docker Compose：mysql / redis / minio / backend / web",
        "统一 UTF-8、Asia/Shanghai、ISO 8601",
        "后端统一权限、状态机和附件服务",
        "支持内网与私有化部署"
    ], 16)
    panel(s, Inches(6.78), Inches(1.75), Inches(5.5), Inches(4.9))
    text(s, Inches(6.98), Inches(2.04), Inches(4.8), Inches(0.3), "小程序真机前必须完成", 22, INK, True)
    bullets(s, Inches(6.94), Inches(2.42), Inches(4.8), Inches(3.8), [
        "把接口地址切换成手机可访问的 HTTPS 域名",
        "在微信后台配置 request/upload/download 合法域名",
        "若启用真实微信登录，配置 AppID / Secret",
        "真机联调时重点验证附件上传与预览"
    ], 16)

    # 14 close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, True)
    text(s, Inches(1.1), Inches(1.2), Inches(3.0), Inches(0.3), "Bridge to Execution", 16, RGBColor(38, 214, 234), True)
    text(s, Inches(1.1), Inches(2.1), Inches(8.9), Inches(1.1), "把问题从一行 Excel，变成可协同、可追责、可关闭、可复盘的现场系统", 29, WHITE, True)
    bullets(s, Inches(1.15), Inches(3.65), Inches(8.1), Inches(2.2), [
        "项目内问题分库，边界清晰。",
        "现场图像化录入，管理层快速理解。",
        "责任人、状态、时间线和证据统一保留。",
        "Web 与小程序共用一套数据与规则。"
    ], 17)
    text(s, Inches(10.9), Inches(0.32), Inches(1.8), Inches(0.24), "FIELD ISSUE OPS", 11, WHITE, True, PP_ALIGN.RIGHT)
    text(s, Inches(12.0), Inches(7.02), Inches(0.3), Inches(0.2), "14", 10, RGBColor(214, 224, 238), False, PP_ALIGN.RIGHT)

    return prs


def main():
    prs = build()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
