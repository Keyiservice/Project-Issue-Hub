from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "ppt_real_assets"
OUT = ROOT / "Project-Issue-Hub-内网部署与小程序公网网关方案-领导IT版.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(24, 33, 43)
SLATE = RGBColor(79, 90, 103)
MUTED = RGBColor(123, 133, 144)
WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(247, 243, 236)
LINE = RGBColor(229, 223, 214)
TEAL = RGBColor(51, 136, 127)
ORANGE = RGBColor(211, 107, 41)
SOFT_TEAL = RGBColor(227, 242, 238)
SOFT_ORANGE = RGBColor(248, 233, 222)
SOFT_BLUE = RGBColor(231, 240, 250)
SOFT_RED = RGBColor(251, 235, 230)
RED = RGBColor(184, 74, 58)
GREEN = RGBColor(45, 126, 91)


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    grid = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.1), Inches(0.1), Inches(13.1), Inches(7.3))
    grid.fill.background()
    grid.line.color.rgb = RGBColor(234, 228, 220)
    grid.line.transparency = 0.5
    grid.line.width = Pt(0.5)

    glow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.3), Inches(-0.7), Inches(4.0), Inches(4.0))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = TEAL
    glow1.fill.transparency = 0.93
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.6), Inches(5.2), Inches(3.0), Inches(3.0))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ORANGE
    glow2.fill.transparency = 0.94
    glow2.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def card(slide, left, top, width, height, fill=WHITE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def title_block(slide, page_no, title, subtitle):
    add_text(slide, Inches(0.72), Inches(0.28), Inches(1.0), Inches(0.2), f"{page_no:02d}", 11, MUTED, True)
    add_text(slide, Inches(1.12), Inches(0.25), Inches(7.5), Inches(0.4), title, 26, NAVY, True)
    add_text(slide, Inches(1.15), Inches(0.72), Inches(9.8), Inches(0.24), subtitle, 12, SLATE)


def add_image(slide, filename, left, top, width, height):
    path = ASSETS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    left = card(slide, Inches(0.58), Inches(0.55), Inches(5.25), Inches(6.2), fill=NAVY)
    left.line.fill.background()
    add_text(slide, Inches(0.92), Inches(0.9), Inches(3.6), Inches(0.2), "PROJECT ISSUE HUB", 11, RGBColor(205, 213, 220))
    add_text(slide, Inches(0.92), Inches(1.28), Inches(4.3), Inches(1.25), "内网核心系统 +\n小程序公网网关方案", 28, WHITE, True)
    add_text(
        slide,
        Inches(0.92),
        Inches(3.1),
        Inches(4.2),
        Inches(1.15),
        "给领导和 IT 的统一说明版本\n目标是数据留内网、管理端走内网、小程序正式可上线。",
        16,
        RGBColor(224, 231, 236),
    )

    pill = card(slide, Inches(0.92), Inches(5.88), Inches(2.95), Inches(0.5), fill=RGBColor(54, 64, 74))
    pill.line.fill.background()
    add_text(slide, Inches(1.13), Inches(5.98), Inches(2.5), Inches(0.16), "Private Core + Public Gateway", 11, WHITE, True)

    add_text(slide, Inches(6.12), Inches(0.95), Inches(5.9), Inches(0.2), "FOR LEADERSHIP & IT", 11, MUTED)
    add_text(
        slide,
        Inches(6.12),
        Inches(1.28),
        Inches(6.1),
        Inches(0.9),
        "核心系统不对公网裸露，\n微信小程序只通过一层受控 HTTPS 网关访问。",
        25,
        NAVY,
        True,
    )
    add_bullets(
        slide,
        Inches(6.16),
        Inches(2.28),
        Inches(5.6),
        Inches(1.3),
        [
            "Web 管理端仅开放公司内网访问。",
            "Backend / MySQL / Redis / MinIO 全部保留在内网。",
            "微信小程序只访问一个公网 HTTPS 域名。",
        ],
        16,
    )

    add_image(slide, "dashboard_real.png", Inches(6.12), Inches(3.78), Inches(3.6), Inches(2.35))
    add_image(slide, "miniapp_home_real.png", Inches(9.98), Inches(3.78), Inches(2.15), Inches(2.35))
    add_text(slide, Inches(10.55), Inches(6.82), Inches(1.55), Inches(0.14), "2026-04", 10.5, MUTED, False, PP_ALIGN.RIGHT)


def add_pain(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "为什么要这样部署", "这页先解决领导和 IT 的同一个疑问：为什么不能简单做成纯内网小程序。")

    card(slide, Inches(0.8), Inches(1.45), Inches(5.7), Inches(5.55), fill=SOFT_ORANGE)
    add_text(slide, Inches(1.04), Inches(1.8), Inches(3.2), Inches(0.22), "当前诉求", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.06),
        Inches(2.18),
        Inches(5.0),
        Inches(4.1),
        [
            "公司希望核心系统和数据尽量留在内网，降低暴露面。",
            "但现场人员和管理层又需要在手机端、外网环境下使用小程序。",
            "同时系统还涉及图片、视频上传与预览，链路不能太脆弱。",
            "最终目标不是演示可访问，而是正式上线后稳定、可审计、可维护。",
        ],
        16,
    )

    card(slide, Inches(6.72), Inches(1.45), Inches(5.8), Inches(5.55), fill=WHITE)
    add_text(slide, Inches(6.98), Inches(1.8), Inches(3.0), Inches(0.22), "关键事实", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(6.99),
        Inches(2.18),
        Inches(5.1),
        Inches(3.7),
        [
            "微信小程序正式版必须配置合法域名，且要求 HTTPS。",
            "手机端必须真实访问到这个域名，纯内网地址无法满足小程序要求。",
            "因此不能让小程序直连纯内网服务，但可以让它访问一个最小化公网网关。",
            "正确思路不是把系统搬到公网，而是只暴露一个受控入口。",
        ],
        16,
    )
    add_text(slide, Inches(7.0), Inches(5.95), Inches(5.0), Inches(0.4), "结论：数据留内网，入口走公网网关，是这套系统最稳妥的落地方式。", 18, RED, True)


def add_architecture(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "推荐架构", "把“谁在用、谁对外、谁留内网”拆清楚，后续 IT 才好实施。")

    boxes = [
        (Inches(0.8), Inches(1.75), Inches(2.45), Inches(1.2), "微信小程序", "现场人员 / 管理层\n手机外网访问", WHITE),
        (Inches(4.1), Inches(1.55), Inches(4.3), Inches(1.55), "公网 HTTPS 网关", "pih-api.company.com\nNginx / WAF / API Gateway", SOFT_TEAL),
        (Inches(9.1), Inches(1.75), Inches(3.45), Inches(1.2), "内网浏览器", "项目经理 / 管理层 / 管理员", WHITE),
    ]
    for left, top, width, height, title, desc, fill in boxes:
        card(slide, left, top, width, height, fill)
        add_text(slide, left + Inches(0.14), top + Inches(0.22), width - Inches(0.28), Inches(0.22), title, 18, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.14), top + Inches(0.58), width - Inches(0.28), Inches(0.5), desc, 12.5, SLATE, False, PP_ALIGN.CENTER)

    for x in [Inches(3.35), Inches(8.55)]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, Inches(2.08), Inches(0.46), Inches(0.42))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()

    card(slide, Inches(1.12), Inches(3.65), Inches(11.0), Inches(2.45), WHITE)
    add_text(slide, Inches(1.35), Inches(3.92), Inches(2.5), Inches(0.22), "内网核心区", 19, NAVY, True)

    core = [
        ("Web 管理端", "Vue3\n仅内网开放", Inches(1.45)),
        ("Backend", "Spring Boot\n统一业务规则", Inches(3.62)),
        ("MySQL", "业务主数据库", Inches(5.83)),
        ("Redis", "缓存 / 会话", Inches(8.0)),
        ("MinIO", "图片 / 视频附件", Inches(10.0)),
    ]
    for title, desc, x in core:
        fill = SOFT_BLUE if title == "Backend" else PAPER
        card(slide, x, Inches(4.35), Inches(1.72), Inches(1.05), fill)
        add_text(slide, x + Inches(0.1), Inches(4.54), Inches(1.5), Inches(0.18), title, 15.5, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.08), Inches(4.84), Inches(1.55), Inches(0.26), desc, 11.2, SLATE, False, PP_ALIGN.CENTER)


def add_boundary(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "什么应该暴露，什么不应该暴露", "这页给 IT 做边界判断，也方便领导理解风险控制点。")

    card(slide, Inches(0.8), Inches(1.45), Inches(5.75), Inches(5.6), SOFT_TEAL)
    card(slide, Inches(6.77), Inches(1.45), Inches(5.75), Inches(5.6), SOFT_RED)

    add_text(slide, Inches(1.05), Inches(1.82), Inches(2.8), Inches(0.22), "允许对公网开放", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(2.18),
        Inches(5.0),
        Inches(4.4),
        [
            "小程序必要接口：登录、项目列表、问题列表、问题详情、评论、状态流转。",
            "附件上传与附件预览接口。",
            "字典只读接口，例如分类、来源、责任部门等。",
            "仅通过一个 HTTPS 域名统一出口，不开放零散端口。",
        ],
        16,
    )

    add_text(slide, Inches(7.02), Inches(1.82), Inches(3.0), Inches(0.22), "不建议对公网开放", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.04),
        Inches(2.18),
        Inches(5.0),
        Inches(4.4),
        [
            "Web 管理后台。",
            "MySQL / Redis / MinIO Console。",
            "管理员接口、用户管理、项目团队、字典管理。",
            "Swagger、运维接口、内部报表接口。",
        ],
        16,
    )


def add_network(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "网络与防火墙建议", "把链路缩到最小，避免“能访问”变成“都能访问”。")

    card(slide, Inches(0.8), Inches(1.45), Inches(12.0), Inches(5.65), WHITE)
    headers = ["来源", "目标", "端口", "建议", "说明"]
    rows = [
        ("公网", "公网网关", "443", "开放", "小程序正式访问入口"),
        ("公网", "内网 Backend", "-", "关闭", "禁止直连业务服务"),
        ("公网", "MySQL / Redis / MinIO Console", "-", "关闭", "禁止直连数据层"),
        ("公网网关", "内网 Backend", "8081", "白名单开放", "仅网关服务器可访问"),
        ("Backend", "MySQL", "3306", "内网开放", "业务数据读写"),
        ("Backend", "Redis", "6379", "内网开放", "缓存 / 会话"),
        ("Backend", "MinIO", "9000", "内网开放", "附件存取"),
    ]
    col_x = [0.98, 2.65, 4.55, 5.75, 7.15]
    col_w = [1.45, 1.6, 0.9, 1.25, 5.1]
    for i, header in enumerate(headers):
        add_text(slide, Inches(col_x[i]), Inches(1.78), Inches(col_w[i]), Inches(0.2), header, 13.5, MUTED, True)

    for idx, row in enumerate(rows):
        y = 2.08 + idx * 0.6
        row_card = card(slide, Inches(0.94), Inches(y), Inches(11.65), Inches(0.48), PAPER if idx % 2 == 0 else WHITE)
        row_card.line.fill.background()
        for i, value in enumerate(row):
            add_text(slide, Inches(col_x[i]), Inches(y + 0.1), Inches(col_w[i]), Inches(0.2), value, 12.5, NAVY if i < 4 else SLATE, i == 3)


def add_security(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "安全与审计控制", "给领导的是放心点，给 IT 的是落地点。")

    items = [
        ("访问控制", "公网只开 443，只代理 /api/，其余路径一律拒绝。", SOFT_TEAL),
        ("身份认证", "小程序走微信绑定 + JWT，Web 走账号密码 + RBAC。", SOFT_BLUE),
        ("附件策略", "图片和视频仍走后端接口中转，不直接暴露 MinIO。", SOFT_ORANGE),
        ("日志审计", "状态流转、分派、优先级、关闭、删除都保留操作日志。", WHITE),
        ("备份策略", "定期备份 MySQL 与 MinIO，升级前先备份再发布。", WHITE),
        ("异常收敛", "可在网关层补限流、IP 白名单、WAF 与审计日志。", WHITE),
    ]
    for idx, (title, desc, fill) in enumerate(items):
        x = Inches(0.8 + (idx % 2) * 6.0)
        y = Inches(1.5 + (idx // 2) * 1.7)
        card(slide, x, y, Inches(5.75), Inches(1.35), fill)
        add_text(slide, x + Inches(0.16), y + Inches(0.18), Inches(1.6), Inches(0.2), title, 16.5, NAVY, True)
        add_text(slide, x + Inches(0.16), y + Inches(0.52), Inches(5.2), Inches(0.45), desc, 13.5, SLATE)


def add_dataflow(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "图片、视频与数据流向", "这一页解决 IT 对附件路径和数据边界的担心。")

    card(slide, Inches(0.8), Inches(1.55), Inches(12.0), Inches(5.4), WHITE)
    steps = [
        ("1. 小程序上传", "现场人员在手机端拍照/录视频后，通过公网 HTTPS 网关上传。", Inches(1.0), Inches(2.2), SOFT_ORANGE),
        ("2. 网关转发", "网关只负责入口，不做业务处理，转发到内网 Backend。", Inches(4.15), Inches(2.2), SOFT_TEAL),
        ("3. 后端入库", "Backend 统一校验权限、大小、类型，然后写入内网 MinIO 与 MySQL。", Inches(7.25), Inches(2.2), SOFT_BLUE),
        ("4. 预览读取", "后续 Web 或小程序预览时，仍通过后端接口读取，不暴露 MinIO 控制台。", Inches(10.2), Inches(2.2), WHITE),
    ]
    for title, desc, x, y, fill in steps:
        card(slide, x, y, Inches(2.55), Inches(1.8), fill)
        add_text(slide, x + Inches(0.12), y + Inches(0.14), Inches(2.1), Inches(0.2), title, 15.5, NAVY, True)
        add_text(slide, x + Inches(0.12), y + Inches(0.46), Inches(2.25), Inches(0.92), desc, 12.5, SLATE)
    for x in [Inches(3.62), Inches(6.73), Inches(9.77)]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, Inches(2.82), Inches(0.35), Inches(0.38))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()
    add_text(slide, Inches(1.02), Inches(5.55), Inches(10.8), Inches(0.3), "建议继续保持“附件统一经后端接口上传和预览”的现有实现，这也是最适合内外网混合部署的方式。", 17, RED, True)


def add_steps(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "实际部署步骤", "这页给 IT 照着执行，也能让领导看到工作量和边界。")

    phases = [
        ("阶段 1：先把内网核心跑稳", ["部署 MySQL / Redis / MinIO / Backend / Web。", "确认内网 Web 正常访问，后台功能和数据链路可用。"]),
        ("阶段 2：部署公网网关", ["准备公网 HTTPS 域名和证书。", "在 DMZ 或边界区部署 Nginx / API Gateway，只代理小程序接口。"]),
        ("阶段 3：配置微信小程序", ["在小程序后台配置 request/upload/download 合法域名。", "填写正式 AppID / Secret，关闭 mock 登录。"]),
        ("阶段 4：真机联调与灰度试用", ["验证登录、录入问题、附件上传、预览、评论、状态流转。", "先选一个真实项目灰度试用，再逐步推广。"]),
    ]
    for idx, (title, bullets) in enumerate(phases):
        x = Inches(0.8 + (idx % 2) * 6.0)
        y = Inches(1.55 + (idx // 2) * 2.3)
        card(slide, x, y, Inches(5.75), Inches(1.95), WHITE)
        add_text(slide, x + Inches(0.16), y + Inches(0.16), Inches(3.6), Inches(0.2), title, 16.5, NAVY, True)
        add_bullets(slide, x + Inches(0.14), y + Inches(0.48), Inches(5.25), Inches(1.15), bullets, 13.5)


def add_release(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "正式上线前必须完成的配置", "这页主要给 IT 和小程序管理员。")

    card(slide, Inches(0.8), Inches(1.5), Inches(5.85), Inches(5.6), SOFT_BLUE)
    add_text(slide, Inches(1.04), Inches(1.82), Inches(2.5), Inches(0.22), "后端与环境变量", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.06),
        Inches(2.18),
        Inches(5.1),
        Inches(4.45),
        [
            "MINIO_PUBLIC_ENDPOINT 指向正式附件访问域名。",
            "REPORT_PUBLIC_BASE_URL 指向正式 Web 域名。",
            "WECHAT_MINIAPP_MOCK_ENABLED 改为 false。",
            "填写正式 AppID / Secret，并按正式域名部署 HTTPS 证书。",
        ],
        16,
    )

    card(slide, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.6), SOFT_ORANGE)
    add_text(slide, Inches(7.14), Inches(1.82), Inches(2.8), Inches(0.22), "微信后台配置", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.16),
        Inches(2.18),
        Inches(4.9),
        Inches(4.45),
        [
            "request 合法域名：公网 API 域名。",
            "uploadFile 合法域名：公网 API 域名。",
            "downloadFile 合法域名：公网 API 域名。",
            "真机联调通过后再提交审核与发布。",
        ],
        16,
    )


def add_value(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "给领导和 IT 的价值分别是什么", "同一套方案，要同时回答业务价值和技术可控性。")

    card(slide, Inches(0.8), Inches(1.55), Inches(5.8), Inches(5.45), WHITE)
    add_text(slide, Inches(1.05), Inches(1.85), Inches(2.2), Inches(0.22), "对领导", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(2.2),
        Inches(5.0),
        Inches(4.2),
        [
            "项目问题不再散落在 Excel 和聊天记录里。",
            "管理层能按项目看总量、超期、重点问题和责任人。",
            "问题从录入到关闭形成闭环，关闭证据可追溯。",
            "推进试点时，不需要把核心系统整体搬到公网。",
        ],
        16,
    )

    card(slide, Inches(6.82), Inches(1.55), Inches(5.68), Inches(5.45), WHITE)
    add_text(slide, Inches(7.05), Inches(1.85), Inches(2.2), Inches(0.22), "对 IT / 运维", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.08),
        Inches(2.2),
        Inches(4.95),
        Inches(4.2),
        [
            "安全边界清楚：公网只留一层入口，其余留在内网。",
            "现有系统代码和部署方式改动最小，维护成本可控。",
            "附件、数据库、后台管理都不需要直接暴露公网。",
            "后续可继续叠加 WAF、限流、审计、监控与备份策略。",
        ],
        16,
    )


def add_realshot(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "当前系统形态", "这页用真实界面说明：系统不是概念方案，已经具备可运行基础。")

    card(slide, Inches(0.75), Inches(1.4), Inches(7.75), Inches(5.65), WHITE)
    card(slide, Inches(8.72), Inches(1.4), Inches(3.85), Inches(5.65), WHITE)
    add_text(slide, Inches(1.0), Inches(1.68), Inches(2.5), Inches(0.22), "Web 管理端", 18, NAVY, True)
    add_image(slide, "dashboard_real.png", Inches(0.96), Inches(2.0), Inches(7.3), Inches(4.6))

    add_text(slide, Inches(8.98), Inches(1.68), Inches(2.0), Inches(0.22), "小程序端", 18, NAVY, True)
    add_image(slide, "miniapp_project_hub_real2.png", Inches(9.05), Inches(2.0), Inches(3.2), Inches(4.6))


def add_summary(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "最终建议", "如果准备进入公司内网正式部署，可以按这个结论执行。")

    card(slide, Inches(0.82), Inches(1.45), Inches(11.9), Inches(1.15), WHITE)
    add_text(slide, Inches(1.08), Inches(1.78), Inches(11.1), Inches(0.45), "Web 管理端留内网，核心服务和数据留内网，小程序只通过一层公网 HTTPS 网关接入。", 22, NAVY, True)

    card(slide, Inches(0.82), Inches(2.95), Inches(5.8), Inches(3.7), SOFT_TEAL)
    add_text(slide, Inches(1.05), Inches(3.25), Inches(2.0), Inches(0.22), "这样做的好处", 19, NAVY, True)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(3.6),
        Inches(5.0),
        Inches(2.6),
        [
            "满足微信小程序正式发布要求。",
            "核心数据不直接暴露公网。",
            "现有系统改造成本低，部署可控。",
            "便于后续推广到更多项目和团队。",
        ],
        15,
    )

    card(slide, Inches(6.82), Inches(2.95), Inches(5.9), Inches(3.7), SOFT_ORANGE)
    add_text(slide, Inches(7.06), Inches(3.25), Inches(2.5), Inches(0.22), "下一步建议", 19, NAVY, True)
    add_bullets(
        slide,
        Inches(7.08),
        Inches(3.6),
        Inches(5.1),
        Inches(2.6),
        [
            "由 IT 确认公网网关部署位、证书和域名方案。",
            "选一个真实项目做灰度试用。",
            "完成真机联调、培训和推广 SOP。",
            "上线后逐步补充监控、备份和运维手册。",
        ],
        15,
    )


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_cover(prs)
    add_pain(1, prs)
    add_architecture(2, prs)
    add_boundary(3, prs)
    add_network(4, prs)
    add_security(5, prs)
    add_dataflow(6, prs)
    add_steps(7, prs)
    add_release(8, prs)
    add_value(9, prs)
    add_realshot(10, prs)
    add_summary(11, prs)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
