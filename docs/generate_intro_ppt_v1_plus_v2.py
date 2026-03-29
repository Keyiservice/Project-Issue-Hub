from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
BASE_PPT = ROOT / "非标设备项目现场问题协同系统-第一版补图增强版.pptx"
OUT_PPT = ROOT / "非标设备项目现场问题协同系统-第一版补图增强版-v2.pptx"
REAL = ROOT / "ppt_real_assets"

INK = RGBColor(24, 34, 43)
MUTED = RGBColor(96, 107, 120)
SUBTLE = RGBColor(139, 150, 162)
PAPER = RGBColor(249, 246, 241)
PANEL = RGBColor(255, 252, 248)
LINE = RGBColor(225, 220, 212)
ORANGE = RGBColor(203, 93, 31)
TEAL = RGBColor(22, 123, 120)


def add_bg(slide, width, height):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.7), Inches(5.15), Inches(3.2), Inches(3.2))
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = ORANGE
    glow_left.fill.transparency = 0.93
    glow_left.line.fill.background()

    glow_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.15), Inches(-0.65), Inches(4.1), Inches(4.1))
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = TEAL
    glow_right.fill.transparency = 0.91
    glow_right.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        text_run = p.add_run()
        text_run.text = item
        text_run.font.name = "Microsoft YaHei"
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = INK


def add_header(slide, page_num, title, subtitle):
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(0.56), Inches(0.5), Inches(0.56), Inches(0.56))
    arc.fill.background()
    arc.line.color.rgb = TEAL
    arc.line.width = Pt(3)

    add_text(slide, Inches(0.95), Inches(0.5), Inches(9.2), Inches(0.42), title, 26, INK, True)
    add_text(slide, Inches(0.98), Inches(0.94), Inches(9.4), Inches(0.35), subtitle, 12, MUTED)
    add_text(slide, Inches(10.75), Inches(0.32), Inches(1.75), Inches(0.22), "REAL MINIAPP SHOTS", 10.5, SUBTLE, False, PP_ALIGN.RIGHT)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(7.07), Inches(12.0), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    add_text(slide, Inches(0.65), Inches(7.08), Inches(5.2), Inches(0.22), "First Edition Enhanced v2", 9.5, SUBTLE)
    add_text(slide, Inches(12.0), Inches(7.02), Inches(0.4), Inches(0.22), str(page_num), 9.5, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    panel.line.width = Pt(1)
    return panel


def add_image_panel(slide, img_name, left, top, width, height):
    add_panel(slide, left, top, width, height)
    path = REAL / img_name
    if path.exists():
        slide.shapes.add_picture(str(path), left + Inches(0.08), top + Inches(0.08), width=width - Inches(0.16), height=height - Inches(0.16))


def add_single_image_slide(prs, page_num, title, subtitle, img_name, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle)
    add_image_panel(slide, img_name, Inches(0.72), Inches(1.55), Inches(7.35), Inches(5.2))
    add_panel(slide, Inches(8.3), Inches(1.55), Inches(4.02), Inches(5.2))
    add_bullets(slide, Inches(8.52), Inches(1.85), Inches(3.6), Inches(4.65), bullets, 14.6)


def add_dual_image_slide(prs, page_num, title, subtitle, left_img, right_img, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle)
    add_image_panel(slide, left_img, Inches(0.72), Inches(1.55), Inches(4.1), Inches(5.05))
    add_image_panel(slide, right_img, Inches(4.98), Inches(1.55), Inches(3.18), Inches(5.05))
    add_panel(slide, Inches(8.35), Inches(1.55), Inches(3.97), Inches(5.05))
    add_bullets(slide, Inches(8.56), Inches(1.82), Inches(3.55), Inches(4.45), bullets, 14.2)


def main():
    prs = Presentation(str(BASE_PPT))
    start_page = len(prs.slides) + 1

    add_single_image_slide(
        prs,
        start_page,
        "小程序原图补充：首页总览",
        "保留真实界面，不改颜色，只补充原图截图。",
        "miniapp_home_real.png",
        [
            "首页已经把高频入口收成项目工作台、OPL 录入、项目问题库、项目统计四个主入口。",
            "页面上方直接显示当前人和当前项目，进入项目后就能进入项目级问题库，不会把不同项目混在一起。",
            "这张图来自微信开发者工具里的真实运行页面，不是重新绘制的示意图。"
        ]
    )

    add_dual_image_slide(
        prs,
        start_page + 1,
        "小程序原图补充：开发现场入口感",
        "用真实运行截图补足第一页里缺少的小程序画面。",
        "miniapp_after_seed_reload.png",
        "miniapp_login_crop.png",
        [
            "左侧保留了开发者工具里的真实运行场景，能看清小程序页面处于真实调试状态。",
            "右侧是已经补进系统介绍里的登录原图，和首页一起构成完整的小程序入口说明。",
            "后续如果你要继续补真机首页、项目工作台、问题详情，我会直接往这版 v2 上继续叠。"
        ]
    )

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
