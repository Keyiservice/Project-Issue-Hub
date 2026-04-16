from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "ppt_real_assets"
OUT = ROOT / "Project-Issue-Hub-Intranet-Only-Overview-EN.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(24, 33, 43)
SLATE = RGBColor(79, 90, 103)
MUTED = RGBColor(123, 133, 144)
WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(247, 243, 236)
LINE = RGBColor(229, 223, 214)
TEAL = RGBColor(51, 136, 127)
ORANGE = RGBColor(211, 107, 41)
SOFT_TEAL = RGBColor(227, 242, 238)
SOFT_ORANGE = RGBColor(248, 233, 222)
SOFT_BLUE = RGBColor(231, 240, 250)
SOFT_RED = RGBColor(251, 235, 230)
RED = RGBColor(184, 74, 58)


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.3), Inches(-0.7), Inches(4.0), Inches(4.0))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = TEAL
    glow1.fill.transparency = 0.93
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.6), Inches(5.2), Inches(3.0), Inches(3.0))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ORANGE
    glow2.fill.transparency = 0.94
    glow2.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def card(slide, left, top, width, height, fill=WHITE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def title_block(slide, page_no, title, subtitle):
    add_text(slide, Inches(0.72), Inches(0.28), Inches(1.0), Inches(0.2), f"{page_no:02d}", 11, MUTED, True)
    add_text(slide, Inches(1.12), Inches(0.25), Inches(9.0), Inches(0.4), title, 26, NAVY, True)
    add_text(slide, Inches(1.15), Inches(0.72), Inches(10.5), Inches(0.24), subtitle, 12, SLATE)


def add_image(slide, filename, left, top, width, height):
    path = ASSETS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    left = card(slide, Inches(0.58), Inches(0.55), Inches(5.25), Inches(6.2), fill=NAVY)
    left.line.fill.background()
    add_text(slide, Inches(0.92), Inches(0.9), Inches(3.6), Inches(0.2), "PROJECT ISSUE HUB", 11, RGBColor(205, 213, 220))
    add_text(slide, Inches(0.92), Inches(1.28), Inches(4.3), Inches(1.25), "Intranet-Only\nDeployment Overview", 28, WHITE, True)
    add_text(
        slide,
        Inches(0.92),
        Inches(3.1),
        Inches(4.2),
        Inches(1.15),
        "For leadership & IT\nAll services stay in the intranet.\nAccess via Web only.",
        16,
        RGBColor(224, 231, 236),
    )

    pill = card(slide, Inches(0.92), Inches(5.88), Inches(2.95), Inches(0.5), fill=RGBColor(54, 64, 74))
    pill.line.fill.background()
    add_text(slide, Inches(1.13), Inches(5.98), Inches(2.5), Inches(0.16), "Private Network Only", 11, WHITE, True)

    add_text(slide, Inches(6.12), Inches(0.95), Inches(5.9), Inches(0.2), "INTRANET DEPLOYMENT", 11, MUTED)
    add_text(
        slide,
        Inches(6.12),
        Inches(1.28),
        Inches(6.1),
        Inches(0.9),
        "All core services remain inside the corporate network.\nUsers access via Web on intranet.",
        25,
        NAVY,
        True,
    )
    add_bullets(
        slide,
        Inches(6.16),
        Inches(2.28),
        Inches(5.6),
        Inches(1.3),
        [
            "No public exposure of backend, databases, or storage.",
            "Single intranet domain for Web access.",
            "Operations stay within internal security policies.",
        ],
        16,
    )

    add_image(slide, "dashboard_real.png", Inches(6.12), Inches(3.78), Inches(3.6), Inches(2.35))
    add_image(slide, "issue_detail_real.png", Inches(9.98), Inches(3.78), Inches(2.15), Inches(2.35))
    add_text(slide, Inches(10.55), Inches(6.82), Inches(1.55), Inches(0.14), "2026-04", 10.5, MUTED, False, PP_ALIGN.RIGHT)


def add_pain(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Why This Matters", "Pain points from Excel-based OPL and the targeted improvements.")

    card(slide, Inches(0.8), Inches(1.45), Inches(5.7), Inches(5.55), fill=SOFT_ORANGE)
    add_text(slide, Inches(1.04), Inches(1.8), Inches(3.2), Inches(0.22), "Current Pain Points", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.06),
        Inches(2.18),
        Inches(5.0),
        Inches(4.1),
        [
            "Issues recorded as plain text, without full现场 context.",
            "Repeated questions about impact, owner, and next step.",
            "Status flow and ownership are not transparent.",
            "Cross-project issues are混 in one list.",
            "Statistics are manual, slow, and inconsistent.",
        ],
        16,
    )

    card(slide, Inches(6.72), Inches(1.45), Inches(5.8), Inches(5.55), fill=WHITE)
    add_text(slide, Inches(6.98), Inches(1.8), Inches(3.0), Inches(0.22), "What We Change", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(6.99),
        Inches(2.18),
        Inches(5.1),
        Inches(3.7),
        [
            "Single intranet web system for all projects.",
            "Clear owner, status, next action, and closure evidence.",
            "Project-level dashboards and reporting.",
            "Unified issue library with standard lifecycle.",
        ],
        16,
    )


def add_solution(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Approach", "Web-first intranet deployment for project issue collaboration.")

    steps = [
        ("1. Capture", "Field / PM enters issue with photos, video, and details.", SOFT_ORANGE),
        ("2. Assign", "Owner and department are明确, with due dates.", SOFT_TEAL),
        ("3. Track", "Status flow with evidence and timeline.", SOFT_BLUE),
        ("4. Close", "Verification and closure evidence retained.", WHITE),
    ]
    for idx, (title, desc, fill) in enumerate(steps):
        x = Inches(0.8 + idx * 3.1)
        card(slide, x, Inches(2.2), Inches(2.85), Inches(2.3), fill)
        add_text(slide, x + Inches(0.14), Inches(2.38), Inches(2.4), Inches(0.2), title, 16.5, NAVY, True)
        add_text(slide, x + Inches(0.14), Inches(2.75), Inches(2.5), Inches(0.6), desc, 12.5, SLATE)

    add_text(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4), "Key idea: One intranet system, unified rules, clear ownership, and full traceability.", 18, RED, True)


def add_architecture(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Intranet Architecture", "All services remain inside the corporate network.")

    card(slide, Inches(0.9), Inches(1.6), Inches(11.8), Inches(4.8), WHITE)
    add_text(slide, Inches(1.12), Inches(1.88), Inches(3.2), Inches(0.22), "Intranet Only", 19, NAVY, True)

    core = [
        ("Web Admin", "Vue3\nIntranet domain", Inches(1.2)),
        ("Backend", "Spring Boot\nUnified rules", Inches(3.8)),
        ("MySQL", "Business DB", Inches(6.25)),
        ("Redis", "Cache / Session", Inches(8.4)),
        ("MinIO", "Attachments", Inches(10.1)),
    ]
    for title, desc, x in core:
        fill = SOFT_BLUE if title == "Backend" else PAPER
        card(slide, x, Inches(2.8), Inches(1.8), Inches(1.2), fill)
        add_text(slide, x + Inches(0.1), Inches(3.05), Inches(1.6), Inches(0.18), title, 15, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.08), Inches(3.35), Inches(1.6), Inches(0.3), desc, 11.5, SLATE, False, PP_ALIGN.CENTER)

    add_text(slide, Inches(1.12), Inches(4.6), Inches(10.8), Inches(0.3), "Access via intranet domain only. No public endpoints required in this version.", 14.5, SLATE)


def add_security(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Security & Operations", "Simple controls aligned with internal security policy.")

    items = [
        ("Access Control", "Only intranet DNS resolution and firewall ACLs allow access.", SOFT_TEAL),
        ("RBAC", "Role-based permission model for pages and actions.", SOFT_BLUE),
        ("Audit Trail", "Status flow, assignment, priority changes, closures all logged.", WHITE),
        ("Backups", "MySQL + MinIO backups on schedule.", WHITE),
    ]
    for idx, (title, desc, fill) in enumerate(items):
        x = Inches(0.8 + (idx % 2) * 6.0)
        y = Inches(1.7 + (idx // 2) * 2.0)
        card(slide, x, y, Inches(5.75), Inches(1.55), fill)
        add_text(slide, x + Inches(0.16), y + Inches(0.18), Inches(1.8), Inches(0.2), title, 16, NAVY, True)
        add_text(slide, x + Inches(0.16), y + Inches(0.55), Inches(5.2), Inches(0.45), desc, 13.5, SLATE)


def add_deployment(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Deployment Steps", "Minimal steps for IT to deliver an intranet-only system.")

    phases = [
        ("Step 1: Deploy Core", ["MySQL / Redis / MinIO / Backend / Web in intranet.", "Verify Web login and key flows."]),
        ("Step 2: Set Intranet DNS", ["Bind `pih-web.intra.company.com` to Web server.", "Restrict access to intranet only."]),
        ("Step 3: Data & Users", ["Import projects and OPL data.", "Create users, roles, and project teams."]),
        ("Step 4: Go Live", ["Run a pilot project first, then expand to more teams."]),
    ]
    for idx, (title, bullets) in enumerate(phases):
        x = Inches(0.8 + (idx % 2) * 6.0)
        y = Inches(1.6 + (idx // 2) * 2.3)
        card(slide, x, y, Inches(5.75), Inches(1.95), WHITE)
        add_text(slide, x + Inches(0.16), y + Inches(0.16), Inches(3.6), Inches(0.2), title, 16.5, NAVY, True)
        add_bullets(slide, x + Inches(0.14), y + Inches(0.48), Inches(5.25), Inches(1.15), bullets, 13.5)


def add_realshot(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Current System", "Real screenshots to show the system is operational.")

    card(slide, Inches(0.75), Inches(1.4), Inches(7.75), Inches(5.65), WHITE)
    card(slide, Inches(8.72), Inches(1.4), Inches(3.85), Inches(5.65), WHITE)
    add_text(slide, Inches(1.0), Inches(1.68), Inches(2.5), Inches(0.22), "Web Dashboard", 18, NAVY, True)
    add_image(slide, "dashboard_real.png", Inches(0.96), Inches(2.0), Inches(7.3), Inches(4.6))

    add_text(slide, Inches(8.98), Inches(1.68), Inches(2.0), Inches(0.22), "Issue Detail", 18, NAVY, True)
    add_image(slide, "issue_detail_real.png", Inches(9.05), Inches(2.0), Inches(3.2), Inches(4.6))


def add_summary(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_block(slide, slide_no, "Recommendation", "Best fit for internal deployment and adoption.")

    card(slide, Inches(0.82), Inches(1.45), Inches(11.9), Inches(1.15), WHITE)
    add_text(slide, Inches(1.08), Inches(1.78), Inches(11.1), Inches(0.45), "Keep all services inside the intranet. Provide Web-only access for internal teams.", 22, NAVY, True)

    card(slide, Inches(0.82), Inches(2.95), Inches(5.8), Inches(3.7), SOFT_TEAL)
    add_text(slide, Inches(1.05), Inches(3.25), Inches(2.0), Inches(0.22), "Benefits", 19, NAVY, True)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(3.6),
        Inches(5.0),
        Inches(2.6),
        [
            "No public exposure of core services.",
            "Faster internal approvals.",
            "Lower operational risk.",
            "Fits internal security policies.",
        ],
        15,
    )

    card(slide, Inches(6.82), Inches(2.95), Inches(5.9), Inches(3.7), SOFT_ORANGE)
    add_text(slide, Inches(7.06), Inches(3.25), Inches(2.5), Inches(0.22), "Next Steps", 19, NAVY, True)
    add_bullets(
        slide,
        Inches(7.08),
        Inches(3.6),
        Inches(5.1),
        Inches(2.6),
        [
            "Confirm intranet servers and DNS plan.",
            "Pilot with one live project.",
            "Finalize training and SOP.",
            "Review backup and monitoring cadence.",
        ],
        15,
    )


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_cover(prs)
    add_pain(1, prs)
    add_solution(2, prs)
    add_architecture(3, prs)
    add_security(4, prs)
    add_deployment(5, prs)
    add_realshot(6, prs)
    add_summary(7, prs)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
