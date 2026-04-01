from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REAL = ROOT / "ppt_real_assets"
OUT = ROOT / "Project-Issue-Hub-部署与发布指南.pptx"

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(20, 28, 38)
SLATE = RGBColor(82, 92, 104)
MUTED = RGBColor(127, 138, 149)
PAPER = RGBColor(247, 243, 236)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(214, 110, 38)
TEAL = RGBColor(56, 142, 132)
SOFT_ORANGE = RGBColor(248, 233, 221)
SOFT_TEAL = RGBColor(227, 242, 238)
LINE = RGBColor(228, 222, 214)
RED = RGBColor(185, 76, 58)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.2), Inches(-1.0), Inches(4.6), Inches(4.6))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = TEAL
    glow1.fill.transparency = 0.92
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.9), Inches(5.2), Inches(3.4), Inches(3.4))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ORANGE
    glow2.fill.transparency = 0.94
    glow2.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def panel(slide, left, top, width, height, fill=WHITE, rounded=True):
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(st, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def title_block(slide, page_no, title, subtitle):
    add_text(slide, Inches(0.74), Inches(0.34), Inches(0.55), Inches(0.25), f"{page_no:02d}", 11, MUTED, True)
    add_text(slide, Inches(1.16), Inches(0.3), Inches(8.0), Inches(0.45), title, 26, NAVY, True)
    add_text(slide, Inches(1.2), Inches(0.78), Inches(9.8), Inches(0.3), subtitle, 12.5, SLATE)


def add_image(slide, filename, left, top, width, height):
    path = REAL / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def stat_card(slide, left, top, title, value, note, fill=WHITE):
    panel(slide, left, top, Inches(2.7), Inches(1.18), fill=fill)
    add_text(slide, left + Inches(0.12), top + Inches(0.12), Inches(2.2), Inches(0.2), title, 12, MUTED, True)
    add_text(slide, left + Inches(0.12), top + Inches(0.4), Inches(1.0), Inches(0.35), value, 24, NAVY, True)
    add_text(slide, left + Inches(0.12), top + Inches(0.8), Inches(2.3), Inches(0.2), note, 10.5, SLATE)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    left = panel(slide, Inches(0.62), Inches(0.58), Inches(5.35), Inches(6.2), fill=NAVY)
    left.line.fill.background()
    add_text(slide, Inches(0.95), Inches(0.92), Inches(3.6), Inches(0.3), "PROJECT ISSUE HUB", 12, RGBColor(203, 212, 220), False)
    add_text(slide, Inches(0.95), Inches(1.3), Inches(4.7), Inches(1.25),
             "系统部署与\n小程序发布指南", 28, WHITE, True)
    add_text(slide, Inches(0.95), Inches(2.55), Inches(4.2), Inches(0.9),
             "覆盖环境要求、Docker 部署、生产上线、\n数据备份、微信小程序发布与运维检查。", 17, RGBColor(228, 234, 240))

    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.86), Inches(2.9), Inches(0.52))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(54, 64, 75)
    pill.line.fill.background()
    add_text(slide, Inches(1.12), Inches(5.96), Inches(2.5), Inches(0.18), "Docker + Web + Mini Program", 11, WHITE, True)

    add_text(slide, Inches(6.28), Inches(0.92), Inches(5.9), Inches(0.3), "DEPLOYMENT & RELEASE PLAYBOOK", 11, MUTED)
    add_text(slide, Inches(6.28), Inches(1.28), Inches(6.0), Inches(0.9),
             "从本地联调到生产上线，\n把部署动作标准化。", 26, NAVY, True)
    add_bullets(slide, Inches(6.34), Inches(2.28), Inches(5.7), Inches(1.5), [
        "明确服务器、端口、域名、账号和配置要求",
        "统一 Linux / Docker Compose 的部署与升级流程",
        "把微信小程序正式发布前提和步骤一次讲清楚",
    ], 16)

    add_image(slide, "dashboard_real.png", Inches(6.3), Inches(4.05), Inches(6.2), Inches(2.55))
    add_text(slide, Inches(10.55), Inches(6.78), Inches(1.65), Inches(0.15), "2026-04-01", 10.5, MUTED, False, PP_ALIGN.RIGHT)


def add_scope(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "本指南覆盖什么", "不是只告诉你一条命令，而是把部署、验收、发布和运维都说完整。")

    topics = [
        ("部署对象", "Web、Backend、MySQL、Redis、MinIO、小程序"),
        ("环境要求", "本地开发、测试机、生产服务器的配置建议"),
        ("配置项", "端口、域名、AppID、Secret、MinIO 对外地址"),
        ("上线流程", "Docker 启动、验证、备份、升级、回滚"),
        ("小程序发布", "开发版、上传、审核、发布、合法域名"),
        ("运维检查", "日志、容量、备份、上传、登录、报表验证"),
    ]
    for idx, (title, desc) in enumerate(topics):
        x = Inches(0.82 + (idx % 3) * 4.08)
        y = Inches(1.55 + (idx // 3) * 2.08)
        panel(slide, x, y, Inches(3.65), Inches(1.58), fill=WHITE)
        add_text(slide, x + Inches(0.16), y + Inches(0.18), Inches(2.8), Inches(0.24), title, 17, NAVY, True)
        add_text(slide, x + Inches(0.16), y + Inches(0.56), Inches(3.15), Inches(0.62), desc, 13.5, SLATE)

    add_text(slide, Inches(0.94), Inches(5.95), Inches(10.9), Inches(0.3),
             "交付文件会同时提供：部署指南 Markdown + 部署与发布 PPT，后续参数调整时可以持续维护。", 14, SLATE)


def add_arch(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "系统部署拓扑", "一套后端规则，同时服务小程序和 Web。")

    cols = [
        ("客户端", "微信小程序\nWeb 管理后台"),
        ("业务服务", "Spring Boot Backend\nJWT / RBAC / 状态机"),
        ("基础设施", "MySQL\nRedis\nMinIO"),
    ]
    xs = [Inches(0.95), Inches(4.52), Inches(8.1)]
    for idx, (title, desc) in enumerate(cols):
        panel(slide, xs[idx], Inches(2.0), Inches(3.1), Inches(2.2))
        add_text(slide, xs[idx] + Inches(0.16), Inches(2.25), Inches(2.0), Inches(0.24), title, 18, NAVY, True)
        add_text(slide, xs[idx] + Inches(0.16), Inches(2.72), Inches(2.6), Inches(1.0), desc, 16, SLATE)
        if idx < 2:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, xs[idx] + Inches(3.15), Inches(2.88), Inches(0.35), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()

    band = panel(slide, Inches(0.95), Inches(5.15), Inches(11.35), Inches(1.0))
    add_bullets(slide, Inches(1.12), Inches(5.38), Inches(10.9), Inches(0.55), [
        "所有业务规则收口在后端：状态流转、权限、编号、超期、导出、附件访问。",
        "小程序负责轻操作，Web 负责重管理，两端共用同一套数据库和附件存储。",
    ], 14)


def add_requirements(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "环境要求", "开发机和正式服务器都要先满足基础条件。")

    panel(slide, Inches(0.78), Inches(1.42), Inches(6.0), Inches(4.95))
    panel(slide, Inches(6.92), Inches(1.42), Inches(5.63), Inches(4.95))

    add_text(slide, Inches(1.02), Inches(1.68), Inches(2.6), Inches(0.25), "本地开发 / 联调", 19, NAVY, True)
    add_bullets(slide, Inches(1.05), Inches(2.06), Inches(5.45), Inches(3.8), [
        "Windows / macOS / Linux 任一桌面环境",
        "Docker Desktop 或 Docker Engine",
        "微信开发者工具",
        "建议 4 核 CPU、16 GB 内存、30 GB 可用空间",
        "推荐 Node.js 20+、Java 17、Git 2.4+",
    ], 15.5)

    add_text(slide, Inches(7.18), Inches(1.68), Inches(2.6), Inches(0.25), "正式服务器", 19, NAVY, True)
    add_bullets(slide, Inches(7.22), Inches(2.06), Inches(5.0), Inches(3.8), [
        "Ubuntu 22.04 / Rocky 9 / CentOS Stream 9",
        "最低 4 核 8 GB，推荐 8 核 16 GB",
        "100 GB 以上 SSD 磁盘",
        "需要可访问镜像源或企业私有仓库",
        "必须具备备份与 HTTPS 域名能力",
    ], 15.5)


def add_ports(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "默认端口与访问入口", "端口来自当前 docker-compose 配置，正式环境可通过 .env 调整。")

    rows = [
        ("Web", "80", "http://localhost", "管理后台"),
        ("Backend", "8081", "http://localhost:8081", "REST API / Swagger"),
        ("MySQL", "3307", "127.0.0.1:3307", "业务数据库"),
        ("Redis", "6380", "127.0.0.1:6380", "缓存"),
        ("MinIO API", "9002", "http://localhost:9002", "附件对象接口"),
        ("MinIO Console", "9003", "http://localhost:9003", "文件管理控制台"),
    ]
    for idx, row in enumerate(rows):
        y = Inches(1.45 + idx * 0.78)
        fill = SOFT_ORANGE if idx % 2 == 0 else WHITE
        panel(slide, Inches(0.86), y, Inches(11.6), Inches(0.58), fill=fill)
        add_text(slide, Inches(1.05), y + Inches(0.12), Inches(1.4), Inches(0.16), row[0], 14, NAVY, True)
        add_text(slide, Inches(2.55), y + Inches(0.12), Inches(1.0), Inches(0.16), row[1], 13, NAVY)
        add_text(slide, Inches(3.62), y + Inches(0.12), Inches(3.35), Inches(0.16), row[2], 13, SLATE)
        add_text(slide, Inches(7.35), y + Inches(0.12), Inches(3.6), Inches(0.16), row[3], 13, SLATE)

    panel(slide, Inches(0.88), Inches(6.35), Inches(5.75), Inches(0.62), fill=WHITE)
    add_text(slide, Inches(1.05), Inches(6.52), Inches(5.3), Inches(0.2),
             "默认管理员账号：admin / 123456", 14.5, NAVY, True)

    panel(slide, Inches(6.74), Inches(6.35), Inches(5.7), Inches(0.62), fill=WHITE)
    add_text(slide, Inches(6.95), Inches(6.52), Inches(5.2), Inches(0.2),
             "报表与附件对外访问依赖 MINIO_PUBLIC_ENDPOINT / REPORT_PUBLIC_BASE_URL", 12.5, SLATE)


def add_local_deploy(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "本地一键启动", "对开发、联调和演示环境，当前项目已经支持 Docker Compose 一键拉起。")

    panel(slide, Inches(0.8), Inches(1.45), Inches(5.0), Inches(4.75))
    add_text(slide, Inches(1.02), Inches(1.72), Inches(1.8), Inches(0.25), "命令与验证", 19, NAVY, True)
    add_text(slide, Inches(1.02), Inches(2.15), Inches(4.35), Inches(1.5),
             "1. docker compose up -d --build\n\n"
             "2. docker compose ps\n\n"
             "3. 打开 http://localhost\n"
             "4. 登录 admin / 123456\n"
             "5. 打开 http://localhost:8081/swagger-ui.html", 16, SLATE)
    add_text(slide, Inches(1.02), Inches(4.6), Inches(4.3), Inches(0.9),
             "说明：首次启动会自动初始化数据库、默认账号、示例项目与示例问题数据。", 14, SLATE)

    add_image(slide, "login_real.png", Inches(6.0), Inches(1.58), Inches(3.0), Inches(2.15))
    add_image(slide, "dashboard_real.png", Inches(8.9), Inches(1.52), Inches(3.48), Inches(2.78))
    add_image(slide, "issue_projects_real.png", Inches(6.0), Inches(4.05), Inches(6.38), Inches(2.2))


def add_prod_env(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "正式环境需要修改哪些配置", "上线不是直接复用开发环境，需要切到正式域名和真实小程序参数。")

    panel(slide, Inches(0.82), Inches(1.48), Inches(5.9), Inches(4.98))
    panel(slide, Inches(6.83), Inches(1.48), Inches(5.7), Inches(4.98))

    add_text(slide, Inches(1.02), Inches(1.74), Inches(2.5), Inches(0.25), ".env 关键项", 19, NAVY, True)
    add_bullets(slide, Inches(1.05), Inches(2.15), Inches(5.3), Inches(3.85), [
        "MINIO_PUBLIC_ENDPOINT：改成正式附件域名",
        "REPORT_PUBLIC_BASE_URL：改成正式 Web 域名",
        "WECHAT_MINIAPP_MOCK_ENABLED=false",
        "WECHAT_MINIAPP_APP_ID=正式 AppID",
        "WECHAT_MINIAPP_SECRET=正式 Secret",
        "按服务器实际情况调整各服务宿主机端口",
    ], 15)

    add_text(slide, Inches(7.05), Inches(1.74), Inches(2.6), Inches(0.25), "上线前必须具备", 19, NAVY, True)
    add_bullets(slide, Inches(7.08), Inches(2.15), Inches(5.05), Inches(3.85), [
        "公网可访问 HTTPS 域名",
        "微信公众平台已配置 request/upload/download 合法域名",
        "MinIO 附件访问地址已可被浏览器和小程序访问",
        "Web、API、文件访问路径已验证",
        "数据库与附件备份策略已准备",
    ], 15)


def add_linux_steps(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "Linux 服务器部署步骤", "推荐用一台 Linux 服务器先完成完整试运行，再视规模拆分。")

    steps = [
        "1. 安装 Docker 与 docker compose",
        "2. git clone 项目源码，复制 .env.example 为 .env",
        "3. 填写正式域名、AppID、Secret、端口",
        "4. 执行 docker compose up -d --build",
        "5. 用 docker compose ps / logs 验证容器状态",
        "6. 验证 Web、API、MinIO、报表导出、小程序登录",
    ]
    add_bullets(slide, Inches(0.95), Inches(1.7), Inches(5.5), Inches(4.25), steps, 15.5)

    panel(slide, Inches(6.2), Inches(1.62), Inches(6.0), Inches(4.85))
    add_text(slide, Inches(6.42), Inches(1.9), Inches(2.4), Inches(0.25), "推荐命令", 19, NAVY, True)
    add_text(slide, Inches(6.42), Inches(2.35), Inches(5.3), Inches(3.55),
             "mkdir -p /opt/project-issue-hub\n"
             "cd /opt/project-issue-hub\n"
             "git clone <repo> .\n"
             "cp .env.example .env\n"
             "docker compose up -d --build\n"
             "docker compose ps\n"
             "docker compose logs --tail=200 backend", 15.5, SLATE, False, PP_ALIGN.LEFT, "Consolas")


def add_ops(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "备份、升级与回滚", "部署完成后，真正决定系统可持续性的，是备份和标准升级动作。")

    panel(slide, Inches(0.82), Inches(1.55), Inches(3.9), Inches(4.72))
    panel(slide, Inches(4.86), Inches(1.55), Inches(3.35), Inches(4.72))
    panel(slide, Inches(8.35), Inches(1.55), Inches(4.15), Inches(4.72))

    add_text(slide, Inches(1.02), Inches(1.82), Inches(1.8), Inches(0.22), "备份", 18, NAVY, True)
    add_bullets(slide, Inches(1.04), Inches(2.18), Inches(3.35), Inches(3.55), [
        "MySQL：每日 mysqldump",
        "MinIO：每日增量 / 每周全量",
        "保留 .env 与 docker-compose 存档",
        "正式环境不要执行 docker compose down -v",
    ], 14.5)

    add_text(slide, Inches(5.05), Inches(1.82), Inches(1.8), Inches(0.22), "升级", 18, NAVY, True)
    add_bullets(slide, Inches(5.08), Inches(2.18), Inches(2.75), Inches(3.55), [
        "先备份数据库",
        "git pull / 切 tag",
        "docker compose up -d --build",
        "验证登录、上传、导出",
    ], 14.5)

    add_text(slide, Inches(8.58), Inches(1.82), Inches(1.8), Inches(0.22), "回滚", 18, NAVY, True)
    add_bullets(slide, Inches(8.62), Inches(2.18), Inches(3.55), Inches(3.55), [
        "回退到上一个 Git tag",
        "重新构建容器",
        "必要时恢复 SQL 备份",
        "上线前建议保留可回滚镜像",
    ], 14.5)


def add_miniapp_pre(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "小程序发布前提", "开发者工具里能跑，不等于真机能跑；正式发布必须满足域名和账号条件。")

    add_image(slide, "miniapp_home_real.png", Inches(0.85), Inches(1.58), Inches(2.8), Inches(4.95))
    panel(slide, Inches(4.0), Inches(1.56), Inches(8.48), Inches(4.95))
    add_bullets(slide, Inches(4.18), Inches(1.9), Inches(7.95), Inches(3.9), [
        "后端必须关闭 mock 登录：WECHAT_MINIAPP_MOCK_ENABLED=false",
        "后端必须配置真实 AppID 和 Secret",
        "接口地址不能再用 127.0.0.1，必须换成 HTTPS 正式域名",
        "微信公众平台必须配置 request / uploadFile / downloadFile 合法域名",
        "小程序正式发布前，要把登录、项目选择、创建问题、上传附件、关闭问题完整走通",
    ], 16)
    add_text(slide, Inches(4.2), Inches(5.95), Inches(7.8), Inches(0.25),
             "开发阶段关闭合法域名校验仅适用于开发者工具，真机与线上版都不适用。", 13.5, RED, True)


def add_miniapp_release(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "小程序发布流程", "按微信标准流程：上传代码 -> 提交审核 -> 发布上线。")

    steps = [
        ("1", "准备正式环境", "后端、域名、合法域名配置全部就绪"),
        ("2", "开发者工具上传", "填写版本号与版本说明，生成开发版"),
        ("3", "后台提交审核", "在微信公众平台版本管理中提交审核"),
        ("4", "审核通过后发布", "发布为线上版本，通知试点团队使用"),
    ]
    for idx, (n, title, desc) in enumerate(steps):
        x = Inches(0.9 + idx * 3.05)
        panel(slide, x, Inches(2.15), Inches(2.7), Inches(2.05))
        badge = panel(slide, x + Inches(0.18), Inches(2.34), Inches(0.42), Inches(0.42), SOFT_ORANGE)
        badge.line.fill.background()
        add_text(slide, x + Inches(0.29), Inches(2.42), Inches(0.16), Inches(0.12), n, 12, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.72), Inches(2.32), Inches(1.6), Inches(0.2), title, 16, NAVY, True)
        add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2.2), Inches(0.6), desc, 12.5, SLATE)
        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(2.78), Inches(2.9), Inches(0.22), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()

    add_image(slide, "miniapp_login_crop.png", Inches(0.95), Inches(4.65), Inches(2.6), Inches(2.0))
    add_image(slide, "miniapp_issue_list_real.png", Inches(3.8), Inches(4.65), Inches(2.6), Inches(2.0))
    add_image(slide, "miniapp_project_hub_real.png", Inches(6.67), Inches(4.65), Inches(2.6), Inches(2.0))
    add_image(slide, "miniapp_home_workspace_real.png", Inches(9.55), Inches(4.65), Inches(2.6), Inches(2.0))


def add_acceptance(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "上线验收清单", "不要只看页面能打开，要按核心链路验收。")

    checks = [
        "Web 可登录，项目问题库正常加载",
        "小程序可登录、可选项目、可创建问题",
        "图片 / 视频附件可上传、可预览",
        "问题可分派、跟进、待验证、关闭",
        "关闭证据可留存",
        "统计页面按项目切换正常",
        "Excel Report 可导出且附件链接可用",
        "项目团队与责任人筛选正常",
    ]
    panel(slide, Inches(0.82), Inches(1.56), Inches(5.55), Inches(4.9))
    add_bullets(slide, Inches(1.02), Inches(1.9), Inches(4.95), Inches(4.1), checks, 15)

    add_image(slide, "issues_real.png", Inches(6.6), Inches(1.62), Inches(2.8), Inches(2.08))
    add_image(slide, "issue_detail_real.png", Inches(9.48), Inches(1.62), Inches(2.7), Inches(2.08))
    add_image(slide, "stats_real.png", Inches(6.6), Inches(3.95), Inches(2.8), Inches(2.08))
    add_image(slide, "users_real.png", Inches(9.48), Inches(3.95), Inches(2.7), Inches(2.08))


def add_rollout(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "推荐推广顺序", "先跑通，再试点，再推广，避免一上来全员切换。")

    phases = [
        ("阶段 1", "内网试运行", "先部署 Web + Backend + 数据库，让项目经理和管理层先用 Web 端管理问题库。"),
        ("阶段 2", "小程序试点", "选择一个项目组先试点录入、附件上传、关闭证据和时间线。"),
        ("阶段 3", "正式推广", "逐步扩大到全部项目，并纳入项目团队、统计报表与管理驾驶舱。"),
    ]
    for idx, item in enumerate(phases):
        x = Inches(0.95 + idx * 4.05)
        fill = SOFT_ORANGE if idx == 0 else SOFT_TEAL if idx == 1 else WHITE
        panel(slide, x, Inches(2.18), Inches(3.45), Inches(2.95), fill=fill)
        add_text(slide, x + Inches(0.16), Inches(2.48), Inches(0.8), Inches(0.2), item[0], 12, MUTED, True)
        add_text(slide, x + Inches(0.16), Inches(2.82), Inches(2.4), Inches(0.25), item[1], 19, NAVY, True)
        add_text(slide, x + Inches(0.16), Inches(3.33), Inches(2.95), Inches(1.12), item[2], 13.5, SLATE)

    add_text(slide, Inches(1.0), Inches(6.2), Inches(10.8), Inches(0.28),
             "建议同时保留：源码仓库、.env 正式配置、数据库备份脚本、发布记录、系统介绍 PPT 和本部署指南。", 14, SLATE)


def add_end(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "最终建议", "先按一台服务器标准部署，再围绕小程序正式发布完成域名和审核。")

    panel(slide, Inches(0.82), Inches(1.55), Inches(12.0), Inches(4.95))
    add_text(slide, Inches(1.05), Inches(1.92), Inches(4.0), Inches(0.28), "执行顺序建议", 20, NAVY, True)
    add_bullets(slide, Inches(1.08), Inches(2.35), Inches(10.8), Inches(3.6), [
        "第一步：在 Linux 服务器用 Docker Compose 跑通 mysql / redis / minio / backend / web。",
        "第二步：把 MINIO_PUBLIC_ENDPOINT、REPORT_PUBLIC_BASE_URL、AppID、Secret 配成正式值。",
        "第三步：配置 HTTPS 域名与微信合法域名，完成真机联调。",
        "第四步：从一个项目开始试点，上线后按项目逐步推广。",
    ], 17)
    add_text(slide, Inches(1.08), Inches(5.95), Inches(10.5), Inches(0.22),
             "交付文件：部署指南 Markdown + 部署与发布指南 PPT，可直接给 IT、项目经理和小程序发布负责人使用。", 13.5, SLATE)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_cover(prs)
    add_scope(prs, 1)
    add_arch(prs, 2)
    add_requirements(prs, 3)
    add_ports(prs, 4)
    add_local_deploy(prs, 5)
    add_prod_env(prs, 6)
    add_linux_steps(prs, 7)
    add_ops(prs, 8)
    add_miniapp_pre(prs, 9)
    add_miniapp_release(prs, 10)
    add_acceptance(prs, 11)
    add_rollout(prs, 12)
    add_end(prs, 13)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
