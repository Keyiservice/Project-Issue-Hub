from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_PPT = ROOT / "非标设备项目现场问题协同系统-老板汇报版.pptx"
REAL = ROOT / "ppt_real_assets"

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(25, 33, 43)
SLATE = RGBColor(78, 89, 102)
MUTED = RGBColor(119, 130, 143)
PAPER = RGBColor(248, 244, 237)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(210, 106, 37)
TEAL = RGBColor(57, 145, 133)
SOFT_ORANGE = RGBColor(248, 232, 220)
SOFT_TEAL = RGBColor(226, 242, 238)
LINE = RGBColor(228, 222, 214)
SUCCESS = RGBColor(37, 123, 92)
WARNING = RGBColor(172, 105, 26)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.1), Inches(-0.8), Inches(4.4), Inches(4.4))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = TEAL
    glow1.fill.transparency = 0.9
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.9), Inches(5.0), Inches(3.2), Inches(3.2))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ORANGE
    glow2.fill.transparency = 0.92
    glow2.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def panel(slide, left, top, width, height, fill=WHITE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def title_block(slide, page_no, title, subtitle):
    add_text(slide, Inches(0.72), Inches(0.38), Inches(8.6), Inches(0.3), f"{page_no:02d}", 12, MUTED, True)
    add_text(slide, Inches(1.15), Inches(0.34), Inches(7.6), Inches(0.5), title, 26, NAVY, True)
    add_text(slide, Inches(1.18), Inches(0.82), Inches(9.1), Inches(0.3), subtitle, 12, SLATE)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    left = panel(slide, Inches(0.62), Inches(0.6), Inches(5.35), Inches(6.2), fill=NAVY)
    left.line.fill.background()
    add_text(slide, Inches(0.95), Inches(0.95), Inches(3.6), Inches(0.35), "FIELD ISSUE OPS", 11, RGBColor(198, 207, 215), False)
    add_text(slide, Inches(0.95), Inches(1.35), Inches(4.5), Inches(1.2), "非标设备项目现场\n问题协同系统", 28, WHITE, True)
    add_text(slide, Inches(0.95), Inches(2.55), Inches(4.25), Inches(1.0), "老板汇报版\n聚焦价值、现状、闭环与推广。", 17, RGBColor(226, 232, 238))

    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.82), Inches(2.6), Inches(0.52))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(53, 63, 74)
    pill.line.fill.background()
    add_text(slide, Inches(1.18), Inches(5.93), Inches(2.1), Inches(0.18), "现场提报 + 管理驾驶舱", 12, WHITE, True)

    add_text(slide, Inches(6.3), Inches(0.92), Inches(5.9), Inches(0.3), "NON-STANDARD EQUIPMENT DELIVERY", 11, MUTED, False)
    add_text(slide, Inches(6.3), Inches(1.28), Inches(5.8), Inches(0.9), "把 Excel OPL\n升级为真正可闭环的协同系统", 26, NAVY, True)
    add_text(slide, Inches(6.32), Inches(2.18), Inches(5.4), Inches(0.6), "目标不是替代表格本身，而是让管理层、项目经理、现场人员看到的是同一套真实状态。", 15, SLATE)

    for idx, item in enumerate([
        ("现场录入", "拍图/视频，少打字"),
        ("项目闭环", "责任、状态、验证可追踪"),
        ("管理看板", "高优先级、超期、未分派一眼看清"),
    ]):
        x = Inches(6.32 + idx * 1.95)
        y = Inches(3.18)
        card = panel(slide, x, y, Inches(1.75), Inches(1.48), WHITE)
        add_text(slide, x + Inches(0.14), y + Inches(0.18), Inches(1.3), Inches(0.25), item[0], 16, NAVY, True)
        add_text(slide, x + Inches(0.14), y + Inches(0.55), Inches(1.45), Inches(0.6), item[1], 12.5, SLATE)

    add_text(slide, Inches(6.32), Inches(6.35), Inches(3.0), Inches(0.2), "汇报对象：管理层 / 项目负责人", 10.5, MUTED)
    add_text(slide, Inches(10.5), Inches(6.35), Inches(2.0), Inches(0.2), "2026-03-25", 10.5, MUTED, False, PP_ALIGN.RIGHT)


def add_summary(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "一句话结论", "这不是一个新的表格工具，而是一套项目现场问题闭环平台。")

    panel(slide, Inches(0.72), Inches(1.38), Inches(12.0), Inches(1.18), fill=WHITE)
    add_text(slide, Inches(0.95), Inches(1.65), Inches(11.2), Inches(0.6),
             "它把“现场发生了什么、谁在负责、当前推进到哪一步、什么时候能关闭”统一到一个系统里，"
             "让现场提报、项目推进、管理判断用的是同一份真实状态。", 21, NAVY, True)

    card_titles = [
        ("解决现状", "Excel 只有文字、责任不清、流转不透明、反复追问。"),
        ("系统形态", "微信小程序负责轻操作，Web 驾驶舱负责重管理，后端统一规则。"),
        ("当前状态", "核心主链路已跑通，具备项目级问题库、闭环、统计、权限和附件能力。"),
        ("管理价值", "把判断从“追问现场”改成“看系统状态”，把协同从被动追踪改成主动推进。"),
    ]
    for idx, (title, desc) in enumerate(card_titles):
        x = Inches(0.72 + (idx % 2) * 6.08)
        y = Inches(2.82 + (idx // 2) * 1.76)
        panel(slide, x, y, Inches(5.92), Inches(1.5), WHITE)
        add_text(slide, x + Inches(0.18), y + Inches(0.18), Inches(1.9), Inches(0.25), title, 17, NAVY, True)
        add_text(slide, x + Inches(0.18), y + Inches(0.52), Inches(5.4), Inches(0.7), desc, 13.8, SLATE)


def add_pains(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "为什么必须做", "当前 Excel OPL 的问题，不在于有没有记录，而在于无法形成管理动作。")

    left = panel(slide, Inches(0.72), Inches(1.45), Inches(5.7), Inches(5.45), fill=WHITE)
    add_text(slide, Inches(0.95), Inches(1.72), Inches(2.3), Inches(0.3), "现状问题", 20, NAVY, True)
    add_bullets(slide, Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.6), [
        "记录只有文字，现场情况不完整，管理层很难在 10 秒内判断本质。",
        "经常反复追问“这是什么问题、影响多大、谁在处理、什么时候能关”。",
        "问题流转、责任归属、关闭验证缺乏统一规则，容易出现失焦和扯皮。",
        "现场和管理层使用的不是同一套状态，统计分析只能靠人工汇总。"
    ], 17)

    right = panel(slide, Inches(6.63), Inches(1.45), Inches(6.1), Inches(5.45), fill=WHITE)
    add_text(slide, Inches(6.88), Inches(1.72), Inches(2.5), Inches(0.3), "系统要达到的结果", 20, NAVY, True)
    add_bullets(slide, Inches(6.92), Inches(2.1), Inches(5.45), Inches(4.55), [
        "现场员工随时按项目提报问题，少打字、多拍图、多拍视频。",
        "项目经理快速识别高优先级、超期、未分派和待验证问题。",
        "管理层打开系统即可看每个项目当前节奏、压力和重点风险。",
        "问题从提报到关闭形成标准闭环，附件、评论、日志和证据可追溯。"
    ], 17)


def add_architecture(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "系统怎么组成", "一套后端规则，承接两端使用场景。")

    roles = [
        ("现场人员", "微信小程序", "按项目提报问题\n上传图片/视频\n跟进评论"),
        ("项目经理/管理层", "Web 管理端", "项目问题库\n驾驶舱分析\n团队配置"),
        ("统一业务中心", "Spring Boot", "状态机\nRBAC\n编号\n超期计算"),
        ("底层基础设施", "MySQL/Redis/MinIO", "业务主库\n缓存\n附件对象存储"),
    ]
    for idx, (who, stack, desc) in enumerate(roles):
        x = Inches(0.75 + idx * 3.05)
        card = panel(slide, x, Inches(2.0), Inches(2.72), Inches(3.2), fill=WHITE)
        add_text(slide, x + Inches(0.14), Inches(2.25), Inches(2.2), Inches(0.24), who, 16.5, NAVY, True)
        badge = panel(slide, x + Inches(0.14), Inches(2.63), Inches(1.25), Inches(0.35), SOFT_TEAL if idx >= 2 else SOFT_ORANGE)
        badge.line.fill.background()
        add_text(slide, x + Inches(0.23), Inches(2.69), Inches(1.0), Inches(0.14), stack, 10.5, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.16), Inches(3.18), Inches(2.28), Inches(1.3), desc, 13, SLATE)
        if idx < len(roles) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(2.83), Inches(3.08), Inches(0.28), Inches(0.38))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()

    footer = panel(slide, Inches(0.88), Inches(5.72), Inches(11.42), Inches(0.65), fill=WHITE)
    add_text(slide, Inches(1.08), Inches(5.9), Inches(10.9), Inches(0.2),
             "核心原则：前端只负责操作体验，业务规则统一放在后端，避免状态流转和权限规则散落在页面里。", 13.6, SLATE)


def add_process(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "问题闭环怎么运行", "重点不是 CRUD，而是按项目、按责任、按状态推进。")

    steps = [
        ("1 选择项目", "先锁定项目，再进入该项目自己的问题库。"),
        ("2 现场提报", "录入标题、描述、图片、视频和时间。"),
        ("3 责任分派", "只允许分派给当前项目团队成员。"),
        ("4 处理中", "通过评论、附件、时间线持续跟进。"),
        ("5 待验证", "提交处理说明，等待现场或项目经理验证。"),
        ("6 关闭归档", "记录关闭人、关闭时间、关闭说明和关闭证据。"),
    ]
    for idx, (title, desc) in enumerate(steps):
        x = Inches(0.82 + idx * 2.05)
        y = Inches(2.05 if idx % 2 == 0 else 3.2)
        panel(slide, x, y, Inches(1.82), Inches(1.24), fill=WHITE)
        add_text(slide, x + Inches(0.12), y + Inches(0.12), Inches(1.45), Inches(0.2), title, 14.5, NAVY, True)
        add_text(slide, x + Inches(0.12), y + Inches(0.42), Inches(1.5), Inches(0.55), desc, 11.5, SLATE)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(1.84), y + Inches(0.45), Inches(0.24), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE if idx % 2 == 0 else TEAL
            arrow.line.fill.background()

    add_text(slide, Inches(0.92), Inches(5.35), Inches(11.0), Inches(0.4),
             "管理价值在于：每个问题都有“当前阶段、责任人、下一步、证据”，而不是停留在模糊描述。", 19, NAVY, True)


def add_web_flow(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "Web 端：项目视角的管理入口", "先选项目，再进入该项目的驾驶舱和问题库。")
    add_image(slide, "issue_projects_real.png", Inches(0.72), Inches(1.5), Inches(5.9), Inches(4.95))
    add_image(slide, "dashboard_real.png", Inches(6.76), Inches(1.5), Inches(5.6), Inches(4.95))
    add_bullets(slide, Inches(0.86), Inches(6.62), Inches(11.2), Inches(0.3), [
        "项目入口页解决“所有项目问题混在一起”的问题，驾驶舱则负责集中呈现当前项目的节奏与风险。",
        "不选项目时可以轮巡项目，适合管理层集中查看；选定项目后则进入项目级管理视角。"
    ], 13.4)


def add_issue_detail(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "Web 端：问题详情就是管理动作入口", "管理层不需要看很多字段，只关心当前阶段、责任人、下一步和证据。")
    add_image(slide, "issue_detail_real.png", Inches(0.72), Inches(1.5), Inches(7.2), Inches(5.15))
    add_bullets(slide, Inches(8.15), Inches(1.86), Inches(4.0), Inches(4.8), [
        "顶部直接看问题本体和附件，避免再去翻评论找现场照片。",
        "右侧只保留当前推进信息，减少和时间线重复的冗余描述。",
        "时间线负责记录已发生动作，状态推进负责推动下一步动作。",
        "关闭问题会记录关闭人、关闭时间、关闭说明和关闭证据。"
    ], 16.5)


def add_miniapp(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "小程序：现场轻操作入口", "登录稳定、首页清晰、入口收敛。")
    add_image(slide, "miniapp_login_crop.png", Inches(0.9), Inches(1.55), Inches(2.65), Inches(5.2))
    add_image(slide, "miniapp_home_real.png", Inches(3.82), Inches(1.55), Inches(2.88), Inches(5.2))
    add_bullets(slide, Inches(7.02), Inches(1.9), Inches(5.0), Inches(4.8), [
        "登录方案采用企业账号为主身份，微信只做绑定和免密入口，保证责任归属稳定。",
        "首页只保留项目工作台、OPL 录入、项目问题库、项目统计四个主入口，不做冗余说明。",
        "这端负责轻操作：提报、跟进、拍图/视频、看我相关的问题。"
    ], 17)


def add_identity(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "身份与组织：为什么不是纯微信登录", "权限、岗位、项目职责必须绑定企业身份。")
    add_image(slide, "users_real.png", Inches(0.72), Inches(1.55), Inches(6.8), Inches(5.15))
    add_bullets(slide, Inches(7.82), Inches(1.92), Inches(4.25), Inches(4.8), [
        "系统主身份是企业账号，微信只做首次绑定和后续免密进入。",
        "系统角色管权限，项目成员岗位管项目内职责，避免把“岗位”和“权限”混在一起。",
        "责任人下拉不是全公司所有人，而是当前项目团队成员。"
    ], 16.6)


def add_stats(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "管理收益：统计和风险从人工追问转为系统可见", "项目经理和管理层可以按项目快速识别重点问题。")
    add_image(slide, "stats_real.png", Inches(0.72), Inches(1.55), Inches(7.3), Inches(5.15))
    add_bullets(slide, Inches(8.25), Inches(1.92), Inches(3.9), Inches(4.8), [
        "看板能直接看到问题总数、未关闭、高优先级、超期、状态分布和趋势。",
        "项目统计必须按项目口径计算，而不是当前页口径，这个逻辑已经收口修正。",
        "这类页面适合周会、评审会和项目例会，不适合现场提报。"
    ], 16.4)


def add_delivery(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "当前已经做到哪一步", "主链路不是 Demo 级结构，已经具备企业使用骨架。")

    left = panel(slide, Inches(0.72), Inches(1.55), Inches(5.85), Inches(5.05), fill=WHITE)
    add_text(slide, Inches(0.96), Inches(1.82), Inches(1.7), Inches(0.25), "已完成", 20, SUCCESS, True)
    add_bullets(slide, Inches(1.02), Inches(2.18), Inches(5.1), Inches(4.0), [
        "项目级问题库：先选项目，再进入项目自己的问题池。",
        "现场提报主链路：图片/视频上传、评论跟进、关闭证据。",
        "Web 管理端：驾驶舱、项目管理、项目团队、用户管理、统计分析。",
        "企业账号 + 微信绑定登录、RBAC、项目团队岗位、Docker 本地部署。"
    ], 16)

    right = panel(slide, Inches(6.82), Inches(1.55), Inches(5.5), Inches(5.05), fill=WHITE)
    add_text(slide, Inches(7.06), Inches(1.82), Inches(1.9), Inches(0.25), "下一步建议", 20, WARNING, True)
    add_bullets(slide, Inches(7.12), Inches(2.18), Inches(4.8), Inches(4.0), [
        "做正式培训版和上线 SOP，降低不同角色的上手门槛。",
        "补真机域名、HTTPS 和微信后台合法域名配置，完成真实发布准备。",
        "继续补更多真实项目数据和项目成员，形成可直接推广的模板项目。",
        "上线后优先观察：超期识别、分派准确性、关闭证据质量。"
    ], 16)


def add_rollout(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "建议的落地方式", "不要一次铺开，先从一个项目把节奏跑顺。")
    phases = [
        ("阶段 1", "单项目试运行", "以 FS21 这类典型项目为试点，先把提报、分派、跟进、关闭跑顺。"),
        ("阶段 2", "项目团队固化", "补齐各项目成员与岗位，确保责任人下拉和项目边界准确。"),
        ("阶段 3", "管理例会接入", "让项目周会直接看驾驶舱和项目统计，不再靠 Excel 汇总。"),
        ("阶段 4", "多项目推广", "沉淀导入模板、权限模板、培训材料，再向更多项目复制。"),
    ]
    for idx, (title, head, desc) in enumerate(phases):
        y = Inches(1.55 + idx * 1.3)
        panel(slide, Inches(0.88), y, Inches(11.45), Inches(0.98), WHITE)
        badge = panel(slide, Inches(1.05), y + Inches(0.16), Inches(1.05), Inches(0.34), SOFT_ORANGE if idx % 2 == 0 else SOFT_TEAL)
        badge.line.fill.background()
        add_text(slide, Inches(1.12), y + Inches(0.2), Inches(0.88), Inches(0.14), title, 11, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, Inches(2.34), y + Inches(0.17), Inches(2.6), Inches(0.2), head, 17, NAVY, True)
        add_text(slide, Inches(4.68), y + Inches(0.16), Inches(6.95), Inches(0.42), desc, 13.5, SLATE)


def add_support(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "管理层需要支持什么", "系统能落地，不只是技术问题，也取决于管理动作。")

    panel(slide, Inches(0.72), Inches(1.55), Inches(12.0), Inches(4.9), WHITE)
    supports = [
        ("1", "明确 OPL 是否以系统为准", "如果系统和 Excel 长期并行且没有主从关系，现场会回到旧习惯。"),
        ("2", "明确项目负责人和项目团队", "项目边界和责任人池必须真实，否则分派和统计都会失真。"),
        ("3", "要求关闭必须带证据", "只有关闭说明和证据形成标准，系统闭环才有意义。"),
        ("4", "把周会切到系统看板", "一旦会议仍看 Excel，系统就会沦为二次录入工具。"),
    ]
    for idx, (num, title, desc) in enumerate(supports):
        x = Inches(0.98 + (idx % 2) * 5.95)
        y = Inches(1.9 + (idx // 2) * 1.55)
        tag = panel(slide, x, y, Inches(0.48), Inches(0.38), SOFT_TEAL)
        tag.line.fill.background()
        add_text(slide, x + Inches(0.14), y + Inches(0.08), Inches(0.18), Inches(0.12), num, 12, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.68), y, Inches(2.9), Inches(0.24), title, 16.5, NAVY, True)
        add_text(slide, x + Inches(0.68), y + Inches(0.34), Inches(4.45), Inches(0.5), desc, 13.5, SLATE)

    add_text(slide, Inches(0.9), Inches(6.75), Inches(10.8), Inches(0.18),
             "如果管理动作明确，这套系统可以很快从“可用”进入“被真正使用”。", 16, NAVY, True)


def add_close(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, page_no, "结论", "建议进入试运行和推广准备阶段。")

    panel(slide, Inches(0.78), Inches(1.65), Inches(11.8), Inches(3.2), WHITE)
    add_text(slide, Inches(1.0), Inches(2.05), Inches(11.0), Inches(0.7),
             "系统已经具备从“项目选择 -> 现场提报 -> 责任分派 -> 跟进推进 -> 验证关闭 -> 管理统计”的完整骨架，"
             "下一步最重要的是按项目试运行、把管理会议切到系统、把关闭证据做实。", 24, NAVY, True)

    add_bullets(slide, Inches(1.0), Inches(3.05), Inches(10.9), Inches(1.2), [
        "建议先用 1 个项目跑顺，再模板化复制到更多项目。",
        "技术栈与部署方式已经适合企业内网和私有化持续迭代。"
    ], 17)

    add_text(slide, Inches(0.98), Inches(5.55), Inches(3.8), Inches(0.3), "谢谢", 28, NAVY, True)
    add_text(slide, Inches(0.98), Inches(5.98), Inches(5.2), Inches(0.3), "Q&A / 下一步可直接进入试运行筹备", 15, SLATE)


def add_image(slide, name, left, top, width, height):
    panel(slide, left, top, width, height, WHITE)
    path = REAL / name
    if path.exists():
        slide.shapes.add_picture(str(path), left + Inches(0.08), top + Inches(0.08), width=width - Inches(0.16), height=height - Inches(0.16))


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_cover(prs)
    add_summary(prs, 1)
    add_pains(prs, 2)
    add_architecture(prs, 3)
    add_process(prs, 4)
    add_web_flow(prs, 5)
    add_issue_detail(prs, 6)
    add_miniapp(prs, 7)
    add_identity(prs, 8)
    add_stats(prs, 9)
    add_delivery(prs, 10)
    add_rollout(prs, 11)
    add_support(prs, 12)
    add_close(prs, 13)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
