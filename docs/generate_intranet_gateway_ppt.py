from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REAL = ROOT / "ppt_real_assets"
OUT = ROOT / "Project-Issue-Hub-内网核心与小程序公网网关部署方案.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(20, 28, 38)
SLATE = RGBColor(82, 92, 104)
MUTED = RGBColor(127, 138, 149)
PAPER = RGBColor(247, 243, 236)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(214, 110, 38)
TEAL = RGBColor(56, 142, 132)
SOFT_ORANGE = RGBColor(248, 233, 221)
SOFT_TEAL = RGBColor(227, 242, 238)
SOFT_BLUE = RGBColor(231, 240, 250)
LINE = RGBColor(228, 222, 214)
RED = RGBColor(185, 76, 58)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()

    glow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.3), Inches(-0.9), Inches(4.6), Inches(4.6))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = TEAL
    glow1.fill.transparency = 0.93
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.7), Inches(5.0), Inches(3.4), Inches(3.4))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ORANGE
    glow2.fill.transparency = 0.94
    glow2.line.fill.background()


def panel(slide, left, top, width, height, fill=WHITE, rounded=True):
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(st, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)

        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Microsoft YaHei"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = TEAL

        body = p.add_run()
        body.text = item
        body.font.name = "Microsoft YaHei"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def title_block(slide, page_no, title, subtitle):
    add_text(slide, Inches(0.72), Inches(0.32), Inches(0.55), Inches(0.2), f"{page_no:02d}", 11, MUTED, True)
    add_text(slide, Inches(1.12), Inches(0.28), Inches(9.2), Inches(0.42), title, 26, NAVY, True)
    add_text(slide, Inches(1.15), Inches(0.76), Inches(10.2), Inches(0.25), subtitle, 12.5, SLATE)


def add_image(slide, filename, left, top, width, height):
    path = REAL / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    left = panel(slide, Inches(0.62), Inches(0.58), Inches(5.35), Inches(6.15), fill=NAVY)
    left.line.fill.background()
    add_text(slide, Inches(0.95), Inches(0.9), Inches(3.8), Inches(0.25), "PROJECT ISSUE HUB", 12, RGBColor(203, 212, 220))
    add_text(slide, Inches(0.95), Inches(1.28), Inches(4.7), Inches(1.2),
             "内网核心系统 +\n小程序公网网关\n部署方案", 28, WHITE, True)
    add_text(slide, Inches(0.95), Inches(3.18), Inches(4.25), Inches(1.0),
             "面向管理层、IT 与运维团队。\n目标是数据留在公司内网，同时让微信小程序可正式上线使用。", 17, RGBColor(228, 234, 240))

    badge = panel(slide, Inches(0.95), Inches(5.9), Inches(3.15), Inches(0.52), fill=RGBColor(54, 64, 75))
    badge.line.fill.background()
    add_text(slide, Inches(1.12), Inches(5.98), Inches(2.8), Inches(0.18), "Private Core + Public Gateway", 11, WHITE, True)

    add_text(slide, Inches(6.28), Inches(0.92), Inches(5.9), Inches(0.25), "FOR LEADERSHIP & IT", 11, MUTED)
    add_text(slide, Inches(6.28), Inches(1.28), Inches(6.0), Inches(0.72),
             "既不把数据库和后台暴露公网，又满足小程序在手机外网环境下访问。", 24, NAVY, True)
    add_bullets(slide, Inches(6.34), Inches(2.15), Inches(5.7), Inches(1.45), [
        "Web 管理端仅在公司内网开放",
        "小程序只访问一个受控的公网 HTTPS 网关",
        "核心后端、MySQL、Redis、MinIO 全部留在内网",
    ], 16)

    add_image(slide, "dashboard_real.png", Inches(6.3), Inches(3.85), Inches(3.0), Inches(2.35))
    add_image(slide, "miniapp_home_real.png", Inches(9.55), Inches(3.84), Inches(2.45), Inches(2.4))
    add_text(slide, Inches(10.7), Inches(6.84), Inches(1.4), Inches(0.14), "2026-04", 10.5, MUTED, False, PP_ALIGN.RIGHT)


def add_pain(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "为什么不能做成“纯内网 + 小程序直连”", "关键约束不在项目代码，而在微信小程序的正式运行规则。")

    panel(slide, Inches(0.82), Inches(1.5), Inches(5.65), Inches(4.8), fill=SOFT_ORANGE)
    add_text(slide, Inches(1.05), Inches(1.82), Inches(3.2), Inches(0.22), "微信小程序正式版的硬约束", 19, NAVY, True)
    add_bullets(slide, Inches(1.08), Inches(2.2), Inches(5.0), Inches(3.55), [
        "请求域名必须配置到微信后台的合法域名列表",
        "正式版必须走 HTTPS，不能再用 127.0.0.1",
        "手机必须能真实访问到这个域名",
        "如果后端只在公司内网、外网手机访问不到，小程序就无法正常工作",
    ], 16)

    panel(slide, Inches(6.66), Inches(1.5), Inches(5.85), Inches(4.8), fill=WHITE)
    add_text(slide, Inches(6.92), Inches(1.82), Inches(3.0), Inches(0.22), "正确的结论", 19, NAVY, True)
    add_bullets(slide, Inches(6.95), Inches(2.2), Inches(5.2), Inches(3.0), [
        "不能让小程序直接访问纯内网服务",
        "可以让小程序访问一个最小化的公网 HTTPS 网关",
        "网关再通过白名单访问内网后端",
        "这样既满足微信规则，也不需要把整套系统暴露公网",
    ], 16)
    add_text(slide, Inches(6.95), Inches(5.35), Inches(5.0), Inches(0.5),
             "一句话：\n数据留在内网，入口放在公网，但只开放一层受控网关。", 18, RED, True)


def add_architecture(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "推荐架构", "这套结构最符合你当前项目：安全边界清晰，改造成本低。")

    # Top actors
    panel(slide, Inches(0.8), Inches(1.65), Inches(2.5), Inches(1.15), fill=WHITE)
    add_text(slide, Inches(1.0), Inches(1.92), Inches(2.1), Inches(0.18), "微信小程序", 20, NAVY, True)
    add_text(slide, Inches(1.0), Inches(2.28), Inches(2.0), Inches(0.28), "现场人员 / 管理层\n手机外网访问", 13.5, SLATE)

    panel(slide, Inches(5.0), Inches(1.45), Inches(3.25), Inches(1.55), fill=SOFT_TEAL)
    add_text(slide, Inches(5.22), Inches(1.78), Inches(2.8), Inches(0.22), "公网 HTTPS 网关", 20, NAVY, True)
    add_text(slide, Inches(5.22), Inches(2.18), Inches(2.7), Inches(0.36), "pih-api.company.com\nNginx / API Gateway / WAF", 13.5, SLATE)

    panel(slide, Inches(9.2), Inches(1.65), Inches(3.3), Inches(1.15), fill=WHITE)
    add_text(slide, Inches(9.45), Inches(1.92), Inches(2.4), Inches(0.18), "公司内网浏览器", 20, NAVY, True)
    add_text(slide, Inches(9.45), Inches(2.28), Inches(2.6), Inches(0.18), "项目经理 / 管理层 / 管理员", 13.5, SLATE)

    for x in [Inches(3.45), Inches(8.45)]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, Inches(2.0), Inches(0.45), Inches(0.42))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()

    panel(slide, Inches(1.2), Inches(3.55), Inches(11.0), Inches(2.2), fill=WHITE)
    add_text(slide, Inches(1.45), Inches(3.8), Inches(3.2), Inches(0.22), "内网核心区", 19, NAVY, True)

    blocks = [
        ("Web 管理端", "Vue 3 / 内网域名", Inches(1.5)),
        ("Backend", "Spring Boot / 统一规则", Inches(4.0)),
        ("MySQL", "业务主数据", Inches(6.55)),
        ("Redis", "缓存 / 会话", Inches(8.45)),
        ("MinIO", "图片 / 视频附件", Inches(10.1)),
    ]
    for title, desc, x in blocks:
        panel(slide, x, Inches(4.28), Inches(1.55), Inches(1.05), fill=SOFT_BLUE if title == "Backend" else PAPER)
        add_text(slide, x + Inches(0.12), Inches(4.47), Inches(1.22), Inches(0.18), title, 16, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.08), Inches(4.82), Inches(1.34), Inches(0.22), desc, 11.5, SLATE, False, PP_ALIGN.CENTER)


def add_exposure(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "什么应该暴露，什么不应该暴露", "这一页给 IT 做边界判断。")

    panel(slide, Inches(0.82), Inches(1.5), Inches(5.7), Inches(4.95), fill=SOFT_TEAL)
    panel(slide, Inches(6.75), Inches(1.5), Inches(5.75), Inches(4.95), fill=SOFT_ORANGE)

    add_text(slide, Inches(1.05), Inches(1.82), Inches(2.6), Inches(0.22), "允许对公网开放", 20, NAVY, True)
    add_bullets(slide, Inches(1.08), Inches(2.18), Inches(5.0), Inches(3.65), [
        "pih-api.company.com 的 443 端口",
        "仅代理小程序需要的 /api/ 接口",
        "附件上传与预览接口",
        "登录、项目列表、问题详情、评论、状态流转等必要接口",
        "可选接入 WAF、限流、审计日志",
    ], 16)

    add_text(slide, Inches(6.98), Inches(1.82), Inches(2.6), Inches(0.22), "不建议对公网开放", 20, NAVY, True)
    add_bullets(slide, Inches(7.02), Inches(2.18), Inches(5.0), Inches(3.65), [
        "Web 管理后台",
        "MySQL 数据库",
        "Redis 缓存",
        "MinIO Console",
        "管理员接口、用户管理、字典管理、项目团队管理接口",
        "Swagger 文档入口",
    ], 16)


def add_network(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "网络与防火墙建议", "让网关只做入口，不让数据服务直接暴露。")

    rows = [
        ("公网 -> 网关", "443", "开放", "小程序正式访问入口"),
        ("公网 -> 内网 Backend", "-", "关闭", "禁止直接访问业务后端"),
        ("公网 -> MySQL / Redis / MinIO Console", "-", "关闭", "禁止直接访问数据层"),
        ("网关 -> 内网 Backend", "8081", "白名单开放", "仅允许网关服务器访问"),
        ("Backend -> MySQL", "3306", "内网开放", "业务数据读写"),
        ("Backend -> Redis", "6379", "内网开放", "缓存 / 会话"),
        ("Backend -> MinIO", "9000", "内网开放", "附件读写"),
    ]

    for idx, row in enumerate(rows):
        y = Inches(1.55 + idx * 0.67)
        fill = WHITE if idx % 2 == 0 else PAPER
        panel(slide, Inches(0.85), y, Inches(11.6), Inches(0.5), fill=fill)
        add_text(slide, Inches(1.02), y + Inches(0.1), Inches(2.0), Inches(0.14), row[0], 13, NAVY, True)
        add_text(slide, Inches(3.5), y + Inches(0.1), Inches(0.9), Inches(0.14), row[1], 12.5, SLATE)
        add_text(slide, Inches(4.6), y + Inches(0.1), Inches(1.4), Inches(0.14), row[2], 12.5, RED if row[2] == "关闭" else TEAL, True)
        add_text(slide, Inches(6.0), y + Inches(0.1), Inches(5.8), Inches(0.14), row[3], 12.5, SLATE)


def add_gateway(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "公网网关怎么落地", "对外只给一个域名和 443，后面由网关反代到内网。")

    panel(slide, Inches(0.82), Inches(1.5), Inches(5.6), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(1.02), Inches(1.82), Inches(2.8), Inches(0.22), "Nginx / 网关职责", 20, NAVY, True)
    add_bullets(slide, Inches(1.05), Inches(2.18), Inches(5.0), Inches(3.65), [
        "只监听公网 443",
        "配置小程序合法域名证书",
        "只放行 /api/ 路径",
        "请求转发到内网 backend:8081",
        "限制上传体积，保留真实 IP",
        "其余路径直接 403，避免把后台暴露出去",
    ], 16)

    panel(slide, Inches(6.68), Inches(1.5), Inches(5.82), Inches(4.95), fill=SOFT_BLUE)
    add_text(slide, Inches(6.9), Inches(1.82), Inches(3.6), Inches(0.22), "已经给你的样例文件", 20, NAVY, True)
    add_text(slide, Inches(6.94), Inches(2.24), Inches(5.0), Inches(0.85),
             "docs/ops/nginx-miniapp-public-gateway.conf\n\n"
             "把其中的内网地址 10.10.20.15:8081 改成你们实际 backend 地址即可。", 15, SLATE, False, PP_ALIGN.LEFT, "Consolas")
    add_text(slide, Inches(6.94), Inches(4.25), Inches(5.0), Inches(0.8),
             "推荐公网域名：\n"
             "pih-api.company.com", 18, NAVY, True)


def add_release(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "小程序正式发布前必须满足的条件", "这一页是小程序上线前的硬门槛。")

    add_image(slide, "miniapp_home_real.png", Inches(0.86), Inches(1.68), Inches(2.6), Inches(4.8))
    panel(slide, Inches(3.85), Inches(1.56), Inches(8.55), Inches(5.0), fill=WHITE)

    add_bullets(slide, Inches(4.08), Inches(1.95), Inches(8.0), Inches(4.0), [
        "后端 .env 切到正式参数：WECHAT_MINIAPP_MOCK_ENABLED=false",
        "配置正式 AppID 与 Secret",
        "接口地址不能再用 127.0.0.1，必须是公网 HTTPS 域名",
        "微信公众平台配置 request / uploadFile / downloadFile 合法域名",
        "真机联调通过：登录、选项目、创建问题、上传图片视频、详情预览、跟进、关闭",
    ], 16)
    add_text(slide, Inches(4.08), Inches(5.95), Inches(7.8), Inches(0.3),
             "建议：先用一个试点项目走完整链路，再提交微信审核。", 15, RED, True)


def add_steps(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "推荐实施顺序", "这套顺序最稳，避免先发小程序、后补基础设施。")

    steps = [
        ("1", "先部署内网核心", "先把 Web、Backend、MySQL、Redis、MinIO 在内网跑稳定。"),
        ("2", "再部署公网网关", "只开放 pih-api.company.com 给小程序。"),
        ("3", "完成真机联调", "确保上传、预览、跟进、关闭都正常。"),
        ("4", "最后小程序审核发布", "审核通过后再做项目试点和逐步推广。"),
    ]
    for idx, (no, title, desc) in enumerate(steps):
        x = Inches(0.9 + idx * 3.05)
        fill = SOFT_ORANGE if idx in (0, 3) else WHITE
        panel(slide, x, Inches(2.15), Inches(2.7), Inches(2.55), fill=fill)
        bubble = panel(slide, x + Inches(0.16), Inches(2.32), Inches(0.44), Inches(0.44), fill=SOFT_TEAL)
        bubble.line.fill.background()
        add_text(slide, x + Inches(0.27), Inches(2.4), Inches(0.2), Inches(0.12), no, 12, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.7), Inches(2.34), Inches(1.7), Inches(0.2), title, 16, NAVY, True)
        add_text(slide, x + Inches(0.18), Inches(2.82), Inches(2.2), Inches(0.75), desc, 12.5, SLATE)
        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(2.8), Inches(3.05), Inches(0.18), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()


def add_summary(slide_no, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, slide_no, "最终建议", "给领导和 IT 的一句话结论。")

    panel(slide, Inches(0.85), Inches(1.65), Inches(11.8), Inches(4.8), fill=WHITE)
    add_text(slide, Inches(1.15), Inches(2.0), Inches(4.0), Inches(0.28), "推荐部署原则", 22, NAVY, True)
    add_text(slide, Inches(1.18), Inches(2.6), Inches(10.8), Inches(1.0),
             "数据与后台留在内网，小程序通过一层最小公网网关接入。", 28, RED, True)
    add_bullets(slide, Inches(1.2), Inches(3.8), Inches(10.6), Inches(1.75), [
        "这能同时满足：公司不希望核心系统暴露公网；现场人员又必须通过微信小程序在外网手机上使用。",
        "对当前项目改造最小，IT 最容易接受，后续扩容到更多项目也最稳。",
    ], 18)
    add_text(slide, Inches(1.18), Inches(6.05), Inches(10.6), Inches(0.2),
             "配套文件：专项部署方案 Markdown + 公网网关 Nginx 样例，可直接交给运维执行。", 14, SLATE)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_cover(prs)
    add_pain(1, prs)
    add_architecture(2, prs)
    add_exposure(3, prs)
    add_network(4, prs)
    add_gateway(5, prs)
    add_release(6, prs)
    add_steps(7, prs)
    add_summary(8, prs)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
