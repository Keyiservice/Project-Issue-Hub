from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
BASE_PPT = ROOT / "非标设备项目现场问题协同系统-系统介绍与使用说明.pptx"
OUT_PPT = ROOT / "非标设备项目现场问题协同系统-第一版补图增强版.pptx"
REAL = ROOT / "ppt_real_assets"

INK = RGBColor(24, 34, 43)
MUTED = RGBColor(96, 107, 120)
SUBTLE = RGBColor(139, 150, 162)
PAPER = RGBColor(249, 246, 241)
PANEL = RGBColor(255, 252, 248)
PANEL_ALT = RGBColor(245, 249, 248)
LINE = RGBColor(225, 220, 212)
ORANGE = RGBColor(203, 93, 31)
TEAL = RGBColor(22, 123, 120)
WHITE = RGBColor(255, 255, 255)


def crop_login():
    src = REAL / "miniapp_devtools_foreground.png"
    out = REAL / "miniapp_login_crop.png"
    if src.exists() and not out.exists():
        img = Image.open(src)
        img.crop((42, 92, 355, 650)).save(out)
    return out


def add_bg(slide, width, height):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    blob1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.9), Inches(-0.7), Inches(4.5), Inches(4.5))
    blob1.fill.solid()
    blob1.fill.fore_color.rgb = TEAL
    blob1.fill.transparency = 0.9
    blob1.line.fill.background()
    blob2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.6), Inches(5.25), Inches(3.0), Inches(3.0))
    blob2.fill.solid()
    blob2.fill.fore_color.rgb = ORANGE
    blob2.fill.transparency = 0.93
    blob2.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Microsoft YaHei"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = TEAL
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color = INK


def add_header(slide, page_num, title, subtitle, width):
    kicker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(0.56), Inches(0.5), Inches(0.56), Inches(0.56))
    kicker.fill.background()
    kicker.line.color.rgb = TEAL
    kicker.line.width = Pt(3)
    add_text(slide, Inches(0.95), Inches(0.5), Inches(8.6), Inches(0.42), title, 26, INK, True)
    add_text(slide, Inches(0.98), Inches(0.94), Inches(9.0), Inches(0.35), subtitle, 12, MUTED)
    add_text(slide, Inches(10.95), Inches(0.32), Inches(1.55), Inches(0.22), "SUPPLEMENT", 10.5, SUBTLE, False, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(7.07), Inches(12.0), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    add_text(slide, Inches(0.65), Inches(7.08), Inches(4.2), Inches(0.22), "First Edition Enhanced", 9.5, SUBTLE)
    add_text(slide, Inches(12.05), Inches(7.02), Inches(0.35), Inches(0.22), str(page_num), 9.5, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=PANEL):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    return shp


def add_image_panel(slide, img_name, left, top, width, height, caption=None):
    add_panel(slide, left, top, width, height)
    path = REAL / img_name
    if path.exists():
        slide.shapes.add_picture(str(path), left + Inches(0.08), top + Inches(0.08), width=width - Inches(0.16), height=height - Inches(0.16))
    if caption:
        add_text(slide, left + Inches(0.12), top + height - Inches(0.34), width - Inches(0.24), Inches(0.22), caption, 10.5, WHITE, False, PP_ALIGN.RIGHT)


def add_two_image_slide(prs, page_num, title, subtitle, left_img, right_img, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle, prs.slide_width)
    add_image_panel(slide, left_img, Inches(0.72), Inches(1.55), Inches(5.85), Inches(4.45))
    add_image_panel(slide, right_img, Inches(6.77), Inches(1.55), Inches(5.55), Inches(4.45), "真实系统截图")
    add_panel(slide, Inches(0.72), Inches(6.15), Inches(11.6), Inches(0.72), PANEL_ALT)
    add_bullets(slide, Inches(0.94), Inches(6.28), Inches(11.1), Inches(0.4), bullets, 13.5)


def add_single_image_slide(prs, page_num, title, subtitle, img_name, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs.slide_width, prs.slide_height)
    add_header(slide, page_num, title, subtitle, prs.slide_width)
    add_image_panel(slide, img_name, Inches(0.72), Inches(1.55), Inches(8.1), Inches(5.05))
    add_panel(slide, Inches(8.98), Inches(1.55), Inches(3.35), Inches(5.05))
    add_bullets(slide, Inches(9.16), Inches(1.84), Inches(2.95), Inches(4.55), bullets, 14.8)


def main():
    crop_login()
    prs = Presentation(str(BASE_PPT))

    start_page = len(prs.slides) + 1
    add_two_image_slide(
        prs,
        start_page,
        "Web 实拍补充：驾驶舱与项目问题库",
        "这两页是管理层和项目经理最常打开的真实界面。",
        "dashboard_real.png",
        "issues_real.png",
        [
            "驾驶舱负责看全局节奏、项目轮巡和重点风险。",
            "项目问题库负责按项目筛选问题，并支持直接录入新问题。",
            "这一组图片补足了第一版 PPT 里“讲清楚了，但图不够”的问题。"
        ],
    )

    add_single_image_slide(
        prs,
        start_page + 1,
        "Web 实拍补充：问题详情页",
        "问题详情是整套系统最核心的一页。",
        "issue_detail_real.png",
        [
            "这里同时承载描述、附件、责任人、进度、时间线和关闭动作。",
            "图片和视频在描述区下方直接预览，不再藏在附件链接里。",
            "关闭问题时会记录关闭人、关闭时间、关闭说明和关闭证据。"
        ],
    )

    add_two_image_slide(
        prs,
        start_page + 2,
        "Web 实拍补充：项目管理与用户管理",
        "项目团队和企业账号体系，是系统跑稳的基础。",
        "projects_real.png",
        "users_real.png",
        [
            "项目管理决定项目主数据和项目团队成员。",
            "用户管理决定企业账号、角色、微信绑定和项目参与关系。",
            "责任人无法随意乱选，是因为系统按“项目团队成员”做了限制。"
        ],
    )

    add_single_image_slide(
        prs,
        start_page + 3,
        "Web 实拍补充：统计分析页",
        "这一页更适合复盘、周会和管理评审。",
        "stats_real.png",
        [
            "统计页偏分析，驾驶舱偏总览。",
            "适合看趋势、结构、超期压力和项目风险分布。",
            "后续还可以继续扩展导出和跨项目对比。"
        ],
    )

    add_single_image_slide(
        prs,
        start_page + 4,
        "小程序实拍补充：真实登录页",
        "当前已补入原始小程序截图，不再使用重绘示意图。",
        "miniapp_login_crop.png",
        [
            "这张图来自微信开发者工具当前真实页面截图，未改颜色。",
            "登录链路是：微信登录 -> 企业账号绑定 -> 首次强制改密。",
            "如果你要，我下一轮继续把小程序首页、项目问题库、详情页也抓成原图补进去。"
        ],
    )

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
